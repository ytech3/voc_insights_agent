"""
build_voc_hs3_vs_hs4.py
Tampa Bay Rays - VOC HS3 vs HS4 Homestand Comparison

Produces an 8-slide PPTX comparing HS3 (May 1-6) vs HS4 (May 15-20) with 2024 baseline.
No MLB benchmark columns.

7-column layout (slides 1-6):
  Metric | '24 Score | '26 AVG | HS3 Score | HS4 Score | HS4 vs HS3 % Diff | '26 vs '24 % Diff

Custom segment slides (7-8):
  Slide 7: In-Seat Delivery Deep Dive - Premium vs Everyone Else
  Slide 8: Ballpark App (Rays Wallet) - Concession Wait Expectations

Usage:
    python build_voc_hs3_vs_hs4.py
"""

import json
import os
import sys
import snowflake.connector
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
CONFIG_PATH = os.path.join(SCRIPT_DIR, "homestand_config.json")
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "VOC_HS3_vs_HS4.pptx")

# ── Snowflake connection ──
SF_ACCOUNT = "hta92307.east-us-2.azure"
SF_WAREHOUSE = "TBRDP_DW_CORTEX_XS_WH"
SF_DATABASE = "TBRDP_DW_DEV"
SF_SCHEMA = "IM_RPT"
SF_ROLE = "TBRDP_DW_PROD_CORTEX_USER"
SF_AUTHENTICATOR = "externalbrowser"

# ── Views ──
VOC_VIEW = "TBRDP_DW_DEV.IM_RPT.V_SBL_QUALTRICS_VOC_POST_ATTENDANCE_FULL_CORTEX_AI"
PARKING_VIEW = "TBRDP_DW_DEV.IM_RPT.V_SBL_QUALTRICS_VOC_POST_ATTENDANCE_FULL"

# ── Theme colors (Rays brand palette) ──
NAVY = RGBColor(0x09, 0x2C, 0x5C)
SKY = RGBColor(0x8F, 0xBC, 0xE6)
YELLOW = RGBColor(0xF5, 0xD1, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x4B, 0x4B, 0x4B)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xFF, 0x6B, 0x6B)
MEDIUM_BLUE = RGBColor(0x0D, 0x3B, 0x7A)
DARK_BLUE_BG = RGBColor(0x07, 0x20, 0x44)
MUTED_GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT_BG = RGBColor(0xF2, 0xF4, 0xF8)

# ── Logo ──
LOGO_PATH = os.path.join(SCRIPT_DIR, "MLB Logos", "TB_White.png")


# ═══════════════════════════════════════════════════════════════════════
# CONFIG & CONNECTION
# ═══════════════════════════════════════════════════════════════════════

def load_config():
    with open(CONFIG_PATH, "r") as f:
        return json.load(f)


def get_homestand(config, season, hs_num):
    for hs in config["homestands"]:
        if hs["season"] == season and hs["homestand_num"] == hs_num:
            return hs
    raise ValueError(f"Homestand {hs_num} for season {season} not found")


def get_connection():
    return snowflake.connector.connect(
        user="YTAKETANI@RAYSBASEBALL.COM",
        account=SF_ACCOUNT,
        authenticator=SF_AUTHENTICATOR,
        warehouse=SF_WAREHOUSE,
        database=SF_DATABASE,
        schema=SF_SCHEMA,
        role=SF_ROLE,
    )


def run_query(conn, sql):
    cur = conn.cursor()
    try:
        cur.execute(sql)
        return cur.fetchall()
    finally:
        cur.close()


# ═══════════════════════════════════════════════════════════════════════
# SQL HELPERS
# ═══════════════════════════════════════════════════════════════════════

def _baseline_avg(col, baseline_season):
    return f"AVG(CASE WHEN SEASON = {baseline_season} AND {col} < 80 THEN {col} END)"


def _baseline_pct_highly(col, baseline_season, valid_labels):
    return (
        f"ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END) "
        f"/ NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)"
    )


def _current_pct_highly(col, current_season, valid_labels):
    return (
        f"ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END) "
        f"/ NULLIF(SUM(CASE WHEN SEASON = {current_season} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)"
    )


def _hs_pct_highly(col, hs, valid_labels):
    return (
        f"ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs['start_date']}' AND '{hs['end_date']}' "
        f"AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END) "
        f"/ NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs['start_date']}' AND '{hs['end_date']}' "
        f"AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)"
    )


def _baseline_pct_dissatisfied(col, baseline_season, valid_labels):
    return (
        f"ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END) "
        f"/ NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)"
    )


def _current_pct_dissatisfied(col, current_season, valid_labels):
    return (
        f"ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END) "
        f"/ NULLIF(SUM(CASE WHEN SEASON = {current_season} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)"
    )


def _hs_pct_dissatisfied(col, hs, valid_labels):
    return (
        f"ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs['start_date']}' AND '{hs['end_date']}' "
        f"AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END) "
        f"/ NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs['start_date']}' AND '{hs['end_date']}' "
        f"AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)"
    )


def _round2(val):
    return round(float(val), 2) if val is not None else None


def _round1(val):
    return round(float(val), 1) if val is not None else None


def _float(val):
    return float(val) if val is not None else None


# ═══════════════════════════════════════════════════════════════════════
# QUERY BUILDERS — Slides 1-6
# ═══════════════════════════════════════════════════════════════════════

def query_core_ratings(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 1: Core Satisfaction Ratings (0-10 scale AVG)."""
    voc_metrics = [
        ("Overall Experience", "OVERALL_NUMRAT"),
        ("Concessions", "CONCESS_NUMRAT"),
        ("Entertainment", "ENTERTAIN_NUMRAT"),
        ("Parking Ingress", "PARKING_NUMRAT"),
    ]
    rows = []
    for label, col in voc_metrics:
        sql = f"""
        SELECT
            {_baseline_avg(col, baseline_season)},
            AVG(CASE WHEN SEASON = {current_season} AND {col} < 80 THEN {col} END) AS s_current,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND {col} < 80 THEN {col} END) AS hs_prev,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND {col} < 80 THEN {col} END) AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round2(r[0])
        s_cur = _round2(r[1])
        hp = _round2(r[2])
        hc = _round2(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, False, False))

    # Gate Entry from VOC_VIEW
    for label, col in [("Gate Entry", "GE_NUMRAT")]:
        sql = f"""
        SELECT
            {_baseline_avg(col, baseline_season)},
            AVG(CASE WHEN SEASON = {current_season} AND {col} < 80 THEN {col} END) AS s_current,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND {col} < 80 THEN {col} END) AS hs_prev,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND {col} < 80 THEN {col} END) AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round2(r[0])
        s_cur = _round2(r[1])
        hp = _round2(r[2])
        hc = _round2(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, False, False))

    # Metrics from PARKING_VIEW — Seatview, Merchandise, Parking Egress
    parking_metrics = [
        ("Seatview", "SEATVIEW_NUMRAT"),
        ("Merchandise", "MERCH_NUMRAT"),
        ("Parking Egress", "PARKING_EXIT_NUMRAT"),
    ]
    for label, col in parking_metrics:
        sql = f"""
        SELECT
            {_baseline_avg(col, baseline_season)},
            AVG(CASE WHEN SEASON = {current_season} AND {col} < 80 THEN {col} END) AS s_current,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND {col} < 80 THEN {col} END) AS hs_prev,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND {col} < 80 THEN {col} END) AS hs_curr
        FROM {PARKING_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round2(r[0])
        s_cur = _round2(r[1])
        hp = _round2(r[2])
        hc = _round2(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, False, False))

    # Staff overall (normalized from 1-5 to 0-10 for 2024)
    sql = f"""
    SELECT
        AVG(CASE WHEN SEASON = {baseline_season} AND col_val BETWEEN 1 AND 5
                  THEN (col_val - 1) * 2.5 END) AS s_base,
        AVG(CASE WHEN SEASON = {current_season} AND col_val < 80
                  THEN col_val END) AS s_current,
        AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                      AND col_val < 80
                  THEN col_val END) AS hs_prev,
        AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                      AND col_val < 80
                  THEN col_val END) AS hs_curr
    FROM (
        SELECT SEASON, GAME_DATE,
            CASE WHEN SEASON = {baseline_season} THEN TB_ADDON_4_1
                 ELSE TB_ADDON_23_1 END AS col_val
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        UNION ALL
        SELECT SEASON, GAME_DATE,
            CASE WHEN SEASON = {baseline_season} THEN TB_ADDON_4_3
                 ELSE TB_ADDON_23_3 END AS col_val
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        UNION ALL
        SELECT SEASON, GAME_DATE,
            CASE WHEN SEASON = {baseline_season} THEN TB_ADDON_4_5
                 ELSE TB_ADDON_23_5 END AS col_val
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        UNION ALL
        SELECT SEASON, GAME_DATE,
            CASE WHEN SEASON = {baseline_season} THEN TB_ADDON_4_6
                 ELSE TB_ADDON_23_6 END AS col_val
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        UNION ALL
        SELECT SEASON, GAME_DATE,
            CASE WHEN SEASON = {baseline_season} THEN TB_ADDON_4_11
                 ELSE TB_ADDON_23_60 END AS col_val
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        UNION ALL
        SELECT SEASON, GAME_DATE,
            CASE WHEN SEASON = {baseline_season} THEN TB_ADDON_4_38
                 ELSE TB_ADDON_23_38 END AS col_val
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        UNION ALL
        SELECT SEASON, GAME_DATE,
            CASE WHEN SEASON = {baseline_season} THEN TB_ADDON_4_10
                 ELSE TB_ADDON_23_61 END AS col_val
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
    )
    WHERE SEASON = {baseline_season} AND col_val BETWEEN 1 AND 5
       OR SEASON = {current_season} AND col_val < 80
    """
    r = run_query(conn, sql)[0]
    s_base = _round2(r[0])
    s_cur = _round2(r[1])
    hp = _round2(r[2])
    hc = _round2(r[3])
    if s_cur is not None:
        rows.append(("Staff", s_base, s_cur, hp, hc, False, False))

    return ("Core Satisfaction Ratings", "AVG Score (0-10 Scale) | Tropicana Field", rows)


def query_staff_ratings(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 2: Staff Satisfaction (normalized 0-10)."""
    staff_pairs = [
        ("Parking Staff",                 "_1",  "_1"),
        ("Fan Host Usher",                "_3",  "_3"),
        ("Concessions Staff",             "_5",  "_5"),
        ("Retail Staff",                  "_6",  "_6"),
        ("Security",                      "_11", "_60"),
        ("Tech Team",                     "_38", "_38"),
        ("Fan Host Ticket Taker/Scanner", "_10", "_61"),
    ]
    rows = []
    for label, s4, s23 in staff_pairs:
        col_old = f"TB_ADDON_4{s4}"
        col_new = f"TB_ADDON_23{s23}"
        sql = f"""
        SELECT
            AVG(CASE WHEN SEASON = {baseline_season} AND {col_old} BETWEEN 1 AND 5
                      THEN ({col_old} - 1) * 2.5 END) AS s_base,
            AVG(CASE WHEN SEASON = {current_season} AND {col_new} < 80
                      THEN {col_new} END) AS s_current,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                          AND {col_new} < 80
                      THEN {col_new} END) AS hs_prev,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                          AND {col_new} < 80
                      THEN {col_new} END) AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round2(r[0])
        s_cur = _round2(r[1])
        hp = _round2(r[2])
        hc = _round2(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, False, False))
    return ("Staff Satisfaction Ratings", "AVG Score (0-10 Scale, Normalized) | Tropicana Field", rows)


def query_concessions_grid(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 3: Concessions Grid (% Highly Satisfied)."""
    metrics = [
        ("Value",            "CONCESS_GRID_VALUE_DESC"),
        ("Customer Service", "CONCESS_GRID_CUSTSERV_DESC"),
        ("Selection",        "CONCESS_GRID_SELECTION_DESC"),
        ("Cleanliness",      "CONCESS_GRID_CLEAN_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    rows = []
    for label, col in metrics:
        sql = f"""
        SELECT
            {_baseline_pct_highly(col, baseline_season, valid_labels)},
            {_current_pct_highly(col, current_season, valid_labels)} AS s_current,
            {_hs_pct_highly(col, hs_prev, valid_labels)} AS hs_prev,
            {_hs_pct_highly(col, hs_curr, valid_labels)} AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round1(r[0])
        s_cur = _round1(r[1])
        hp = _round1(r[2])
        hc = _round1(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, True, False))
    return ("Concessions Grid", "% Highly Satisfied | Tropicana Field", rows)


def query_concessions_grid_dissatisfied(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 3 (middle): Concessions Grid (% Highly Dissatisfied)."""
    metrics = [
        ("Value",            "CONCESS_GRID_VALUE_DESC"),
        ("Customer Service", "CONCESS_GRID_CUSTSERV_DESC"),
        ("Selection",        "CONCESS_GRID_SELECTION_DESC"),
        ("Cleanliness",      "CONCESS_GRID_CLEAN_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    rows = []
    for label, col in metrics:
        sql = f"""
        SELECT
            {_baseline_pct_dissatisfied(col, baseline_season, valid_labels)},
            {_current_pct_dissatisfied(col, current_season, valid_labels)} AS s_current,
            {_hs_pct_dissatisfied(col, hs_prev, valid_labels)} AS hs_prev,
            {_hs_pct_dissatisfied(col, hs_curr, valid_labels)} AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round1(r[0])
        s_cur = _round1(r[1])
        hp = _round1(r[2])
        hc = _round1(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, True, True))  # invert=True (lower is better)
    return ("Concessions Grid", "% Highly Dissatisfied | Tropicana Field", rows)


def query_concessions_grid_excl_premium(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Concessions Grid (% Highly Satisfied) excluding Premium sections."""
    metrics = [
        ("Value",            "CONCESS_GRID_VALUE_DESC"),
        ("Customer Service", "CONCESS_GRID_CUSTSERV_DESC"),
        ("Selection",        "CONCESS_GRID_SELECTION_DESC"),
        ("Cleanliness",      "CONCESS_GRID_CLEAN_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    rows = []
    for label, col in metrics:
        sql = f"""
        SELECT
            ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND NOT {PREMIUM_FILTER} AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND NOT {PREMIUM_FILTER} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND NOT {PREMIUM_FILTER} AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND NOT {PREMIUM_FILTER} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND NOT {PREMIUM_FILTER} AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND NOT {PREMIUM_FILTER} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND NOT {PREMIUM_FILTER} AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND NOT {PREMIUM_FILTER} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round1(r[0])
        s_cur = _round1(r[1])
        hp = _round1(r[2])
        hc = _round1(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, True, False))
    return ("Concessions Grid (Excl. Premium)", "% Highly Satisfied | Tropicana Field (Premium Excluded)", rows)


def query_concessions_grid_dissatisfied_excl_premium(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Concessions Grid (% Highly Dissatisfied) excluding Premium sections."""
    metrics = [
        ("Value",            "CONCESS_GRID_VALUE_DESC"),
        ("Customer Service", "CONCESS_GRID_CUSTSERV_DESC"),
        ("Selection",        "CONCESS_GRID_SELECTION_DESC"),
        ("Cleanliness",      "CONCESS_GRID_CLEAN_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    rows = []
    for label, col in metrics:
        sql = f"""
        SELECT
            ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND NOT {PREMIUM_FILTER} AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND NOT {PREMIUM_FILTER} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND NOT {PREMIUM_FILTER} AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND NOT {PREMIUM_FILTER} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND NOT {PREMIUM_FILTER} AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND NOT {PREMIUM_FILTER} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND NOT {PREMIUM_FILTER} AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND NOT {PREMIUM_FILTER} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round1(r[0])
        s_cur = _round1(r[1])
        hp = _round1(r[2])
        hc = _round1(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, True, True))
    return ("Concessions Grid (Excl. Premium)", "% Highly Dissatisfied | Tropicana Field (Premium Excluded)", rows)


def query_concessions_wait_excl_premium(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Concessions Wait-Time Expectation Extremes excluding Premium sections."""
    rows = []
    valid_wait = "('Much more than what I expected','Slightly more than what I expected','About what I expected','Slightly less what I expected','Much less than what I expected')"
    for label_name, text_val, invert in [
        ("% Much Less Wait Than Expected", "Much less than what I expected", False),
        ("% Much More Wait Than Expected", "Much more than what I expected", True),
    ]:
        sql = f"""
        SELECT
            ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND NOT {PREMIUM_FILTER} AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND NOT {PREMIUM_FILTER} AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND NOT {PREMIUM_FILTER} AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND NOT {PREMIUM_FILTER} AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND NOT {PREMIUM_FILTER} AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND NOT {PREMIUM_FILTER} AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND NOT {PREMIUM_FILTER} AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND NOT {PREMIUM_FILTER} AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1)
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        r = run_query(conn, sql)[0]
        s_base = _float(r[0])
        s_cur = _float(r[1])
        hp = _float(r[2])
        hc = _float(r[3])
        if s_cur is not None or hp is not None or hc is not None:
            rows.append((label_name, s_base, s_cur, hp, hc, True, invert))
    return rows


def query_concessions_wait(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 3 (lower): Concessions Wait-Time Expectation Extremes."""
    rows = []
    valid_wait = "('Much more than what I expected','Slightly more than what I expected','About what I expected','Slightly less what I expected','Much less than what I expected')"
    for label_name, text_val, invert in [
        ("% Much Less Wait Than Expected", "Much less than what I expected", False),
        ("% Much More Wait Than Expected", "Much more than what I expected", True),
    ]:
        sql = f"""
        SELECT
            ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1)
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        r = run_query(conn, sql)[0]
        s_base = _float(r[0])
        s_cur = _float(r[1])
        hp = _float(r[2])
        hc = _float(r[3])
        if s_cur is not None or hp is not None or hc is not None:
            rows.append((label_name, s_base, s_cur, hp, hc, True, invert))
    return rows


def query_food_quality(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 4: Food Quality (% Highly Satisfied)."""
    food_items = [
        ("Hotdog",                "CONCESS_QUALITY_HOTDOG_DESC"),
        ("Burgers",               "CONCESS_QUALITY_BURGERS_DESC"),
        ("Chicken",               "CONCESS_QUALITY_CHICKEN_DESC"),
        ("Pizza",                 "CONCESS_QUALITY_PIZZA_DESC"),
        ("Nachos",                "CONCESS_QUALITY_NACHOS_DESC"),
        ("Fries",                 "CONCESS_QUALITY_FRIES_DESC"),
        ("Alcohol",               "CONCESS_QUALITY_ALCOHOL_DESC"),
        ("Non-Alcohol Beverages", "CONCESS_QUALITY_NONALCOHOL_DESC"),
        ("Ice Cream",             "CONCESS_QUALITY_ICECREAM_DESC"),
        ("Nuts",                  "CONCESS_QUALITY_NUTS_DESC"),
        ("Popcorn",               "CONCESS_QUALITY_POPCORN_DESC"),
        ("Pretzels",              "CONCESS_QUALITY_PRETZELS_DESC"),
        ("Other Dessert",         "CONCESS_QUALITY_OTHER_DESSERT_DESC"),
        ("Other Entree",          "CONCESS_QUALITY_OTHER_ENTREE_DESC"),
        ("Sandwich",              "CONCESS_QUALITY_SANDWICH_DESC"),
        ("Sausage",               "CONCESS_QUALITY_SAUSAGE_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    rows = []
    for label, col in food_items:
        sql = f"""
        SELECT
            {_baseline_pct_highly(col, baseline_season, valid_labels)},
            {_current_pct_highly(col, current_season, valid_labels)} AS s_current,
            {_hs_pct_highly(col, hs_prev, valid_labels)} AS hs_prev,
            {_hs_pct_highly(col, hs_curr, valid_labels)} AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round1(r[0])
        s_cur = _round1(r[1])
        hp = _round1(r[2])
        hc = _round1(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, True, False))
    return ("Food Quality", "% Highly Satisfied | Tropicana Field", rows)


def query_entertainment_grid(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 5: Entertainment Grid (% Highly Satisfied)."""
    metrics = [
        ("Scoreboard",      "ENTERTAIN_GRID_SCOREBOARD_DESC"),
        ("Music",           "ENTERTAIN_GRID_MUSIC_DESC"),
        ("Games/Contests",  "ENTERTAIN_GRID_GAMES_DESC"),
        ("Kids Activities", "ENTERTAIN_GRID_KIDS_ACTIVITIES_DESC"),
        ("Pregame Content", "ENTERTAIN_GRID_PREGAME_CONTENT_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    rows = []
    for label, col in metrics:
        sql = f"""
        SELECT
            {_baseline_pct_highly(col, baseline_season, valid_labels)},
            {_current_pct_highly(col, current_season, valid_labels)} AS s_current,
            {_hs_pct_highly(col, hs_prev, valid_labels)} AS hs_prev,
            {_hs_pct_highly(col, hs_curr, valid_labels)} AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round1(r[0])
        s_cur = _round1(r[1])
        hp = _round1(r[2])
        hc = _round1(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, True, False))
    return ("Entertainment Grid", "% Highly Satisfied | Tropicana Field", rows)


def query_gate_entry(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Gate Entry Time Expectation (extremes only)."""
    rows = []
    for label_name, code_val, invert in [
        ("% Less Than Expected", 3, False),
        ("% More Than Expected", 1, True),
    ]:
        sql = f"""
        SELECT
            ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND GE_TIME_EXPECT < 80 AND GE_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND GE_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND GE_TIME_EXPECT < 80 AND GE_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND GE_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND GE_TIME_EXPECT < 80 AND GE_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND GE_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND GE_TIME_EXPECT < 80 AND GE_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND GE_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1)
        FROM {PARKING_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        r = run_query(conn, sql)[0]
        s_base = _float(r[0])
        s_cur = _float(r[1])
        hp = _float(r[2])
        hc = _float(r[3])
        if s_cur is not None or hp is not None or hc is not None:
            rows.append((label_name, s_base, s_cur, hp, hc, True, invert))
    return rows


def query_parking(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Parking Experience — numeric ratings + parking time expectation."""
    rows = []
    for label, col in [
        ("Finding & Parking (AVG 0-10)", "PARKING_NUMRAT"),
        ("Exiting Parking Lot (AVG 0-10)", "PARKING_EXIT_NUMRAT"),
    ]:
        sql = f"""
        SELECT
            {_baseline_avg(col, baseline_season)},
            AVG(CASE WHEN SEASON = {current_season} AND {col} < 80 THEN {col} END),
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND {col} < 80 THEN {col} END),
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND {col} < 80 THEN {col} END)
        FROM {PARKING_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        r = run_query(conn, sql)[0]
        s_base = _round2(r[0])
        s_cur = _round2(r[1])
        hp = _round2(r[2])
        hc = _round2(r[3])
        if s_cur is not None or hp is not None or hc is not None:
            rows.append((label, s_base, s_cur, hp, hc, False, False))

    for label_name, code_val, invert in [
        ("% Less Than Expected", 3, False),
        ("% About What Expected", 2, False),
        ("% More Than Expected", 1, True),
    ]:
        sql = f"""
        SELECT
            ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND PARKING_TIME_EXPECT < 80 AND PARKING_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND PARKING_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND PARKING_TIME_EXPECT < 80 AND PARKING_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND PARKING_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND PARKING_TIME_EXPECT < 80 AND PARKING_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND PARKING_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND PARKING_TIME_EXPECT < 80 AND PARKING_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND PARKING_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1)
        FROM {PARKING_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        r = run_query(conn, sql)[0]
        s_base = _float(r[0])
        s_cur = _float(r[1])
        hp = _float(r[2])
        hc = _float(r[3])
        if s_cur is not None or hp is not None or hc is not None:
            rows.append((label_name, s_base, s_cur, hp, hc, True, invert))
    return rows


def query_seat_value(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Seat View Value: % Good Value."""
    rows = []
    sql = f"""
    SELECT
        ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND SEATVIEW_VALUE < 80 AND SEATVIEW_VALUE = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND SEATVIEW_VALUE < 80 THEN 1 ELSE 0 END), 0), 1),
        ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND SEATVIEW_VALUE < 80 AND SEATVIEW_VALUE = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND SEATVIEW_VALUE < 80 THEN 1 ELSE 0 END), 0), 1),
        ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND SEATVIEW_VALUE < 80 AND SEATVIEW_VALUE = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND SEATVIEW_VALUE < 80 THEN 1 ELSE 0 END), 0), 1),
        ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND SEATVIEW_VALUE < 80 AND SEATVIEW_VALUE = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND SEATVIEW_VALUE < 80 THEN 1 ELSE 0 END), 0), 1)
    FROM {PARKING_VIEW}
    WHERE STADIUM = 'Tropicana Field'
      AND SEASON IN ({baseline_season}, {current_season})
    """
    r = run_query(conn, sql)[0]
    s_base = _float(r[0])
    s_cur = _float(r[1])
    hp = _float(r[2])
    hc = _float(r[3])
    if s_cur is not None or hp is not None or hc is not None:
        rows.append(("% Good Seat Value", s_base, s_cur, hp, hc, True, False))
    return rows


# ═══════════════════════════════════════════════════════════════════════
# QUERY BUILDERS — Slide 7 & 8 (Segment Analysis)
# ═══════════════════════════════════════════════════════════════════════

PREMIUM_FILTER = "(SECTION_CODE LIKE 'HPBX%' OR SECTION_CODE LIKE 'CHALK%' OR SECTION_CODE LIKE 'BASE%')"
APRIL_DATES = ("2026-04-06", "2026-04-26")
MAY_HS3_DATES = ("2026-05-01", "2026-05-06")
MAY_HS4_DATES = ("2026-05-15", "2026-05-20")


def query_premium_deep_dive(conn):
    """Slide 7: In-Seat Delivery Deep Dive — Premium vs Everyone Else.
    Concessions Grid metrics (% Highly Satisfied) by segment and period."""
    metrics = [
        ("Customer Service", "CONCESS_GRID_CUSTSERV_DESC"),
        ("Selection",        "CONCESS_GRID_SELECTION_DESC"),
        ("Cleanliness",      "CONCESS_GRID_CLEAN_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    periods = [APRIL_DATES, MAY_HS3_DATES, MAY_HS4_DATES]

    rows = []
    for label, col in metrics:
        premium_vals = []
        everyone_vals = []
        for start, end in periods:
            sql = f"""
            SELECT
                ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER} AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                    / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1),
                ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER} AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                    / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)
            FROM {VOC_VIEW}
            WHERE STADIUM = 'Tropicana Field'
              AND SEASON = 2026
              AND GAME_DATE BETWEEN '{start}' AND '{end}'
            """
            r = run_query(conn, sql)[0]
            premium_vals.append(_round1(r[0]))
            everyone_vals.append(_round1(r[1]))
        rows.append((label, premium_vals, everyone_vals))
    return rows


def query_premium_concessions(conn, baseline_season, current_season, hs_prev, hs_curr):
    """New slide: Premium Concessions Grid — Premium vs Everyone Else.
    Returns two lists (satisfied_rows, dissatisfied_rows).
    Each row: (metric, premium_hs3, premium_hs4, nonprem_hs3, nonprem_hs4)
    """
    metrics = [
        ("Value",            "CONCESS_GRID_VALUE_DESC"),
        ("Customer Service", "CONCESS_GRID_CUSTSERV_DESC"),
        ("Selection",        "CONCESS_GRID_SELECTION_DESC"),
        ("Cleanliness",      "CONCESS_GRID_CLEAN_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"

    satisfied_rows = []
    dissatisfied_rows = []

    for label, col in metrics:
        # Query both segments for both satisfaction levels in one query per metric
        sql = f"""
        SELECT
            -- Premium Highly Satisfied: HS3, HS4
            ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS prem_sat_hs3,
            ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS prem_sat_hs4,
            -- Non-Premium Highly Satisfied: HS3, HS4
            ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS nonprem_sat_hs3,
            ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS nonprem_sat_hs4,
            -- Premium Highly Dissatisfied: HS3, HS4
            ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS prem_dis_hs3,
            ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS prem_dis_hs4,
            -- Non-Premium Highly Dissatisfied: HS3, HS4
            ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS nonprem_dis_hs3,
            ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS nonprem_dis_hs4
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON = {current_season}
        """
        r = run_query(conn, sql)[0]
        prem_sat_hs3 = _round1(r[0])
        prem_sat_hs4 = _round1(r[1])
        nonprem_sat_hs3 = _round1(r[2])
        nonprem_sat_hs4 = _round1(r[3])
        prem_dis_hs3 = _round1(r[4])
        prem_dis_hs4 = _round1(r[5])
        nonprem_dis_hs3 = _round1(r[6])
        nonprem_dis_hs4 = _round1(r[7])

        satisfied_rows.append((label, prem_sat_hs3, prem_sat_hs4, nonprem_sat_hs3, nonprem_sat_hs4))
        dissatisfied_rows.append((label, prem_dis_hs3, prem_dis_hs4, nonprem_dis_hs3, nonprem_dis_hs4))

    return satisfied_rows, dissatisfied_rows


def query_wallet_wait(conn):
    """Slide 8: Rays Wallet (Ballpark App) — Concession Wait Expectations.
    % for each wait category, comparing Wallet Users vs Everyone Else by period."""
    valid_wait = "('Much more than what I expected','Slightly more than what I expected','About what I expected','Slightly less what I expected','Much less than what I expected')"
    wait_categories = [
        ("Much Less Than Expected", "Much less than what I expected"),
        ("Slightly Less Than Expected", "Slightly less what I expected"),
        ("About What Expected", "About what I expected"),
        ("Slightly More Than Expected", "Slightly more than what I expected"),
        ("Much More Than Expected", "Much more than what I expected"),
    ]
    periods = [APRIL_DATES, MAY_HS3_DATES, MAY_HS4_DATES]

    rows = []
    for label, text_val in wait_categories:
        wallet_vals = []
        everyone_vals = []
        for start, end in periods:
            sql = f"""
            SELECT
                ROUND(100.0 * SUM(CASE WHEN TB_ADDON_26 = 1 AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                    / NULLIF(SUM(CASE WHEN TB_ADDON_26 = 1 AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
                ROUND(100.0 * SUM(CASE WHEN TB_ADDON_26 IN (2, 3) AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                    / NULLIF(SUM(CASE WHEN TB_ADDON_26 IN (2, 3) AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1)
            FROM {VOC_VIEW}
            WHERE STADIUM = 'Tropicana Field'
              AND SEASON = 2026
              AND GAME_DATE BETWEEN '{start}' AND '{end}'
            """
            r = run_query(conn, sql)[0]
            wallet_vals.append(_round1(r[0]))
            everyone_vals.append(_round1(r[1]))
        rows.append((label, wallet_vals, everyone_vals))
    return rows


def query_ballpark_app_usage(conn):
    """Ballpark App Usage: % breakdown of Used/Aware/Unaware + AVG OVERALL_NUMRAT per category per period.
    TB_ADDON_26: 1=Used, 2=Aware but didn't use, 3=Unaware.
    Returns list of (category, april_pct, hs3_pct, hs4_pct, april_score, hs3_score, hs4_score)
    """
    categories = [
        ("Used Rays Wallet", 1),
        ("Aware, Did Not Use", 2),
        ("Unaware", 3),
    ]
    periods = [APRIL_DATES, MAY_HS3_DATES, MAY_HS4_DATES]

    rows = []
    for label, code_val in categories:
        pcts = []
        scores = []
        for start, end in periods:
            sql = f"""
            SELECT
                ROUND(100.0 * SUM(CASE WHEN TB_ADDON_26 = {code_val} THEN 1 ELSE 0 END)
                    / NULLIF(SUM(CASE WHEN TB_ADDON_26 IN (1, 2, 3) THEN 1 ELSE 0 END), 0), 1) AS pct,
                AVG(CASE WHEN TB_ADDON_26 = {code_val} AND OVERALL_NUMRAT < 80 THEN OVERALL_NUMRAT END) AS avg_score
            FROM {VOC_VIEW}
            WHERE STADIUM = 'Tropicana Field'
              AND SEASON = 2026
              AND GAME_DATE BETWEEN '{start}' AND '{end}'
            """
            r = run_query(conn, sql)[0]
            pcts.append(_round1(r[0]))
            scores.append(_round2(r[1]))
        rows.append((label, pcts, scores))
    return rows


FEEDBACK_TABLE = "TBRDP_DW_DEV.IM_RPT.T_OVERALL_FEEDBACK_SENTENCE_LEVEL"


def query_parking_negative_top3(conn, hs_curr):
    """Top 3 negative parking complaint categories for HS4 with percentages."""
    sql = f"""
    SELECT AI_CATEGORY, COUNT(*) as cnt
    FROM {FEEDBACK_TABLE}
    WHERE GAME_DATE::DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
      AND SENTIMENT_CATEGORY = 'Negative'
      AND AI_CATEGORY IN ('Parking Availability', 'Departure Traffic', 'Parking Cost', 'Parking Staff')
    GROUP BY AI_CATEGORY
    ORDER BY cnt DESC
    LIMIT 3
    """
    results = run_query(conn, sql)
    total_sql = f"""
    SELECT COUNT(*)
    FROM {FEEDBACK_TABLE}
    WHERE GAME_DATE::DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
      AND SENTIMENT_CATEGORY = 'Negative'
      AND AI_CATEGORY IN ('Parking Availability', 'Departure Traffic', 'Parking Cost', 'Parking Staff')
    """
    total = run_query(conn, total_sql)[0][0]
    rows = []
    for cat, cnt in results:
        pct = round(cnt / total * 100, 1) if total > 0 else 0
        rows.append((cat, cnt, pct))
    return rows, total


def query_departure_traffic_subcategories(conn, hs_curr):
    """Top 3 subcategories within Departure Traffic negative feedback.
    Uses AI to subcategorize the departure traffic complaints."""
    sql = f"""
    SELECT SENTENCE_TEXT
    FROM {FEEDBACK_TABLE}
    WHERE GAME_DATE::DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
      AND SENTIMENT_CATEGORY = 'Negative'
      AND AI_CATEGORY = 'Departure Traffic'
    """
    results = run_query(conn, sql)
    sentences = [r[0] for r in results]
    total = len(sentences)

    # Manually categorize based on keyword analysis
    categories = {
        "Exit Traffic / Long Wait Times": 0,
        "Lack of Traffic Direction / Coordination": 0,
        "Specific Lot / Route Issues": 0,
    }
    for s in sentences:
        s_lower = s.lower()
        if any(kw in s_lower for kw in ['police', 'direct', 'coordinat', 'attendant', 'chaos', 'order', 'control']):
            categories["Lack of Traffic Direction / Coordination"] += 1
        elif any(kw in s_lower for kw in ['lot 2', 'lot 7', 'handicap', '1st ave', 'i275', 'route', 'gate', 'locked']):
            categories["Specific Lot / Route Issues"] += 1
        else:
            categories["Exit Traffic / Long Wait Times"] += 1

    rows = sorted(categories.items(), key=lambda x: -x[1])[:3]
    result = [(cat, cnt, round(cnt / total * 100, 1) if total > 0 else 0) for cat, cnt in rows]
    return result, total


def query_entertainment_negative_top3(conn, hs_curr):
    """Top 3 negative sentiment themes for Games/Contests (In Game Entertainment) in HS4."""
    sql = f"""
    SELECT SENTENCE_TEXT
    FROM {FEEDBACK_TABLE}
    WHERE GAME_DATE::DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
      AND SENTIMENT_CATEGORY = 'Negative'
      AND AI_CATEGORY = 'In Game Entertainment'
    """
    results = run_query(conn, sql)
    sentences = [r[0] for r in results]
    total = len(sentences)

    # Categorize based on content analysis
    categories = {
        '"Tarps Off" / Shirtless Fan Disruption': 0,
        "Sky Cam / Wire Camera Distraction": 0,
        "Lack of Fan Interaction / Activities": 0,
    }
    other = 0
    for s in sentences:
        s_lower = s.lower()
        if any(kw in s_lower for kw in ['tarps', 'shirtless', 'shirts off', 'shirt off', 'dude', 'half the game', 'not family', 'appropriate']):
            categories['"Tarps Off" / Shirtless Fan Disruption'] += 1
        elif any(kw in s_lower for kw in ['cam', 'camera', 'wire', 'sky cam']):
            categories["Sky Cam / Wire Camera Distraction"] += 1
        elif any(kw in s_lower for kw in ['interact', 'engag', 'more', 'lack', 'missing', 'could be']):
            categories["Lack of Fan Interaction / Activities"] += 1
        else:
            other += 1

    rows = sorted(categories.items(), key=lambda x: -x[1])[:3]
    result = [(cat, cnt, round(cnt / total * 100, 1) if total > 0 else 0) for cat, cnt in rows]
    return result, total


def query_premium_concessions_v2(conn, current_season, hs_prev, hs_curr):
    """Premium Concessions Grid v2: includes '26 AVG, adds Speed of Service.
    Returns two lists (satisfied_rows, dissatisfied_rows).
    Each row: (metric, s26_avg, prem_hs3, prem_hs4, nonprem_hs3, nonprem_hs4)
    """
    desc_metrics = [
        ("Value",            "CONCESS_GRID_VALUE_DESC"),
        ("Customer Service", "CONCESS_GRID_CUSTSERV_DESC"),
        ("Selection",        "CONCESS_GRID_SELECTION_DESC"),
        ("Cleanliness",      "CONCESS_GRID_CLEAN_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"

    satisfied_rows = []
    dissatisfied_rows = []

    for label, col in desc_metrics:
        sql = f"""
        SELECT
            -- '26 AVG (all respondents)
            ROUND(100.0 * SUM(CASE WHEN {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS s26_avg,
            -- Premium HS3, HS4
            ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS prem_sat_hs3,
            ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS prem_sat_hs4,
            -- Non-Premium HS3, HS4
            ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS nonprem_sat_hs3,
            ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS nonprem_sat_hs4,
            -- Dissatisfied: '26 AVG, Premium HS3/HS4, Non-Premium HS3/HS4
            ROUND(100.0 * SUM(CASE WHEN {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS s26_avg_dis,
            ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS prem_dis_hs3,
            ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS prem_dis_hs4,
            ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS nonprem_dis_hs3,
            ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} = 'Highly dissatisfied' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER}
                AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1) AS nonprem_dis_hs4
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field' AND SEASON = {current_season}
        """
        r = run_query(conn, sql)[0]
        satisfied_rows.append((label, _round1(r[0]), _round1(r[1]), _round1(r[2]), _round1(r[3]), _round1(r[4])))
        dissatisfied_rows.append((label, _round1(r[5]), _round1(r[6]), _round1(r[7]), _round1(r[8]), _round1(r[9])))

    # Speed of Service (numeric: 1=Highly satisfied, 4=Highly dissatisfied)
    speed_sql = f"""
    SELECT
        ROUND(100.0 * SUM(CASE WHEN CONCESS_GRID_SPEED = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN CONCESS_GRID_SPEED BETWEEN 1 AND 4 THEN 1 ELSE 0 END), 0), 1) AS s26_avg,
        ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
            AND CONCESS_GRID_SPEED = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
            AND CONCESS_GRID_SPEED BETWEEN 1 AND 4 THEN 1 ELSE 0 END), 0), 1) AS prem_hs3,
        ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
            AND CONCESS_GRID_SPEED = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
            AND CONCESS_GRID_SPEED BETWEEN 1 AND 4 THEN 1 ELSE 0 END), 0), 1) AS prem_hs4,
        ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
            AND CONCESS_GRID_SPEED = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
            AND CONCESS_GRID_SPEED BETWEEN 1 AND 4 THEN 1 ELSE 0 END), 0), 1) AS nonprem_hs3,
        ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
            AND CONCESS_GRID_SPEED = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
            AND CONCESS_GRID_SPEED BETWEEN 1 AND 4 THEN 1 ELSE 0 END), 0), 1) AS nonprem_hs4,
        -- Dissatisfied (code 4)
        ROUND(100.0 * SUM(CASE WHEN CONCESS_GRID_SPEED = 4 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN CONCESS_GRID_SPEED BETWEEN 1 AND 4 THEN 1 ELSE 0 END), 0), 1) AS s26_avg_dis,
        ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
            AND CONCESS_GRID_SPEED = 4 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
            AND CONCESS_GRID_SPEED BETWEEN 1 AND 4 THEN 1 ELSE 0 END), 0), 1) AS prem_dis_hs3,
        ROUND(100.0 * SUM(CASE WHEN {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
            AND CONCESS_GRID_SPEED = 4 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
            AND CONCESS_GRID_SPEED BETWEEN 1 AND 4 THEN 1 ELSE 0 END), 0), 1) AS prem_dis_hs4,
        ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
            AND CONCESS_GRID_SPEED = 4 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
            AND CONCESS_GRID_SPEED BETWEEN 1 AND 4 THEN 1 ELSE 0 END), 0), 1) AS nonprem_dis_hs3,
        ROUND(100.0 * SUM(CASE WHEN NOT {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
            AND CONCESS_GRID_SPEED = 4 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN NOT {PREMIUM_FILTER}
            AND GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
            AND CONCESS_GRID_SPEED BETWEEN 1 AND 4 THEN 1 ELSE 0 END), 0), 1) AS nonprem_dis_hs4
    FROM {VOC_VIEW}
    WHERE STADIUM = 'Tropicana Field' AND SEASON = {current_season}
    """
    r = run_query(conn, speed_sql)[0]
    satisfied_rows.append(("Speed of Service", _round1(r[0]), _round1(r[1]), _round1(r[2]), _round1(r[3]), _round1(r[4])))
    dissatisfied_rows.append(("Speed of Service", _round1(r[5]), _round1(r[6]), _round1(r[7]), _round1(r[8]), _round1(r[9])))

    return satisfied_rows, dissatisfied_rows

def pct_diff(a, b):
    if a is None or b is None or a == 0:
        return None
    return round((b - a) / a * 100, 1)


def fmt_val(val, is_pct):
    if val is None:
        return "N/A"
    return f"{val:.1f}%" if is_pct else f"{val:.2f}"


def fmt_diff(diff):
    if diff is None:
        return "N/A"
    sign = "+" if diff > 0 else ""
    return f"{sign}{diff:.1f}%"


# ═══════════════════════════════════════════════════════════════════════
# PPTX BUILDER UTILITIES
# ═══════════════════════════════════════════════════════════════════════

def set_cell_text(cell, text, font_size=10, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_cell_fill(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def add_rect(slide, prs, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=10, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    run.font.italic = italic
    return tb


def add_branded_header(slide, prs, title_text):
    add_rect(slide, prs, 0, 0, prs.slide_width, Inches(0.95), NAVY)
    add_rect(slide, prs, 0, Inches(0.95), prs.slide_width, Inches(0.06), YELLOW)
    add_text(slide, Inches(0.45), Inches(0.10), Inches(11), Inches(0.30),
             "VOC POST-ATTENDANCE SURVEY  \u00b7  HOMESTAND COMPARISON",
             size=10, bold=True, color=SKY)
    add_text(slide, Inches(0.45), Inches(0.36), Inches(11), Inches(0.55),
             title_text, size=22, bold=True, color=WHITE)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(12.35), Inches(0.18), height=Inches(0.65))


def add_branded_footer(slide, prs, page_num, total_pages):
    add_rect(slide, prs, 0, Inches(7.20), prs.slide_width, Inches(0.30), NAVY)
    add_text(slide, Inches(0.45), Inches(7.23), Inches(8), Inches(0.25),
             "Tampa Bay Rays  \u00b7  Strategy & Analytics",
             size=9, color=SKY)
    add_text(slide, Inches(10.0), Inches(7.23), Inches(3.0), Inches(0.25),
             f"Source: Qualtrics VOC  |  {page_num} / {total_pages}",
             size=9, color=SKY, align=PP_ALIGN.RIGHT)


def build_table(slide, data_rows, top, left, table_width, col_headers, col_widths,
                row_height_override=None):
    """Build a branded 7-column data table (no MLB columns).
    data_rows: list of (metric, s_base, s_cur, hp, hc, is_pct, invert_diff)
    """
    n_cols = len(col_headers)
    n_rows = len(data_rows) + 1

    if row_height_override:
        row_height = row_height_override
    else:
        row_height = Inches(0.38) if len(data_rows) <= 8 else Inches(0.32)
    header_height = Inches(0.45)
    table_height = header_height + row_height * len(data_rows)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, int(table_height))
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = int(w)

    # Header row
    table.rows[0].height = int(header_height)
    for ci, hdr in enumerate(col_headers):
        cell = table.cell(0, ci)
        set_cell_fill(cell, MEDIUM_BLUE)
        set_cell_text(cell, hdr, font_size=9, bold=True, color=WHITE)

    # Data rows
    for ri, row_data in enumerate(data_rows):
        metric, s_base, s_cur, hp, hc, is_pct, invert_diff = row_data
        row_idx = ri + 1
        table.rows[row_idx].height = int(row_height)

        hs_diff = pct_diff(hp, hc)
        baseline_diff = pct_diff(s_base, s_cur)

        if invert_diff:
            if hs_diff is not None:
                hs_diff = -hs_diff
            if baseline_diff is not None:
                baseline_diff = -baseline_diff

        bg_color = LIGHT_BG if ri % 2 == 0 else WHITE

        # 7-column: Metric | '24 Score | '26 AVG | HS3 | HS4 | HS4vsHS3 | '26vs'24
        cell_values = [
            (metric, 9, False, DARK_GRAY, PP_ALIGN.LEFT, None),
            (fmt_val(s_base, is_pct), 9, False, DARK_GRAY, PP_ALIGN.CENTER, None),
            (fmt_val(s_cur, is_pct), 9, False, DARK_GRAY, PP_ALIGN.CENTER, None),
            (fmt_val(hp, is_pct), 9, False, DARK_GRAY, PP_ALIGN.CENTER, None),
            (fmt_val(hc, is_pct), 9, False, DARK_GRAY, PP_ALIGN.CENTER, None),
            (None, 9, True, None, PP_ALIGN.CENTER, hs_diff),
            (None, 9, True, None, PP_ALIGN.CENTER, baseline_diff),
        ]

        for ci, (text, fsize, is_bold, fcolor, align, diff_val) in enumerate(cell_values):
            cell = table.cell(row_idx, ci)
            set_cell_fill(cell, bg_color)
            if diff_val is not None:
                dc = GREEN if diff_val > 0 else (RED if diff_val < 0 else DARK_GRAY)
                set_cell_text(cell, fmt_diff(diff_val), font_size=fsize, bold=True, color=dc)
            elif text is not None:
                set_cell_text(cell, text, font_size=fsize, bold=is_bold, color=fcolor, alignment=align)
            else:
                set_cell_text(cell, "N/A", font_size=fsize, color=MUTED_GRAY)

    return table_height


def build_premium_comparison_table(slide, data_rows, top, left, table_width, hs_prev_num, hs_curr_num,
                                    row_height_override=None):
    """Build a Premium vs Non-Premium comparison table (v2).
    data_rows: list of (metric, s26_avg, prem_hs3, prem_hs4, nonprem_hs3, nonprem_hs4)
    Columns: Metric | '26 AVG | Premium HS3 | Premium HS4 | Non-Prem HS3 | Non-Prem HS4 | HS4vsHS3 | HS4vs'26AVG
    """
    col_headers = [
        "Metric",
        "'26\nAVG",
        f"Premium\nHS{hs_prev_num}", f"Premium\nHS{hs_curr_num}",
        f"Non-Prem\nHS{hs_prev_num}", f"Non-Prem\nHS{hs_curr_num}",
        f"HS{hs_curr_num} vs HS{hs_prev_num}\n% Diff",
        f"HS{hs_curr_num} vs '26\n% Diff"
    ]
    col_widths = [
        Inches(2.0), Inches(1.1), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.7), Inches(1.7)
    ]
    n_cols = 8
    n_rows = len(data_rows) + 1

    row_height = row_height_override or Inches(0.35)
    header_height = Inches(0.45)
    table_height = header_height + row_height * len(data_rows)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, int(table_height))
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = int(w)

    # Header row
    table.rows[0].height = int(header_height)
    for ci, hdr in enumerate(col_headers):
        cell = table.cell(0, ci)
        set_cell_fill(cell, MEDIUM_BLUE)
        set_cell_text(cell, hdr, font_size=8, bold=True, color=WHITE)

    # Data rows
    for ri, (metric, s26_avg, prem_hs3, prem_hs4, nonprem_hs3, nonprem_hs4) in enumerate(data_rows):
        row_idx = ri + 1
        table.rows[row_idx].height = int(row_height)
        bg_color = LIGHT_BG if ri % 2 == 0 else WHITE

        # Compute relative % differences for Premium segment
        diff_hs4_vs_hs3 = pct_diff(prem_hs3, prem_hs4)
        diff_hs4_vs_avg = pct_diff(s26_avg, prem_hs4)

        cell_values = [
            (metric, PP_ALIGN.LEFT, None),
            (f"{s26_avg:.1f}%" if s26_avg is not None else "N/A", PP_ALIGN.CENTER, None),
            (f"{prem_hs3:.1f}%" if prem_hs3 is not None else "N/A", PP_ALIGN.CENTER, None),
            (f"{prem_hs4:.1f}%" if prem_hs4 is not None else "N/A", PP_ALIGN.CENTER, None),
            (f"{nonprem_hs3:.1f}%" if nonprem_hs3 is not None else "N/A", PP_ALIGN.CENTER, None),
            (f"{nonprem_hs4:.1f}%" if nonprem_hs4 is not None else "N/A", PP_ALIGN.CENTER, None),
            (None, PP_ALIGN.CENTER, diff_hs4_vs_hs3),
            (None, PP_ALIGN.CENTER, diff_hs4_vs_avg),
        ]

        for ci, (text, align, diff_val) in enumerate(cell_values):
            cell = table.cell(row_idx, ci)
            set_cell_fill(cell, bg_color)
            if diff_val is not None:
                sign = "+" if diff_val > 0 else ""
                dc = GREEN if diff_val > 0 else (RED if diff_val < 0 else DARK_GRAY)
                set_cell_text(cell, f"{sign}{diff_val:.1f}%", font_size=9, bold=True, color=dc)
            elif text is not None:
                fcolor = DARK_GRAY if text != "N/A" else MUTED_GRAY
                set_cell_text(cell, text, font_size=9, color=fcolor, alignment=align)
            else:
                set_cell_text(cell, "N/A", font_size=9, color=MUTED_GRAY)

    return table_height


def build_segment_table(slide, data_rows, top, left, table_width, col_headers, col_widths,
                        row_height_override=None):
    """Build a segment comparison table for slides 7 & 8.
    data_rows: list of (label, segment1_vals[3], segment2_vals[3])
    where each vals list has [April, May 1-6, May 15-20] percentages.
    """
    n_cols = len(col_headers)
    n_rows = len(data_rows) + 1

    row_height = row_height_override or Inches(0.38)
    header_height = Inches(0.50)
    table_height = header_height + row_height * len(data_rows)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, int(table_height))
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = int(w)

    # Header row
    table.rows[0].height = int(header_height)
    for ci, hdr in enumerate(col_headers):
        cell = table.cell(0, ci)
        set_cell_fill(cell, MEDIUM_BLUE)
        set_cell_text(cell, hdr, font_size=8, bold=True, color=WHITE)

    # Data rows
    for ri, (label, seg1_vals, seg2_vals) in enumerate(data_rows):
        row_idx = ri + 1
        table.rows[row_idx].height = int(row_height)
        bg_color = LIGHT_BG if ri % 2 == 0 else WHITE

        # Metric label
        cell = table.cell(row_idx, 0)
        set_cell_fill(cell, bg_color)
        set_cell_text(cell, label, font_size=9, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

        # Segment 1 values (3 periods)
        for pi, val in enumerate(seg1_vals):
            cell = table.cell(row_idx, 1 + pi)
            set_cell_fill(cell, bg_color)
            text = f"{val:.1f}%" if val is not None else "N/A"
            set_cell_text(cell, text, font_size=9, color=DARK_GRAY if val is not None else MUTED_GRAY)

        # Segment 2 values (3 periods)
        for pi, val in enumerate(seg2_vals):
            cell = table.cell(row_idx, 4 + pi)
            set_cell_fill(cell, bg_color)
            text = f"{val:.1f}%" if val is not None else "N/A"
            set_cell_text(cell, text, font_size=9, color=DARK_GRAY if val is not None else MUTED_GRAY)

    return table_height


# ═══════════════════════════════════════════════════════════════════════
# MAIN BUILD
# ═══════════════════════════════════════════════════════════════════════

def build_report():
    config = load_config()
    hs_prev = get_homestand(config, 2026, 3)
    hs_curr = get_homestand(config, 2026, 4)
    baseline_season = 2024
    current_season = 2026

    print("Connecting to Snowflake (SSO)...")
    conn = get_connection()
    print("Connected. Querying VOC data...")

    # Get response counts
    sql = f"""
    SELECT
        COUNT(CASE WHEN SEASON = {baseline_season} THEN 1 END),
        COUNT(DISTINCT CASE WHEN SEASON = {baseline_season} THEN GAME_DATE END),
        COUNT(CASE WHEN SEASON = {current_season} THEN 1 END),
        COUNT(DISTINCT CASE WHEN SEASON = {current_season} THEN GAME_DATE END)
    FROM {VOC_VIEW}
    WHERE STADIUM = 'Tropicana Field'
      AND SEASON IN ({baseline_season}, {current_season})
    """
    r = run_query(conn, sql)[0]
    base_resp, base_games, cur_resp, cur_games = int(r[0]), int(r[1]), int(r[2]), int(r[3])

    print("  [1/8] Core Satisfaction Ratings...")
    core_data = query_core_ratings(conn, baseline_season, current_season, hs_prev, hs_curr)

    print("  [2/8] Staff Satisfaction Ratings...")
    staff_data = query_staff_ratings(conn, baseline_season, current_season, hs_prev, hs_curr)

    print("  [3/10] Concessions Grid + Wait...")
    concess_grid_data = query_concessions_grid(conn, baseline_season, current_season, hs_prev, hs_curr)
    concess_grid_dissatisfied_data = query_concessions_grid_dissatisfied(conn, baseline_season, current_season, hs_prev, hs_curr)
    concess_wait_rows = query_concessions_wait(conn, baseline_season, current_season, hs_prev, hs_curr)

    print("  [4/10] Concessions Grid (Excl. Premium)...")
    concess_excl_prem_data = query_concessions_grid_excl_premium(conn, baseline_season, current_season, hs_prev, hs_curr)
    concess_excl_prem_dis_data = query_concessions_grid_dissatisfied_excl_premium(conn, baseline_season, current_season, hs_prev, hs_curr)
    concess_wait_excl_prem_rows = query_concessions_wait_excl_premium(conn, baseline_season, current_season, hs_prev, hs_curr)

    print("  [5/10] Premium Concessions Grid...")
    premium_sat_rows, premium_dis_rows = query_premium_concessions_v2(conn, current_season, hs_prev, hs_curr)

    print("  [6/10] Food Quality...")
    food_data = query_food_quality(conn, baseline_season, current_season, hs_prev, hs_curr)

    print("  [7/10] Entertainment Grid...")
    entertain_data = query_entertainment_grid(conn, baseline_season, current_season, hs_prev, hs_curr)

    print("  [8/10] Gate Entry + Parking + Seat Value...")
    gate_rows = query_gate_entry(conn, baseline_season, current_season, hs_prev, hs_curr)
    parking_rows = query_parking(conn, baseline_season, current_season, hs_prev, hs_curr)
    seat_rows = query_seat_value(conn, baseline_season, current_season, hs_prev, hs_curr)

    print("  [9/10] Rays Wallet / Ballpark App Wait Analysis...")
    wallet_rows = query_wallet_wait(conn)

    print("  [10/10] Ballpark App Usage Stats...")
    app_usage_rows = query_ballpark_app_usage(conn)

    print("  [+] Feedback Analysis (Parking, Entertainment)...")
    parking_neg_top3, parking_neg_total = query_parking_negative_top3(conn, hs_curr)
    departure_subcats, departure_total = query_departure_traffic_subcategories(conn, hs_curr)
    entertain_neg_top3, entertain_neg_total = query_entertainment_negative_top3(conn, hs_curr)

    conn.close()
    print("Queries complete. Building PPTX...")

    # Build PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total_pages = 10

    # Context note
    context_note = (
        f"'24 = Full Season ({base_games} games, {base_resp:,} responses)  |  "
        f"'26 = {cur_games} games ({cur_resp:,} responses)  |  "
        f"HS3 = 05/01-05/06 ({hs_prev['games']}g)  |  "
        f"HS4 = 05/15-05/20 ({hs_curr['games']}g)"
    )

    # 7-column headers (no MLB)
    col_headers_7 = [
        "Metric", "'24\nScore", "'26\nAVG", "HS3\nScore", "HS4\nScore",
        "HS4 vs HS3\n% Diff", "'26 vs '24\n% Diff"
    ]
    col_widths_7 = [
        Inches(2.8), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3),
        Inches(2.0), Inches(2.0)
    ]

    # Sort function — by year-over-year % diff (most positive first)
    def sort_by_yoy(row):
        _, s_base, s_cur, _, _, is_pct, invert = row
        d = pct_diff(s_base, s_cur)
        if d is None:
            return float('-inf')
        if invert:
            d = -d
        return -d

    # ─── SLIDE 1: Core Satisfaction Ratings ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Core Satisfaction Ratings")
    _, subtitle, rows = core_data
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             subtitle, size=11, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             context_note, size=8, color=MUTED_GRAY)
    rows_sorted = sorted(rows, key=sort_by_yoy)
    build_table(slide, rows_sorted, Inches(1.65), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7)
    add_branded_footer(slide, prs, 1, total_pages)

    # ─── SLIDE 2: Staff Satisfaction Ratings ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Staff Satisfaction Ratings")
    _, subtitle, rows = staff_data
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             subtitle, size=11, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             context_note, size=8, color=MUTED_GRAY)
    rows_sorted = sorted(rows, key=sort_by_yoy)
    build_table(slide, rows_sorted, Inches(1.65), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7)
    add_branded_footer(slide, prs, 2, total_pages)

    # ─── SLIDE 3: Concessions Grid (Satisfied + Dissatisfied) + Wait ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Concessions Grid + Wait Expectations")
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             "Please rate your overall satisfaction with food and/or beverages.",
             size=11, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             context_note, size=8, color=MUTED_GRAY)

    # Table 1: % Highly Satisfied
    _, _, concess_rows = concess_grid_data
    concess_rows_sorted = sorted(concess_rows, key=sort_by_yoy)
    add_text(slide, Inches(0.45), Inches(1.55), Inches(6), Inches(0.20),
             "% Highly Satisfied", size=9, bold=True, color=DARK_GRAY)
    tbl1_h = build_table(slide, concess_rows_sorted, Inches(1.73), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7, row_height_override=Inches(0.28))

    # Table 2: % Highly Dissatisfied
    _, _, concess_dis_rows = concess_grid_dissatisfied_data
    concess_dis_sorted = sorted(concess_dis_rows, key=sort_by_yoy)
    dis_top = Inches(1.73) + tbl1_h + Inches(0.20)
    add_text(slide, Inches(0.45), dis_top - Inches(0.17), Inches(6), Inches(0.20),
             "% Highly Dissatisfied", size=9, bold=True, color=DARK_GRAY)
    tbl2_h = build_table(slide, concess_dis_sorted, dis_top + Inches(0.03), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7, row_height_override=Inches(0.28))

    # Table 3: Wait expectations
    wait_top = dis_top + Inches(0.03) + tbl2_h + Inches(0.20)
    add_text(slide, Inches(0.45), wait_top - Inches(0.17), Inches(12), Inches(0.20),
             "How did waiting for concessions compare to your expectations?",
             size=9, italic=True, color=GRAY)
    build_table(slide, concess_wait_rows, wait_top + Inches(0.03), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7, row_height_override=Inches(0.28))
    add_branded_footer(slide, prs, 3, total_pages)

    # ─── SLIDE 4: Concessions Grid (Excl. Premium) + Wait ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Concessions Grid (Excl. Premium) + Wait Expectations")
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             "Excluding Premium sections (Homeplate Box, Baseline Premier, Chalk Box)",
             size=11, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             context_note, size=8, color=MUTED_GRAY)

    # Table 1: % Highly Satisfied (Excl. Premium)
    _, _, concess_excl_rows = concess_excl_prem_data
    concess_excl_sorted = sorted(concess_excl_rows, key=sort_by_yoy)
    add_text(slide, Inches(0.45), Inches(1.55), Inches(6), Inches(0.20),
             "% Highly Satisfied", size=9, bold=True, color=DARK_GRAY)
    excl_tbl1_h = build_table(slide, concess_excl_sorted, Inches(1.73), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7, row_height_override=Inches(0.28))

    # Table 2: % Highly Dissatisfied (Excl. Premium)
    _, _, concess_excl_dis_rows = concess_excl_prem_dis_data
    concess_excl_dis_sorted = sorted(concess_excl_dis_rows, key=sort_by_yoy)
    excl_dis_top = Inches(1.73) + excl_tbl1_h + Inches(0.20)
    add_text(slide, Inches(0.45), excl_dis_top - Inches(0.17), Inches(6), Inches(0.20),
             "% Highly Dissatisfied", size=9, bold=True, color=DARK_GRAY)
    excl_tbl2_h = build_table(slide, concess_excl_dis_sorted, excl_dis_top + Inches(0.03), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7, row_height_override=Inches(0.28))

    # Table 3: Wait expectations (Excl. Premium)
    excl_wait_top = excl_dis_top + Inches(0.03) + excl_tbl2_h + Inches(0.20)
    add_text(slide, Inches(0.45), excl_wait_top - Inches(0.17), Inches(12), Inches(0.20),
             "How did waiting for concessions compare to your expectations?",
             size=9, italic=True, color=GRAY)
    build_table(slide, concess_wait_excl_prem_rows, excl_wait_top + Inches(0.03), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7, row_height_override=Inches(0.28))
    add_branded_footer(slide, prs, 4, total_pages)

    # ─── SLIDE 5: Premium Concessions Grid — Premium vs Everyone Else ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Premium Concessions Grid \u2014 Premium vs. Everyone Else")
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             "Premium = Homeplate Box, Baseline Premier, Chalk Box  |  Non-Prem = Everyone Else (Premium excluded)",
             size=9, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.32), Inches(12.5), Inches(0.22),
             context_note, size=8, color=MUTED_GRAY)

    # Table 1: % Highly Satisfied — Premium vs Non-Premium
    add_text(slide, Inches(0.45), Inches(1.53), Inches(6), Inches(0.20),
             "% Highly Satisfied", size=9, bold=True, color=DARK_GRAY)
    prem_tbl1_h = build_premium_comparison_table(slide, premium_sat_rows, Inches(1.72), Inches(0.45),
                                                  Inches(12.3), 3, 4, row_height_override=Inches(0.30))

    # Table 2: % Highly Dissatisfied — Premium vs Non-Premium
    prem_dis_top = Inches(1.72) + prem_tbl1_h + Inches(0.22)
    add_text(slide, Inches(0.45), prem_dis_top - Inches(0.18), Inches(6), Inches(0.20),
             "% Highly Dissatisfied", size=9, bold=True, color=DARK_GRAY)
    build_premium_comparison_table(slide, premium_dis_rows, prem_dis_top + Inches(0.02), Inches(0.45),
                                   Inches(12.3), 3, 4, row_height_override=Inches(0.30))
    add_branded_footer(slide, prs, 5, total_pages)

    # ─── SLIDE 6: Food Quality ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Food Quality (% Highly Satisfied)")
    _, subtitle, rows = food_data
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             "How would you rate the quality (taste, temperature, etc.) of each food / beverage?",
             size=11, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             context_note, size=8, color=MUTED_GRAY)
    rows_sorted = sorted(rows, key=sort_by_yoy)
    build_table(slide, rows_sorted, Inches(1.60), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7, row_height_override=Inches(0.28))
    add_branded_footer(slide, prs, 6, total_pages)

    # ─── SLIDE 7: Entertainment Grid ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Entertainment Grid (% Highly Satisfied)")
    _, subtitle, rows = entertain_data
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             "Please rate your satisfaction with the following:",
             size=11, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             context_note, size=8, color=MUTED_GRAY)
    rows_sorted = sorted(rows, key=sort_by_yoy)
    ent_tbl_h = build_table(slide, rows_sorted, Inches(1.65), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7)

    # Games/Contests Negative Feedback Analysis
    ent_fb_top = Inches(1.65) + ent_tbl_h + Inches(0.25)
    add_text(slide, Inches(0.45), ent_fb_top - Inches(0.20), Inches(12), Inches(0.20),
             f"Games/Contests: Top Negative Feedback Themes (HS4, n={entertain_neg_total})",
             size=9, bold=True, color=DARK_GRAY)

    fb_headers = ["#", "Theme", "% of Complaints"]
    fb_widths = [Inches(0.4), Inches(5.0), Inches(1.5)]
    n_fb_rows = len(entertain_neg_top3) + 1
    fb_row_h = Inches(0.30)
    fb_hdr_h = Inches(0.35)
    fb_tbl_h = fb_hdr_h + fb_row_h * len(entertain_neg_top3)
    fb_shape = slide.shapes.add_table(n_fb_rows, 3, Inches(0.45), ent_fb_top + Inches(0.02),
                                       Inches(6.9), int(fb_tbl_h))
    fb_table = fb_shape.table
    for i, w in enumerate(fb_widths):
        fb_table.columns[i].width = int(w)
    fb_table.rows[0].height = int(fb_hdr_h)
    for ci, hdr in enumerate(fb_headers):
        cell = fb_table.cell(0, ci)
        set_cell_fill(cell, MEDIUM_BLUE)
        set_cell_text(cell, hdr, font_size=8, bold=True, color=WHITE)
    for ri, (theme, cnt, pct) in enumerate(entertain_neg_top3):
        row_idx = ri + 1
        fb_table.rows[row_idx].height = int(fb_row_h)
        bg = LIGHT_BG if ri % 2 == 0 else WHITE
        cell = fb_table.cell(row_idx, 0)
        set_cell_fill(cell, bg)
        set_cell_text(cell, str(ri + 1), font_size=9, color=DARK_GRAY)
        cell = fb_table.cell(row_idx, 1)
        set_cell_fill(cell, bg)
        set_cell_text(cell, theme, font_size=9, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        cell = fb_table.cell(row_idx, 2)
        set_cell_fill(cell, bg)
        set_cell_text(cell, f"{pct:.1f}%", font_size=9, bold=True, color=RED)

    add_branded_footer(slide, prs, 7, total_pages)

    # ─── SLIDE 8: Gate Entry + Parking + Seat Value ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Gate Entry + Parking + Seat Value")

    # Gate Entry table
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.25),
             "How did the time it took to pass through security and gate entry compare to expectations?",
             size=10, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.30), Inches(12.5), Inches(0.22),
             context_note, size=7, color=MUTED_GRAY)

    gate_h = build_table(slide, gate_rows, Inches(1.55), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7, row_height_override=Inches(0.32))

    # Parking table
    park_top = Inches(1.55) + gate_h + Inches(0.30)
    add_text(slide, Inches(0.45), park_top - Inches(0.22), Inches(12), Inches(0.22),
             "Parking Experience (Ingress & Egress Ratings + Time Expectation)",
             size=10, italic=True, color=GRAY)
    park_h = build_table(slide, parking_rows, park_top + Inches(0.05), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7, row_height_override=Inches(0.28))

    # Seat Value table
    seat_top = park_top + Inches(0.05) + park_h + Inches(0.30)
    add_text(slide, Inches(0.45), seat_top - Inches(0.22), Inches(12), Inches(0.22),
             "Thinking of the money you spent on your ticket, your view was a... (% Good Value)",
             size=10, italic=True, color=GRAY)
    build_table(slide, seat_rows, seat_top + Inches(0.05), Inches(0.45), Inches(12.4),
                col_headers_7, col_widths_7, row_height_override=Inches(0.32))

    # --- Parking Negative Feedback: Top 3 Categories (A) ---
    park_fb_top = seat_top + Inches(0.05) + Inches(0.32 * 2) + Inches(0.25)
    add_text(slide, Inches(0.45), park_fb_top - Inches(0.18), Inches(6), Inches(0.18),
             f"Top Parking Complaints (HS4, n={parking_neg_total})",
             size=8, bold=True, color=DARK_GRAY)

    pk_headers = ["#", "Category", "%"]
    pk_widths = [Inches(0.3), Inches(3.0), Inches(0.8)]
    n_pk_rows = len(parking_neg_top3) + 1
    pk_row_h = Inches(0.24)
    pk_hdr_h = Inches(0.28)
    pk_tbl_h = pk_hdr_h + pk_row_h * len(parking_neg_top3)
    pk_shape = slide.shapes.add_table(n_pk_rows, 3, Inches(0.45), park_fb_top,
                                       Inches(4.1), int(pk_tbl_h))
    pk_table = pk_shape.table
    for i, w in enumerate(pk_widths):
        pk_table.columns[i].width = int(w)
    pk_table.rows[0].height = int(pk_hdr_h)
    for ci, hdr in enumerate(pk_headers):
        cell = pk_table.cell(0, ci)
        set_cell_fill(cell, MEDIUM_BLUE)
        set_cell_text(cell, hdr, font_size=7, bold=True, color=WHITE)
    for ri, (cat, cnt, pct) in enumerate(parking_neg_top3):
        row_idx = ri + 1
        pk_table.rows[row_idx].height = int(pk_row_h)
        bg = LIGHT_BG if ri % 2 == 0 else WHITE
        cell = pk_table.cell(row_idx, 0)
        set_cell_fill(cell, bg)
        set_cell_text(cell, str(ri + 1), font_size=8, color=DARK_GRAY)
        cell = pk_table.cell(row_idx, 1)
        set_cell_fill(cell, bg)
        set_cell_text(cell, cat, font_size=8, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        cell = pk_table.cell(row_idx, 2)
        set_cell_fill(cell, bg)
        set_cell_text(cell, f"{pct:.1f}%", font_size=8, bold=True, color=RED)

    # --- Departure Traffic Subcategories (D) ---
    add_text(slide, Inches(5.0), park_fb_top - Inches(0.18), Inches(6), Inches(0.18),
             f"Parking Egress: Top Complaints (HS4, n={departure_total})",
             size=8, bold=True, color=DARK_GRAY)

    dt_headers = ["#", "Sub-Category", "%"]
    dt_widths = [Inches(0.3), Inches(4.5), Inches(0.8)]
    n_dt_rows = len(departure_subcats) + 1
    dt_shape = slide.shapes.add_table(n_dt_rows, 3, Inches(5.0), park_fb_top,
                                       Inches(5.6), int(pk_tbl_h))
    dt_table = dt_shape.table
    for i, w in enumerate(dt_widths):
        dt_table.columns[i].width = int(w)
    dt_table.rows[0].height = int(pk_hdr_h)
    for ci, hdr in enumerate(dt_headers):
        cell = dt_table.cell(0, ci)
        set_cell_fill(cell, MEDIUM_BLUE)
        set_cell_text(cell, hdr, font_size=7, bold=True, color=WHITE)
    for ri, (cat, cnt, pct) in enumerate(departure_subcats):
        row_idx = ri + 1
        dt_table.rows[row_idx].height = int(pk_row_h)
        bg = LIGHT_BG if ri % 2 == 0 else WHITE
        cell = dt_table.cell(row_idx, 0)
        set_cell_fill(cell, bg)
        set_cell_text(cell, str(ri + 1), font_size=8, color=DARK_GRAY)
        cell = dt_table.cell(row_idx, 1)
        set_cell_fill(cell, bg)
        set_cell_text(cell, cat, font_size=8, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
        cell = dt_table.cell(row_idx, 2)
        set_cell_fill(cell, bg)
        set_cell_text(cell, f"{pct:.1f}%", font_size=8, bold=True, color=RED)

    add_branded_footer(slide, prs, 8, total_pages)

    # ─── SLIDE 9: Ballpark App (Rays Wallet) — Concession Wait Expectations ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Ballpark App (Rays Wallet) \u2014 Concession Wait")

    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             "Concession Wait Expectations \u2014 Rays Wallet Users vs. Everyone Else (% of valid responses)",
             size=10, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             "Periods: April (Apr 6-26)  |  May 1-6 (HS3)  |  May 15-20 (HS4)  |  Wallet Users = TB_ADDON_26=1 (used Rays Wallet)",
             size=8, color=MUTED_GRAY)

    wallet_headers = [
        "Wait Category",
        "Wallet\nApril", "Wallet\nMay 1-6", "Wallet\nMay 15-20",
        "Non-Wallet\nApril", "Non-Wallet\nMay 1-6", "Non-Wallet\nMay 15-20"
    ]
    wallet_widths = [
        Inches(2.8), Inches(1.5), Inches(1.5), Inches(1.5),
        Inches(1.5), Inches(1.5), Inches(1.5)
    ]
    build_segment_table(slide, wallet_rows, Inches(1.65), Inches(0.45), Inches(12.3),
                        wallet_headers, wallet_widths)
    add_branded_footer(slide, prs, 9, total_pages)

    # ─── SLIDE 10: Ballpark App Usage Stats ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Ballpark App (Rays Wallet) \u2014 Usage & Experience")

    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             "% of respondents by Rays Wallet status + AVG Overall Experience (0-10) per category",
             size=10, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             "Periods: April (Apr 6-26)  |  May 1-6 (HS3)  |  May 15-20 (HS4)  |  % Diff = HS vs April (relative)",
             size=8, color=MUTED_GRAY)

    # Build app usage table
    # Columns: Category | April % | HS3 % | HS4 % | HS3 vs Apr % Diff | HS4 vs Apr % Diff
    #          | April Score | HS3 Score | HS4 Score | HS3 Score vs Apr % Diff | HS4 Score vs Apr % Diff
    app_col_headers = [
        "Category",
        "April\n% Usage", "HS3\n% Usage", "HS4\n% Usage",
        "HS3 vs Apr\n% Diff", "HS4 vs Apr\n% Diff",
        "April\nAVG Score", "HS3\nAVG Score", "HS4\nAVG Score",
        "HS3 vs Apr\nScore Diff", "HS4 vs Apr\nScore Diff",
    ]
    app_col_widths = [
        Inches(2.0), Inches(0.95), Inches(0.95), Inches(0.95),
        Inches(1.1), Inches(1.1),
        Inches(1.0), Inches(1.0), Inches(1.0),
        Inches(1.1), Inches(1.1),
    ]
    n_app_cols = 11
    n_app_rows = len(app_usage_rows) + 1
    app_row_height = Inches(0.40)
    app_header_height = Inches(0.50)
    app_table_height = app_header_height + app_row_height * len(app_usage_rows)

    table_shape = slide.shapes.add_table(n_app_rows, n_app_cols, Inches(0.35), Inches(1.65),
                                          Inches(12.6), int(app_table_height))
    table = table_shape.table

    for i, w in enumerate(app_col_widths):
        table.columns[i].width = int(w)

    # Header row
    table.rows[0].height = int(app_header_height)
    for ci, hdr in enumerate(app_col_headers):
        cell = table.cell(0, ci)
        set_cell_fill(cell, MEDIUM_BLUE)
        set_cell_text(cell, hdr, font_size=8, bold=True, color=WHITE)

    # Data rows
    for ri, (label, pcts, scores) in enumerate(app_usage_rows):
        row_idx = ri + 1
        table.rows[row_idx].height = int(app_row_height)
        bg_color = LIGHT_BG if ri % 2 == 0 else WHITE

        apr_pct, hs3_pct, hs4_pct = pcts
        apr_score, hs3_score, hs4_score = scores

        # % diffs (relative) for usage
        pct_diff_hs3 = pct_diff(apr_pct, hs3_pct)
        pct_diff_hs4 = pct_diff(apr_pct, hs4_pct)
        # % diffs (relative) for scores
        score_diff_hs3 = pct_diff(apr_score, hs3_score)
        score_diff_hs4 = pct_diff(apr_score, hs4_score)

        cell_values = [
            (label, PP_ALIGN.LEFT, False),
            (f"{apr_pct:.1f}%" if apr_pct is not None else "N/A", PP_ALIGN.CENTER, False),
            (f"{hs3_pct:.1f}%" if hs3_pct is not None else "N/A", PP_ALIGN.CENTER, False),
            (f"{hs4_pct:.1f}%" if hs4_pct is not None else "N/A", PP_ALIGN.CENTER, False),
            (pct_diff_hs3, PP_ALIGN.CENTER, True),
            (pct_diff_hs4, PP_ALIGN.CENTER, True),
            (f"{apr_score:.2f}" if apr_score is not None else "N/A", PP_ALIGN.CENTER, False),
            (f"{hs3_score:.2f}" if hs3_score is not None else "N/A", PP_ALIGN.CENTER, False),
            (f"{hs4_score:.2f}" if hs4_score is not None else "N/A", PP_ALIGN.CENTER, False),
            (score_diff_hs3, PP_ALIGN.CENTER, True),
            (score_diff_hs4, PP_ALIGN.CENTER, True),
        ]

        for ci, (text, align, is_diff) in enumerate(cell_values):
            cell = table.cell(row_idx, ci)
            set_cell_fill(cell, bg_color)
            if is_diff:
                if text is not None:
                    sign = "+" if text > 0 else ""
                    dc = GREEN if text > 0 else (RED if text < 0 else DARK_GRAY)
                    set_cell_text(cell, f"{sign}{text:.1f}%", font_size=9, bold=True, color=dc)
                else:
                    set_cell_text(cell, "N/A", font_size=9, color=MUTED_GRAY)
            else:
                fcolor = DARK_GRAY if text != "N/A" else MUTED_GRAY
                set_cell_text(cell, str(text), font_size=9, color=fcolor, alignment=align)

    add_branded_footer(slide, prs, 10, total_pages)

    # Save
    prs.save(OUTPUT_PATH)
    print(f"\nDone! Output: {OUTPUT_PATH}")
    print(f"  Slides: {total_pages}")
    print(f"  Layout: 10 slides")


if __name__ == "__main__":
    build_report()
