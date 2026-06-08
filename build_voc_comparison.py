"""
build_voc_comparison.py
Tampa Bay Rays - VOC Homestand Comparison PPTX Generator

Usage:
    python build_voc_comparison.py --hs 1 2                     # Compare HS1 vs HS2 (with 2024 baseline)
    python build_voc_comparison.py --hs 2 3                     # Compare HS2 vs HS3
    python build_voc_comparison.py --hs 2 3 --no-baseline       # Without 2024 column
    python build_voc_comparison.py --hs 2 3 --baseline-season 2025  # Different baseline

Requires: pip install snowflake-connector-python python-pptx

Config: Edit homestand_config.json to add new homestand dates.
"""

import argparse
import json
import os
import sys
import snowflake.connector
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
CONFIG_PATH = os.path.join(SCRIPT_DIR, "homestand_config.json")

# ── Snowflake connection defaults (uses SSO) ──
SF_ACCOUNT = "hta92307.east-us-2.azure"
SF_WAREHOUSE = "TBRDP_DW_CORTEX_XS_WH"
SF_DATABASE = "TBRDP_DW_DEV"
SF_SCHEMA = "IM_RPT"
SF_ROLE = "TBRDP_DW_PROD_CORTEX_USER"
SF_AUTHENTICATOR = "externalbrowser"

# ── View ──
VOC_VIEW = "TBRDP_DW_DEV.IM_RPT.V_SBL_QUALTRICS_VOC_POST_ATTENDANCE_FULL_CORTEX_AI"
PARKING_VIEW = "TBRDP_DW_DEV.IM_RPT.V_SBL_QUALTRICS_VOC_POST_ATTENDANCE_FULL"

# ── Theme colors (Rays brand palette) ──
NAVY = RGBColor(0x09, 0x2C, 0x5C)
SKY = RGBColor(0x8F, 0xBC, 0xE6)
YELLOW = RGBColor(0xF5, 0xD1, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x4B, 0x4B, 0x4B)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xFF, 0x6B, 0x6B)
MEDIUM_BLUE = RGBColor(0x0D, 0x3B, 0x7A)
DARK_BLUE_BG = RGBColor(0x07, 0x20, 0x44)
MUTED_GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT_BG = RGBColor(0xF2, 0xF4, 0xF8)

# ── Logo ──
LOGO_PATH = os.path.join(SCRIPT_DIR, "MLB Logos", "TB_White.png")


# ═══════════════════════════════════════════════════════════════════════
# CONFIG
# ═══════════════════════════════════════════════════════════════════════

def load_config():
    with open(CONFIG_PATH, "r") as f:
        return json.load(f)


def get_homestand(config, season, hs_num):
    for hs in config["homestands"]:
        if hs["season"] == season and hs["homestand_num"] == hs_num:
            return hs
    raise ValueError(f"Homestand {hs_num} for season {season} not found in {CONFIG_PATH}")


# ═══════════════════════════════════════════════════════════════════════
# SNOWFLAKE CONNECTION
# ═══════════════════════════════════════════════════════════════════════

def get_connection():
    return snowflake.connector.connect(
        user="YTAKETANI@RAYSBASEBALL.COM",
        account=SF_ACCOUNT,
        authenticator=SF_AUTHENTICATOR,
        warehouse=SF_WAREHOUSE,
        database=SF_DATABASE,
        schema=SF_SCHEMA,
        role=SF_ROLE,
    )


def run_query(conn, sql):
    """Execute SQL and return list of tuples."""
    cur = conn.cursor()
    try:
        cur.execute(sql)
        return cur.fetchall()
    finally:
        cur.close()


# ═══════════════════════════════════════════════════════════════════════
# QUERY BUILDERS — one per slide category
# ═══════════════════════════════════════════════════════════════════════

def query_core_ratings(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 1: Core Satisfaction Ratings (0-10 scale AVG)."""
    metrics = [
        ("Overall Experience", "OVERALL_NUMRAT"),
        ("Concessions", "CONCESS_NUMRAT"),
        ("Parking", "PARKING_NUMRAT"),
        ("Entertainment", "ENTERTAIN_NUMRAT"),
    ]
    rows = []
    for label, col in metrics:
        sql = f"""
        SELECT
            {_baseline_avg(col, baseline_season)},
            AVG(CASE WHEN SEASON = {current_season} AND {col} < 80 THEN {col} END) AS s_current,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND {col} < 80 THEN {col} END) AS hs_prev,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND {col} < 80 THEN {col} END) AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round2(r[0])
        s_cur = _round2(r[1])
        hp = _round2(r[2])
        hc = _round2(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, False, False))
    return ("Core Satisfaction Ratings", "AVG Score (0-10 Scale) | Tropicana Field", rows)


def query_staff_ratings(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 2: Staff Satisfaction (normalized 0-10)."""
    staff_pairs = [
        ("Parking Staff",                 "_1",  "_1"),
        ("Fan Host Usher",                "_3",  "_3"),
        ("Concessions Staff",             "_5",  "_5"),
        ("Retail Staff",                  "_6",  "_6"),
        ("Security",                      "_11", "_60"),
        ("Tech Team",                     "_38", "_38"),
        ("Fan Host Ticket Taker/Scanner", "_10", "_61"),
    ]
    rows = []
    for label, s4, s23 in staff_pairs:
        col_old = f"TB_ADDON_4{s4}"
        col_new = f"TB_ADDON_23{s23}"
        sql = f"""
        SELECT
            AVG(CASE WHEN SEASON = {baseline_season} AND {col_old} BETWEEN 1 AND 5
                      THEN ({col_old} - 1) * 2.5 END) AS s_base,
            AVG(CASE WHEN SEASON = {current_season} AND {col_new} < 80
                      THEN {col_new} END) AS s_current,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                          AND {col_new} < 80
                      THEN {col_new} END) AS hs_prev,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                          AND {col_new} < 80
                      THEN {col_new} END) AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round2(r[0])
        s_cur = _round2(r[1])
        hp = _round2(r[2])
        hc = _round2(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, False, False))
    return ("Staff Satisfaction Ratings", "AVG Score (0-10 Scale, Normalized) | Tropicana Field", rows)


def query_concessions_grid(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 3: Concessions Grid (% Highly Satisfied)."""
    metrics = [
        ("Value",            "CONCESS_GRID_VALUE_DESC"),
        ("Customer Service", "CONCESS_GRID_CUSTSERV_DESC"),
        ("Selection",        "CONCESS_GRID_SELECTION_DESC"),
        ("Cleanliness",      "CONCESS_GRID_CLEAN_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    rows = []
    for label, col in metrics:
        sql = f"""
        SELECT
            {_baseline_pct_highly(col, baseline_season, valid_labels)},
            {_current_pct_highly(col, current_season, valid_labels)} AS s_current,
            {_hs_pct_highly(col, hs_prev, valid_labels)} AS hs_prev,
            {_hs_pct_highly(col, hs_curr, valid_labels)} AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round1(r[0])
        s_cur = _round1(r[1])
        hp = _round1(r[2])
        hc = _round1(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, True, False))

    return ("Concessions Grid", "% Highly Satisfied | Tropicana Field", rows)


def query_concessions_wait(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 4: Concessions Wait-Time Expectation Extremes (% of valid responses)."""
    rows = []
    valid_wait = "('Much more than what I expected','Slightly more than what I expected','About what I expected','Slightly less what I expected','Much less than what I expected')"
    for label_name, text_val, invert in [
        ("% Much Less Wait Than Expected", "Much less than what I expected", False),
        ("% Much More Wait Than Expected", "Much more than what I expected", True),
    ]:
        sql = f"""
        SELECT
            ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1)
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        r = run_query(conn, sql)[0]
        s_base = _float(r[0])
        s_cur = _float(r[1])
        hp = _float(r[2])
        hc = _float(r[3])
        if s_cur is not None or hp is not None or hc is not None:
            rows.append((label_name, s_base, s_cur, hp, hc, True, invert))

    return ("Concessions Wait Expectation", "% of Valid Responses (Extremes) | Tropicana Field", rows)


def query_entertainment_grid(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 4: Entertainment Grid (% Highly Satisfied)."""
    metrics = [
        ("Scoreboard",      "ENTERTAIN_GRID_SCOREBOARD_DESC"),
        ("Music",           "ENTERTAIN_GRID_MUSIC_DESC"),
        ("Games/Contests",  "ENTERTAIN_GRID_GAMES_DESC"),
        ("Kids Activities", "ENTERTAIN_GRID_KIDS_ACTIVITIES_DESC"),
        ("Pregame Content", "ENTERTAIN_GRID_PREGAME_CONTENT_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    rows = []
    for label, col in metrics:
        sql = f"""
        SELECT
            {_baseline_pct_highly(col, baseline_season, valid_labels)},
            {_current_pct_highly(col, current_season, valid_labels)} AS s_current,
            {_hs_pct_highly(col, hs_prev, valid_labels)} AS hs_prev,
            {_hs_pct_highly(col, hs_curr, valid_labels)} AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round1(r[0])
        s_cur = _round1(r[1])
        hp = _round1(r[2])
        hc = _round1(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, True, False))
    return ("Entertainment Grid", "% Highly Satisfied | Tropicana Field", rows)


def query_food_quality(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Slide 5: Food Quality (% Highly Satisfied)."""
    food_items = [
        ("Hotdog",                "CONCESS_QUALITY_HOTDOG_DESC"),
        ("Burgers",               "CONCESS_QUALITY_BURGERS_DESC"),
        ("Chicken",               "CONCESS_QUALITY_CHICKEN_DESC"),
        ("Pizza",                 "CONCESS_QUALITY_PIZZA_DESC"),
        ("Nachos",                "CONCESS_QUALITY_NACHOS_DESC"),
        ("Fries",                 "CONCESS_QUALITY_FRIES_DESC"),
        ("Alcohol",               "CONCESS_QUALITY_ALCOHOL_DESC"),
        ("Non-Alcohol Beverages", "CONCESS_QUALITY_NONALCOHOL_DESC"),
        ("Ice Cream",             "CONCESS_QUALITY_ICECREAM_DESC"),
        ("Nuts",                  "CONCESS_QUALITY_NUTS_DESC"),
        ("Popcorn",               "CONCESS_QUALITY_POPCORN_DESC"),
        ("Pretzels",              "CONCESS_QUALITY_PRETZELS_DESC"),
        ("Other Dessert",         "CONCESS_QUALITY_OTHER_DESSERT_DESC"),
        ("Other Entree",          "CONCESS_QUALITY_OTHER_ENTREE_DESC"),
        ("Sandwich",              "CONCESS_QUALITY_SANDWICH_DESC"),
        ("Sausage",               "CONCESS_QUALITY_SAUSAGE_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    rows = []
    for label, col in food_items:
        sql = f"""
        SELECT
            {_baseline_pct_highly(col, baseline_season, valid_labels)},
            {_current_pct_highly(col, current_season, valid_labels)} AS s_current,
            {_hs_pct_highly(col, hs_prev, valid_labels)} AS hs_prev,
            {_hs_pct_highly(col, hs_curr, valid_labels)} AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round1(r[0])
        s_cur = _round1(r[1])
        hp = _round1(r[2])
        hc = _round1(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, s_base, s_cur, hp, hc, True, False))
    return ("Food Quality", "% Highly Satisfied | Tropicana Field", rows)


def query_gate_entry(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Gate Entry Time Expectation (extremes only).
    Uses numeric GE_TIME_EXPECT column (1=More, 2=About, 3=Less) with < 80 junk filter."""
    rows = []
    for label_name, code_val, invert in [
        ("% Less Than Expected", 3, False),
        ("% More Than Expected", 1, True),
    ]:
        sql = f"""
        SELECT
            ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND GE_TIME_EXPECT < 80 AND GE_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND GE_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND GE_TIME_EXPECT < 80 AND GE_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND GE_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND GE_TIME_EXPECT < 80 AND GE_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND GE_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND GE_TIME_EXPECT < 80 AND GE_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND GE_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1)
        FROM {PARKING_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        r = run_query(conn, sql)[0]
        s_base = _float(r[0])
        s_cur = _float(r[1])
        hp = _float(r[2])
        hc = _float(r[3])
        if s_cur is not None or hp is not None or hc is not None:
            rows.append((label_name, s_base, s_cur, hp, hc, True, invert))

    return ("Gate Entry Experience", "% of Valid Responses (Time Expectation) | Tropicana Field", rows)


def query_parking(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Parking Experience — numeric ratings + parking time expectation.
    Uses PARKING_VIEW which has PARKING_EXIT_NUMRAT and numeric PARKING_TIME_EXPECT."""
    rows = []

    # Parking numeric ratings (0-10 AVG)
    for label, col in [
        ("Finding & Parking (AVG 0-10)", "PARKING_NUMRAT"),
        ("Exiting Parking Lot (AVG 0-10)", "PARKING_EXIT_NUMRAT"),
    ]:
        sql = f"""
        SELECT
            {_baseline_avg(col, baseline_season)},
            AVG(CASE WHEN SEASON = {current_season} AND {col} < 80 THEN {col} END),
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND {col} < 80 THEN {col} END),
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND {col} < 80 THEN {col} END)
        FROM {PARKING_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        r = run_query(conn, sql)[0]
        s_base = _round2(r[0])
        s_cur = _round2(r[1])
        hp = _round2(r[2])
        hc = _round2(r[3])
        if s_cur is not None or hp is not None or hc is not None:
            rows.append((label, s_base, s_cur, hp, hc, False, False))

    # Parking Time expectation categories (1=More, 2=About, 3=Less)
    for label_name, code_val, invert in [
        ("% Less Than Expected", 3, False),
        ("% About What Expected", 2, False),
        ("% More Than Expected", 1, True),
    ]:
        sql = f"""
        SELECT
            ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND PARKING_TIME_EXPECT < 80 AND PARKING_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND PARKING_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND PARKING_TIME_EXPECT < 80 AND PARKING_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND PARKING_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND PARKING_TIME_EXPECT < 80 AND PARKING_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND PARKING_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND PARKING_TIME_EXPECT < 80 AND PARKING_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND PARKING_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1)
        FROM {PARKING_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        r = run_query(conn, sql)[0]
        s_base = _float(r[0])
        s_cur = _float(r[1])
        hp = _float(r[2])
        hc = _float(r[3])
        if s_cur is not None or hp is not None or hc is not None:
            rows.append((label_name, s_base, s_cur, hp, hc, True, invert))

    return ("Parking Experience", "AVG Ratings + Time Expectation | Tropicana Field", rows)


def query_seat_value(conn, baseline_season, current_season, hs_prev, hs_curr):
    """Seat View Value: % Good Value (code 1).
    Codes: 1=Good value, 2=Fair value, 3=Poor value; >= 80 junk."""
    rows = []
    sql = f"""
    SELECT
        ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND SEATVIEW_VALUE < 80 AND SEATVIEW_VALUE = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND SEATVIEW_VALUE < 80 THEN 1 ELSE 0 END), 0), 1),
        ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND SEATVIEW_VALUE < 80 AND SEATVIEW_VALUE = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND SEATVIEW_VALUE < 80 THEN 1 ELSE 0 END), 0), 1),
        ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND SEATVIEW_VALUE < 80 AND SEATVIEW_VALUE = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND SEATVIEW_VALUE < 80 THEN 1 ELSE 0 END), 0), 1),
        ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND SEATVIEW_VALUE < 80 AND SEATVIEW_VALUE = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND SEATVIEW_VALUE < 80 THEN 1 ELSE 0 END), 0), 1)
    FROM {PARKING_VIEW}
    WHERE STADIUM = 'Tropicana Field'
      AND SEASON IN ({baseline_season}, {current_season})
    """
    r = run_query(conn, sql)[0]
    s_base = _float(r[0])
    s_cur = _float(r[1])
    hp = _float(r[2])
    hc = _float(r[3])
    if s_cur is not None or hp is not None or hc is not None:
        rows.append(("% Good Seat Value", s_base, s_cur, hp, hc, True, False))

    return ("Seat View Value", "% Good Value for Money Spent on Ticket | Tropicana Field", rows)


# ═══════════════════════════════════════════════════════════════════════
# SQL HELPERS
# ═══════════════════════════════════════════════════════════════════════

def _baseline_avg(col, baseline_season):
    return f"AVG(CASE WHEN SEASON = {baseline_season} AND {col} < 80 THEN {col} END)"


def _baseline_pct_highly(col, baseline_season, valid_labels):
    return (
        f"ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END) "
        f"/ NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)"
    )


def _current_pct_highly(col, current_season, valid_labels):
    return (
        f"ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END) "
        f"/ NULLIF(SUM(CASE WHEN SEASON = {current_season} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)"
    )


def _hs_pct_highly(col, hs, valid_labels):
    return (
        f"ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs['start_date']}' AND '{hs['end_date']}' "
        f"AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END) "
        f"/ NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs['start_date']}' AND '{hs['end_date']}' "
        f"AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)"
    )


def _round2(val):
    return round(float(val), 2) if val is not None else None


def _round1(val):
    return round(float(val), 1) if val is not None else None


def _float(val):
    return float(val) if val is not None else None


# ═══════════════════════════════════════════════════════════════════════
# FORMATTING HELPERS
# ═══════════════════════════════════════════════════════════════════════

def pct_diff(a, b):
    """% change from a to b."""
    if a is None or b is None or a == 0:
        return None
    return round((b - a) / a * 100, 1)


def fmt_val(val, is_pct):
    if val is None:
        return "N/A"
    return f"{val:.1f}%" if is_pct else f"{val:.2f}"


def fmt_diff(diff):
    if diff is None:
        return "N/A"
    sign = "+" if diff > 0 else ""
    return f"{sign}{diff:.1f}%"


# ═══════════════════════════════════════════════════════════════════════
# PPTX BUILDER
# ═══════════════════════════════════════════════════════════════════════

def set_cell_text(cell, text, font_size=10, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_cell_fill(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def add_rect(slide, prs, x, y, w, h, fill, line=None):
    """Add a filled rectangle shape."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=10, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    """Add a text box."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    run.font.italic = italic
    return tb


def add_branded_header(slide, prs, title_text, subtitle_text):
    """Add Rays-branded header: navy bar + yellow accent stripe + logo + title."""
    # Navy header bar
    add_rect(slide, prs, 0, 0, prs.slide_width, Inches(0.95), NAVY)
    # Yellow accent stripe
    add_rect(slide, prs, 0, Inches(0.95), prs.slide_width, Inches(0.06), YELLOW)
    # Eyebrow text
    add_text(slide, Inches(0.45), Inches(0.10), Inches(11), Inches(0.30),
             "VOC POST-ATTENDANCE SURVEY  ·  HOMESTAND COMPARISON",
             size=10, bold=True, color=SKY)
    # Title
    add_text(slide, Inches(0.45), Inches(0.36), Inches(11), Inches(0.55),
             title_text, size=22, bold=True, color=WHITE)
    # Logo
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(12.35), Inches(0.18), height=Inches(0.65))


def add_branded_footer(slide, prs, page_num, total_pages):
    """Add Rays-branded footer: navy bar + attribution text."""
    add_rect(slide, prs, 0, Inches(7.20), prs.slide_width, Inches(0.30), NAVY)
    add_text(slide, Inches(0.45), Inches(7.23), Inches(8), Inches(0.25),
             "Tampa Bay Rays  ·  Strategy & Analytics",
             size=9, color=SKY)
    add_text(slide, Inches(10.0), Inches(7.23), Inches(3.0), Inches(0.25),
             f"Source: Qualtrics VOC  |  {page_num} / {total_pages}",
             size=9, color=SKY, align=PP_ALIGN.RIGHT)


def build_pptx(all_slides, output_path, include_baseline, baseline_season,
               current_season, hs_prev_num, hs_curr_num, hs_prev, hs_curr,
               baseline_responses, baseline_games, current_responses, current_games):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bl_tag = str(baseline_season)[-2:]
    cur_tag = str(current_season)[-2:]

    if include_baseline:
        col_headers = [
            "Metric",
            f"'{bl_tag} Score",
            f"'{cur_tag} AVG",
            f"HS{hs_prev_num} Score",
            f"HS{hs_curr_num} Score",
            f"HS{hs_curr_num} vs HS{hs_prev_num}\n% Diff",
            f"'{cur_tag} vs '{bl_tag}\n% Diff",
        ]
        n_cols = 7
        col_widths = [Inches(2.8), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3), Inches(2.0), Inches(2.0)]
    else:
        col_headers = [
            "Metric",
            f"'{cur_tag} AVG",
            f"HS{hs_prev_num} Score",
            f"HS{hs_curr_num} Score",
            f"HS{hs_curr_num} vs HS{hs_prev_num}\n% Diff",
        ]
        n_cols = 5
        col_widths = [Inches(3.2), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.8)]

    # Context note
    parts = []
    if include_baseline:
        parts.append(f"'{bl_tag} = Full Season ({baseline_games} games, {baseline_responses:,} responses)")
    total_hs_games = hs_prev["games"] + hs_curr["games"]
    parts.append(f"'{cur_tag} = {hs_curr_num} Homestands ({total_hs_games} games, {current_responses:,} responses)")
    parts.append(f"HS{hs_prev_num} = {hs_prev['start_date'][5:].replace('-','/')} ({hs_prev['games']}g)")
    parts.append(f"HS{hs_curr_num} = {hs_curr['start_date'][5:].replace('-','/')} ({hs_curr['games']}g)")
    context_note = "  |  ".join(parts)

    total_pages = len(all_slides)

    for slide_idx, slide_info in enumerate(all_slides):
        title_text, subtitle_text, data_rows = slide_info

        # Sort rows by '26 vs '24 % Diff (most positive first) when baseline included
        if include_baseline:
            def sort_key(row):
                _, s_base, s_cur, _, _, is_pct, invert_diff = row
                d = pct_diff(s_base, s_cur)
                if d is None:
                    return float('-inf')
                if invert_diff:
                    d = -d
                return -d
            data_rows = sorted(data_rows, key=sort_key)

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # Branded header (navy bar + yellow stripe + logo)
        add_branded_header(slide, prs, title_text, subtitle_text)

        # Subtitle line (below yellow stripe)
        add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
                 subtitle_text, size=12, italic=True, color=GRAY)

        # Context note
        add_text(slide, Inches(0.45), Inches(1.38), Inches(12), Inches(0.25),
                 context_note, size=9, color=MUTED_GRAY)

        # Table
        n_rows = len(data_rows) + 1
        row_height = Inches(0.38) if len(data_rows) <= 8 else Inches(0.32)
        header_height = Inches(0.50)
        table_height = header_height + row_height * len(data_rows)
        top = Inches(1.72)
        left = Inches(0.45)

        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, Inches(12.4), int(table_height))
        table = table_shape.table

        for i, w in enumerate(col_widths):
            table.columns[i].width = int(w)

        # Header row
        table.rows[0].height = int(header_height)
        for ci, hdr in enumerate(col_headers):
            cell = table.cell(0, ci)
            set_cell_fill(cell, MEDIUM_BLUE)
            set_cell_text(cell, hdr, font_size=11, bold=True, color=WHITE)

        # Data rows
        for ri, (metric, s_base, s_cur, hp, hc, is_pct, invert_diff) in enumerate(data_rows):
            row_idx = ri + 1
            table.rows[row_idx].height = int(row_height)
            hs_diff = pct_diff(hp, hc)
            baseline_diff = pct_diff(s_base, s_cur) if include_baseline else None

            # For inverted metrics, flip diff sign for display and color coding
            if invert_diff:
                if hs_diff is not None:
                    hs_diff = -hs_diff
                if baseline_diff is not None:
                    baseline_diff = -baseline_diff

            bg_color = LIGHT_BG if ri % 2 == 0 else WHITE

            if include_baseline:
                cell_data = [
                    (metric, 10, False, DARK_GRAY, PP_ALIGN.LEFT, None),
                    (fmt_val(s_base, is_pct), 10, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                    (fmt_val(s_cur, is_pct), 10, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                    (fmt_val(hp, is_pct), 10, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                    (fmt_val(hc, is_pct), 10, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                    (None, 10, True, None, PP_ALIGN.CENTER, hs_diff),
                    (None, 10, True, None, PP_ALIGN.CENTER, baseline_diff),
                ]
            else:
                cell_data = [
                    (metric, 10, False, DARK_GRAY, PP_ALIGN.LEFT, None),
                    (fmt_val(s_cur, is_pct), 10, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                    (fmt_val(hp, is_pct), 10, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                    (fmt_val(hc, is_pct), 10, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                    (None, 10, True, None, PP_ALIGN.CENTER, hs_diff),
                ]

            for ci, (text, fsize, is_diff_col, fcolor, align, diff_val) in enumerate(cell_data):
                cell = table.cell(row_idx, ci)
                set_cell_fill(cell, bg_color)
                if is_diff_col:
                    if diff_val is not None:
                        dc = GREEN if diff_val > 0 else (RED if diff_val < 0 else DARK_GRAY)
                        set_cell_text(cell, fmt_diff(diff_val), font_size=fsize, bold=True, color=dc)
                    else:
                        set_cell_text(cell, "N/A", font_size=fsize, color=MUTED_GRAY)
                else:
                    set_cell_text(cell, text, font_size=fsize, bold=False, color=fcolor, alignment=align)

        # Branded footer (navy bar + Strategy & Analytics)
        add_branded_footer(slide, prs, slide_idx + 1, total_pages)

    prs.save(output_path)
    print(f"Saved: {output_path}")


# ═══════════════════════════════════════════════════════════════════════
# RESPONSE COUNT HELPER
# ═══════════════════════════════════════════════════════════════════════

def get_response_counts(conn, baseline_season, current_season):
    sql = f"""
    SELECT
        COUNT(CASE WHEN SEASON = {baseline_season} THEN 1 END) AS base_responses,
        COUNT(DISTINCT CASE WHEN SEASON = {baseline_season} THEN GAME_DATE END) AS base_games,
        COUNT(CASE WHEN SEASON = {current_season} THEN 1 END) AS cur_responses,
        COUNT(DISTINCT CASE WHEN SEASON = {current_season} THEN GAME_DATE END) AS cur_games
    FROM {VOC_VIEW}
    WHERE STADIUM = 'Tropicana Field'
      AND SEASON IN ({baseline_season}, {current_season})
    """
    r = run_query(conn, sql)[0]
    return int(r[0]), int(r[1]), int(r[2]), int(r[3])


# ═══════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════

def main():
    parser = argparse.ArgumentParser(description="VOC Homestand Comparison PPTX Generator")
    parser.add_argument("--hs", nargs=2, type=int, required=True,
                        help="Two consecutive homestand numbers to compare, e.g. --hs 1 2")
    parser.add_argument("--season", type=int, default=2026,
                        help="Current season (default: 2026)")
    parser.add_argument("--baseline-season", type=int, default=2024,
                        help="Baseline season for year-over-year comparison (default: 2024)")
    parser.add_argument("--no-baseline", action="store_true",
                        help="Omit the baseline season column and year-over-year %% Diff column")
    parser.add_argument("--output", type=str, default=None,
                        help="Output PPTX path (default: VOC_HS<N-1>_vs_HS<N>.pptx)")
    args = parser.parse_args()

    hs_prev_num, hs_curr_num = args.hs
    current_season = args.season
    baseline_season = args.baseline_season
    include_baseline = not args.no_baseline

    config = load_config()
    hs_prev = get_homestand(config, current_season, hs_prev_num)
    hs_curr = get_homestand(config, current_season, hs_curr_num)

    if args.output:
        output_path = args.output
    else:
        output_path = os.path.join(SCRIPT_DIR, f"VOC_HS{hs_prev_num}_vs_HS{hs_curr_num}.pptx")

    print(f"Connecting to Snowflake (SSO)...")
    conn = get_connection()
    print(f"Connected. Querying VOC data...")

    base_resp, base_games, cur_resp, cur_games = get_response_counts(conn, baseline_season, current_season)

    print("  [1/9] Core Satisfaction Ratings...")
    slide1 = query_core_ratings(conn, baseline_season, current_season, hs_prev, hs_curr)
    print("  [2/9] Staff Satisfaction Ratings...")
    slide2 = query_staff_ratings(conn, baseline_season, current_season, hs_prev, hs_curr)
    print("  [3/9] Concessions Grid...")
    slide3 = query_concessions_grid(conn, baseline_season, current_season, hs_prev, hs_curr)
    print("  [4/9] Concessions Wait Expectation...")
    slide4 = query_concessions_wait(conn, baseline_season, current_season, hs_prev, hs_curr)
    print("  [5/9] Entertainment Grid...")
    slide5 = query_entertainment_grid(conn, baseline_season, current_season, hs_prev, hs_curr)
    print("  [6/9] Food Quality...")
    slide6 = query_food_quality(conn, baseline_season, current_season, hs_prev, hs_curr)
    print("  [7/9] Gate Entry Experience...")
    slide7 = query_gate_entry(conn, baseline_season, current_season, hs_prev, hs_curr)
    print("  [8/9] Parking Experience...")
    slide8 = query_parking(conn, baseline_season, current_season, hs_prev, hs_curr)
    print("  [9/9] Seat View Value...")
    slide9 = query_seat_value(conn, baseline_season, current_season, hs_prev, hs_curr)

    conn.close()
    print("Queries complete. Building PPTX...")

    all_slides = [slide1, slide2, slide3, slide4, slide5, slide6, slide7, slide8, slide9]

    build_pptx(
        all_slides, output_path, include_baseline, baseline_season,
        current_season, hs_prev_num, hs_curr_num, hs_prev, hs_curr,
        base_resp, base_games, cur_resp, cur_games,
    )

    print(f"\nDone! Output: {output_path}")
    print(f"  Slides: {len(all_slides)}")
    print(f"  Layout: {'7 columns (with baseline)' if include_baseline else '5 columns (no baseline)'}")
    print(f"  Sorted by: {'year-over-year % Diff (most positive first)' if include_baseline else 'HS vs HS % Diff (most positive first)'}")


if __name__ == "__main__":
    main()
