"""
5-slide consultative brief — v3 (4 buckets).

Slide 1: Thesis  ·  why each bucket matters (business + creative), 4 buckets at a glance
Slide 2: Concepts 1-2  (Families, Multi-Gen)  2 copy angles + 2 visuals each
Slide 3: Concepts 3-4  (Couples, Social)  2 copy angles + 2 visuals each
Slide 4: Build spec for Segments 1-2
Slide 5: Build spec for Segments 3-4

Audience:  5,565 (2026 season through 15 games, with BUYER_TYPE exclusion)
Source:    `mlb-dataeng-prod.wheelhouse_rays.qualtrics_voc_post_attendance_full`
Output:    HS12_HighSat_Marketing_Brief_v3.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

OUT  = r"C:/Users/ytaketani/voc_insights_agent/sales_marketing_opp/VOC_mktg_campaign/HS12_HighSat_Marketing_Brief_v3.pptx"
LOGO = r"C:/Users/ytaketani/voc_insights_agent/MLB Logos/TB_White.png"

NAVY   = RGBColor(0x09, 0x2C, 0x5C)
SKY    = RGBColor(0x8F, 0xBC, 0xE6)
YELLOW = RGBColor(0xF5, 0xD1, 0x30)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x4B, 0x4B, 0x4B)
LIGHT  = RGBColor(0xF2, 0xF4, 0xF8)
CODEBG = RGBColor(0xEC, 0xEF, 0xF4)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
TOTAL_PAGES = 5


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=12, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text if isinstance(text, list) else [text]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.name = font; r.font.size = Pt(size); r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_header(slide, title,
               eyebrow="HOMESTAND 12  ·  VOC HIGH-SATISFACTION CAMPAIGN BRIEF  ·  v3"):
    add_rect(slide, 0, 0, prs.slide_width, Inches(0.95), NAVY)
    add_rect(slide, 0, Inches(0.95), prs.slide_width, Inches(0.06), YELLOW)
    add_text(slide, Inches(0.45), Inches(0.10), Inches(11), Inches(0.30),
             eyebrow, size=10, bold=True, color=SKY)
    add_text(slide, Inches(0.45), Inches(0.36), Inches(11), Inches(0.55),
             title, size=22, bold=True, color=WHITE)
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(12.35), Inches(0.18), height=Inches(0.65))


def add_footer(slide, page):
    add_rect(slide, 0, Inches(7.20), prs.slide_width, Inches(0.30), NAVY)
    add_text(slide, Inches(0.45), Inches(7.23), Inches(8), Inches(0.25),
             "Tampa Bay Rays  ·  Strategy & Analytics  ·  Prepared for Marketing",
             size=9, color=SKY)
    add_text(slide, Inches(11.0), Inches(7.23), Inches(2.0), Inches(0.25),
             f"{page} / {TOTAL_PAGES}", size=9, color=SKY, align=PP_ALIGN.RIGHT)


# =========================================================================
# SLIDE 1 — Why each bucket matters (business opportunity + creative case)
# =========================================================================
s1 = prs.slides.add_slide(BLANK)
add_header(s1, "Four audience buckets — why each one matters for targeted touchpoints")

# Audience definition strip
add_rect(s1, Inches(0.45), Inches(1.20), Inches(12.45), Inches(0.95), LIGHT)
add_text(s1, Inches(0.65), Inches(1.30), Inches(3.5), Inches(0.30),
         "THE AUDIENCE  +  THE TEST", size=10, bold=True, color=NAVY)
add_text(s1, Inches(0.65), Inches(1.55), Inches(12.0), Inches(0.55),
         ["5,565 qualified fans (2026 season, 15 games through May 3)  ·  Filter:  OVERALL_NUMRAT in (8,9,10)  AND  TB_ADDON_6 != 1  AND  plan-holder BUYER_TYPEs excluded",
          "Test design:  External control = non-VOC fans (the majority).  Treatment = these 4 buckets, integrated as Touch 2 in an existing journey."],
         size=11, bold=False, color=GRAY)

# WHY EACH BUCKET MATTERS panel (left)
add_rect(s1, Inches(0.45), Inches(2.30), Inches(5.90), Inches(4.75), WHITE, line=NAVY)
add_rect(s1, Inches(0.45), Inches(2.30), Inches(5.90), Inches(0.45), NAVY)
add_text(s1, Inches(0.60), Inches(2.36), Inches(5.7), Inches(0.35),
         "WHY EACH BUCKET MATTERS", size=11, bold=True, color=WHITE)

bucket_rationale = [
    ("1. Families (kids)",
     "84% rated 9-10. Over-indexes on ice cream (37%), chicken (28%), popcorn (30%) — kid-friendly F&B signals a memorable outing. Only 20% are prior buyers: high engagement + low conversion = major acquisition headroom. Creative must feel like a family highlight reel."),
    ("2. Multi-Gen Reunion",
     "84% rated 9-10, median age 60. The only bucket where the emotional hook is generational — 'bring your parents back.' No other brand touchpoint speaks to this bond. 25% prior buyers with high avidity means they convert when the message resonates."),
    ("3. Couples",
     "Largest bucket (2,016). 67% are 55+ with disposable income, 50% buy alcohol, only 27% prior buyers. Massive acquisition gap in a high-spend demo. Creative must feel like a date-night invitation, not a sports ad."),
    ("4. Social / Crew",
     "Lowest prior-buyer rate (17%), lowest avidity, lowest team-as-favorite (15%). This is pure incremental upside — if we move this group even slightly, it is the biggest volume lever in the program. Creative must trigger social proof and FOMO."),
]
y = Inches(2.85)
for hd, body in bucket_rationale:
    add_text(s1, Inches(0.60), y, Inches(5.55), Inches(0.25),
             hd, size=10, bold=True, color=NAVY)
    add_text(s1, Inches(0.60), y + Inches(0.25), Inches(5.55), Inches(0.80),
             body, size=8.5, color=GRAY)
    y += Inches(1.10)

# 4 buckets at-a-glance (right)
add_rect(s1, Inches(6.60), Inches(2.30), Inches(6.30), Inches(4.75), WHITE, line=NAVY)
add_rect(s1, Inches(6.60), Inches(2.30), Inches(6.30), Inches(0.45), NAVY)
add_text(s1, Inches(6.75), Inches(2.36), Inches(6.0), Inches(0.35),
         "FOUR CUTS  ·  MUTUALLY EXCLUSIVE  ·  5,565 FANS",
         size=11, bold=True, color=WHITE)

segs_glance = [
    ("1", "Families (kids)",       "858",   "Median 46  ·  parents w/ kids under 18"),
    ("2", "Multi-Gen Reunion",     "1,310", "Median 60  ·  adult kids or other family"),
    ("3", "Couples",               "2,016", "Median 64  ·  spouse / partner present"),
    ("4", "Social / Crew",         "1,845", "Median 62  ·  friends, business, solo, other"),
]
y = Inches(2.95)
for num, name, n, sub in segs_glance:
    add_rect(s1, Inches(6.75), y, Inches(0.65), Inches(0.85), YELLOW)
    add_text(s1, Inches(6.75), y, Inches(0.65), Inches(0.85),
             num, size=24, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s1, Inches(7.60), y + Inches(0.05), Inches(3.6), Inches(0.35),
             name, size=13, bold=True, color=NAVY)
    add_text(s1, Inches(7.60), y + Inches(0.40), Inches(3.6), Inches(0.40),
             sub, size=10, color=GRAY)
    add_text(s1, Inches(11.30), y + Inches(0.05), Inches(1.50), Inches(0.85),
             [n, "fans"], size=15, bold=True, color=NAVY,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.95)

# Footer note
add_text(s1, Inches(0.45), Inches(7.05), Inches(12.4), Inches(0.18),
         "Priority — first match wins:  Families > Multi-Gen > Couples > Social.  ~57-139 unique fans per game per segment — sufficient for ticket-purchase lift detection.",
         size=8.5, color=GRAY, align=PP_ALIGN.CENTER)
add_footer(s1, 1)


# =========================================================================
# Concept card builder — updated for 2 copy angles + 2 visuals
# =========================================================================
def segment_card_v2(slide, x, y, w, h, num, name, fans, profile_lines,
                    copy_angles, visual_lines):
    """
    copy_angles: list of 2 strings
    visual_lines: list of 2 strings
    """
    add_rect(slide, x, y, w, h, WHITE, line=NAVY)
    add_rect(slide, x, y, Inches(0.85), h, NAVY)
    add_text(slide, x, y + Inches(0.15), Inches(0.85), Inches(0.60),
             num, size=28, bold=True, color=YELLOW,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x, y + Inches(0.80), Inches(0.85), Inches(0.45),
             [fans, "FANS"], size=9, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    inner_x = x + Inches(1.00)
    inner_w = w - Inches(1.10)

    add_text(slide, inner_x, y + Inches(0.04), inner_w, Inches(0.28),
             name, size=13, bold=True, color=NAVY)
    add_text(slide, inner_x, y + Inches(0.32), inner_w, Inches(0.14),
             "WHO THEY ARE", size=7.5, bold=True, color=SKY)
    add_text(slide, inner_x, y + Inches(0.46), inner_w, Inches(0.44),
             profile_lines, size=8, color=GRAY)

    # Copy angles band (2 angles side by side)
    band_y = y + Inches(0.93)
    add_rect(slide, inner_x, band_y, inner_w, Inches(0.55), LIGHT)
    add_text(slide, inner_x + Inches(0.08), band_y + Inches(0.02), Inches(1.00), Inches(0.18),
             "COPY ANGLES", size=7, bold=True, color=NAVY)
    add_text(slide, inner_x + Inches(0.08), band_y + Inches(0.20), inner_w / 2 - Inches(0.15), Inches(0.32),
             "A:  " + copy_angles[0], size=8, bold=True, color=NAVY)
    add_text(slide, inner_x + inner_w / 2, band_y + Inches(0.20), inner_w / 2 - Inches(0.10), Inches(0.32),
             "B:  " + copy_angles[1], size=8, bold=True, color=NAVY)

    # Visual directions band (2 visuals stacked)
    band_y2 = y + Inches(1.52)
    add_rect(slide, inner_x, band_y2, inner_w, Inches(0.70), LIGHT)
    add_text(slide, inner_x + Inches(0.08), band_y2 + Inches(0.02), Inches(1.00), Inches(0.18),
             "VISUAL DIR.", size=7, bold=True, color=NAVY)
    add_text(slide, inner_x + Inches(0.08), band_y2 + Inches(0.20), inner_w - Inches(0.16), Inches(0.24),
             "A:  " + visual_lines[0], size=7.5, color=GRAY)
    add_text(slide, inner_x + Inches(0.08), band_y2 + Inches(0.44), inner_w - Inches(0.16), Inches(0.24),
             "B:  " + visual_lines[1], size=7.5, color=GRAY)


# Layout for concept cards
card_w = Inches(12.45)
concept_top = Inches(1.15)
concept_h = Inches(2.28)
concept_gap = Inches(0.10)

# =========================================================================
# SLIDE 2 — Concepts 1 & 2 (Families, Multi-Gen)
# =========================================================================
s2 = prs.slides.add_slide(BLANK)
add_header(s2, "Concepts 1 & 2  ·  Families with kids, Multi-Gen Reunion")

segment_card_v2(
    s2, Inches(0.45), concept_top, card_w, concept_h,
    "1", "Families with kids", "858",
    profile_lines=[
        "Median age 46  ·  84% rated 9-10  ·  Only 20% prior buyers  ·  High kid-friendly F&B spend",
        "Top foods:  Non-alc 66%  ·  Alcohol 50%  ·  Hot dog 43%  ·  Ice cream 37%  ·  Popcorn 30%  ·  Chicken 28%"
    ],
    copy_angles=[
        '"The face they make on a Home Run is worth the trip."',
        '"Ice cream drips. Popcorn flies. They won\'t remember the score — they\'ll remember this."'
    ],
    visual_lines=[
        "Kid in mini-helmet with ice cream cone, parent in matching jersey behind. Wide section shot with families in foreground.",
        "Parent and kid sharing a popcorn bucket, kid standing on seat, game lights in background. F&B-forward family moment."
    ],
)

mg_y = concept_top + concept_h + concept_gap
segment_card_v2(
    s2, Inches(0.45), mg_y, card_w, concept_h,
    "2", "Multi-Gen Reunion — adult kids or other family", "1,310",
    profile_lines=[
        "Median age 60  ·  57% are 55+  ·  84% rated 9-10  ·  25% prior buyers  ·  High avidity",
        "Top foods:  Non-alc 51%  ·  Alcohol 49%  ·  Hot dog 38%  ·  Pretzels 20%  ·  Popcorn 19%"
    ],
    copy_angles=[
        '"The team your parents raised you on. The one you\'ll bring them back to."',
        '"Same seats. New stories. Every generation adds a chapter."'
    ],
    visual_lines=[
        "VERTICAL age-diversity — adult son/daughter beside parent, matching jerseys, generational moment. NOT a friend-group shot.",
        "Three generations in a row — grandparent, parent, adult child — all reacting to same play. Pretzels and hot dogs in hand."
    ],
)

# Separator-rule callout (condensed to one line)
sep_y = mg_y + concept_h + Inches(0.08)
add_rect(s2, Inches(0.45), sep_y, card_w, Inches(0.55), LIGHT)
add_rect(s2, Inches(0.45), sep_y, Inches(0.15), Inches(0.55), YELLOW)
add_text(s2, Inches(0.75), sep_y + Inches(0.04), Inches(11.8), Inches(0.22),
         "SEPARATOR RULE — Multi-Gen vs Social  (apply to every Multi-Gen creative)",
         size=9, bold=True, color=NAVY)
add_text(s2, Inches(0.75), sep_y + Inches(0.26), Inches(11.8), Inches(0.26),
         ["Multi-Gen = vertical age diversity in frame (generational bond, heritage tone)  |  Social = horizontal age peers (lateral bond, present-tense energy)"],
         size=8, color=GRAY)
add_footer(s2, 2)


# =========================================================================
# SLIDE 3 — Concepts 3 & 4 (Couples, Social)
# =========================================================================
s3 = prs.slides.add_slide(BLANK)
add_header(s3, "Concepts 3 & 4  ·  Couples, Social / Crew")

segment_card_v2(
    s3, Inches(0.45), concept_top, card_w, concept_h,
    "3", "Couples — partner attended (any age)", "2,016",
    profile_lines=[
        "Median age 64  ·  67% are 55+  ·  82% rated 9-10  ·  Only 27% prior buyers  ·  High disposable income",
        "Top foods:  Alcohol 50%  ·  Non-alc 44%  ·  Hot dog 38%  ·  Pretzels 22%  ·  Nuts 18%  ·  Fries 15%"
    ],
    copy_angles=[
        '"Catch a game with your favorite teammate."',
        '"Two cold ones, one perfect night. No reservations required."'
    ],
    visual_lines=[
        "Two beers + shared pretzel on the seat tray, lights illuminating over the Trop. Imagery should read 50-65 to match audience.",
        "Close-up of two beers clinking with sunset field in background. Couple leaning into each other. Date-night energy, 50s-60s read."
    ],
)

soc_y = concept_top + concept_h + concept_gap
segment_card_v2(
    s3, Inches(0.45), soc_y, card_w, concept_h,
    "4", "Social / Crew — friends, business, solo, other (catch-all)", "1,845",
    profile_lines=[
        "Median age 62  ·  Only 17% prior buyers (lowest)  ·  Lowest avidity  ·  Biggest incremental upside",
        "Top foods:  Alcohol 32%  ·  Non-alc 25%  ·  Hot dog 22%  ·  Pretzels 10%  ·  Popcorn 9%  ·  Lowest F&B engagement overall"
    ],
    copy_angles=[
        '"Round up the crew. The story isn\'t over."',
        '"The best group plans start with \'who\'s in?\'"'
    ],
    visual_lines=[
        "HORIZONTAL age peers — friend group same cohort, cheers/group reaction, beers up, action on field. Energy and present-tense.",
        "Phone screen showing a group chat with 'I'm in' replies, stadium blurred in background. Social-proof trigger, FOMO energy."
    ],
)

add_footer(s3, 3)


# =========================================================================
# Build-spec helpers (slides 4-5)
# =========================================================================
SQL_FONT = "Consolas"

def add_audience_base_strip(slide):
    x, y, w, h = Inches(0.45), Inches(1.15), Inches(12.45), Inches(1.05)
    add_rect(slide, x, y, w, h, LIGHT, line=NAVY)
    add_rect(slide, x, y, Inches(2.30), h, NAVY)
    add_text(slide, x + Inches(0.10), y, Inches(2.20), h,
             ["AUDIENCE BASE", "(applied to all 4)"],
             size=10, bold=True, color=YELLOW,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(2.45), y + Inches(0.04), w - Inches(2.55), h - Inches(0.08),
             ["FROM `mlb-dataeng-prod.wheelhouse_rays.qualtrics_voc_post_attendance_full`",
              "WHERE  season IN (2026)",
              "  AND  OVERALL_NUMRAT IN (8, 9, 10)",
              "  AND  COALESCE(TB_ADDON_6, 0) != 1",
              "  AND  BUYER_TYPE NOT IN ('Full Season Ticket','Sponsor','Suncoast Credit Union Member Offer',",
              "       'Flexible Season Member Ticket','Season Member Additional Ticket','Partial Season Discount',",
              "       'Half Season Discount','Full Season Discount','Exchange Sponsor','Sponsor Suite Blue','Sponsor Suite Yellow')"],
             size=8, color=NAVY, font=SQL_FONT)


def add_build_card(slide, x, y, w, h, num, name, fans, sql_lines, columns):
    add_rect(slide, x, y, w, h, WHITE, line=NAVY)
    add_rect(slide, x, y, w, Inches(0.40), NAVY)
    add_rect(slide, x, y, Inches(0.55), Inches(0.40), YELLOW)
    add_text(slide, x, y, Inches(0.55), Inches(0.40),
             num, size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + Inches(0.65), y, w - Inches(2.0), Inches(0.40),
             name, size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, x + w - Inches(1.45), y, Inches(1.40), Inches(0.40),
             f"{fans} fans", size=11, bold=True, color=YELLOW,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    body_y = y + Inches(0.45)
    body_h = h - Inches(0.50)

    sql_w = Inches(7.30)
    add_rect(slide, x + Inches(0.10), body_y, sql_w, Inches(0.22), CODEBG)
    add_text(slide, x + Inches(0.18), body_y, sql_w, Inches(0.22),
             "ADDITIONAL FILTERS  (in addition to the audience base above)",
             size=8, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, x + Inches(0.10), body_y + Inches(0.22), sql_w, body_h - Inches(0.22), CODEBG)
    add_text(slide, x + Inches(0.18), body_y + Inches(0.24), sql_w - Inches(0.16), body_h - Inches(0.26),
             sql_lines, size=8, color=NAVY, font=SQL_FONT)

    col_x = x + Inches(7.55)
    col_w = w - Inches(7.65)
    add_rect(slide, col_x, body_y, col_w, Inches(0.22), CODEBG)
    add_text(slide, col_x + Inches(0.10), body_y, col_w, Inches(0.22),
             "KEY COLUMNS  (segment-specific)",
             size=8, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, col_x, body_y + Inches(0.22), col_w, body_h - Inches(0.22), CODEBG)
    add_text(slide, col_x + Inches(0.10), body_y + Inches(0.24), col_w - Inches(0.20), body_h - Inches(0.26),
             columns, size=8.5, color=NAVY, font=SQL_FONT)


def add_always_needed_strip(slide):
    y = Inches(7.02)
    add_text(slide, Inches(0.45), y, Inches(12.4), Inches(0.18),
             ["ALWAYS-NEEDED COLUMNS:  ATTENDING_ID  ·  EMAIL  ·  FIRST_NAME  ·  LAST_NAME  ·  CITY  ·  STATE  ·  AGE  ·  OVERALL_NUMRAT  ·  TB_ADDON_6  ·  BUYER_TYPE  ·  season"],
             size=8, bold=True, color=GRAY, align=PP_ALIGN.CENTER)


SQL_S1 = [
    "AND (",
    "      COALESCE(ATTEND_WITH_CATEGORY_HS_KIDS, 0)     = 1",
    "   OR COALESCE(ATTEND_WITH_CATEGORY_NON_HS_KIDS, 0) = 1",
    ")",
]
COL_S1 = [
    "Inclusion flags",
    "  · ATTEND_WITH_CATEGORY_HS_KIDS",
    "  · ATTEND_WITH_CATEGORY_NON_HS_KIDS",
    "",
    "Optional creative inputs",
    "  · ATTEND_KIDS_AGES_*  (kid age buckets)",
    "  · CONCESS_TYPE_ICECREAM  (over-indexes)",
    "  · CONCESS_TYPE_CHICKEN   (over-indexes)",
]

SQL_S2 = [
    "AND COALESCE(ATTEND_WITH_CATEGORY_HS_KIDS, 0)     != 1",
    "AND COALESCE(ATTEND_WITH_CATEGORY_NON_HS_KIDS, 0) != 1",
    "AND (",
    "      COALESCE(ATTEND_WITH_CATEGORY_ADULT_KIDS, 0) = 1",
    "   OR COALESCE(ATTEND_WITH_CATEGORY_OTHERFAM, 0)   = 1",
    ")",
]
COL_S2 = [
    "Inclusion flags",
    "  · ATTEND_WITH_CATEGORY_ADULT_KIDS",
    "  · ATTEND_WITH_CATEGORY_OTHERFAM",
    "",
    "Priority exclusions  (Segment 1)",
    "  · ATTEND_WITH_CATEGORY_HS_KIDS",
    "  · ATTEND_WITH_CATEGORY_NON_HS_KIDS",
]

SQL_S3 = [
    "AND COALESCE(ATTEND_WITH_CATEGORY_HS_KIDS, 0)     != 1",
    "AND COALESCE(ATTEND_WITH_CATEGORY_NON_HS_KIDS, 0) != 1",
    "AND COALESCE(ATTEND_WITH_CATEGORY_ADULT_KIDS, 0)  != 1",
    "AND COALESCE(ATTEND_WITH_CATEGORY_OTHERFAM, 0)    != 1",
    "AND COALESCE(ATTEND_WITH_CATEGORY_SPOUSE, 0)       = 1",
]
COL_S3 = [
    "Inclusion flag",
    "  · ATTEND_WITH_CATEGORY_SPOUSE",
    "",
    "Priority exclusions  (Segments 1-2)",
    "  · ATTEND_WITH_CATEGORY_HS_KIDS",
    "  · ATTEND_WITH_CATEGORY_NON_HS_KIDS",
    "  · ATTEND_WITH_CATEGORY_ADULT_KIDS",
    "  · ATTEND_WITH_CATEGORY_OTHERFAM",
]

SQL_S4 = [
    "AND COALESCE(ATTEND_WITH_CATEGORY_HS_KIDS, 0)     != 1",
    "AND COALESCE(ATTEND_WITH_CATEGORY_NON_HS_KIDS, 0) != 1",
    "AND COALESCE(ATTEND_WITH_CATEGORY_ADULT_KIDS, 0)  != 1",
    "AND COALESCE(ATTEND_WITH_CATEGORY_OTHERFAM, 0)    != 1",
    "AND COALESCE(ATTEND_WITH_CATEGORY_SPOUSE, 0)      != 1",
    "-- catch-all: friends / business / alone / other / unknown",
]
COL_S4 = [
    "No inclusion flag — pure catch-all",
    "(every fan that fails Segments 1-3 lands here)",
    "",
    "Priority exclusions  (Segments 1-3)",
    "  · all 4 from Segment 3, plus:",
    "  · ATTEND_WITH_CATEGORY_SPOUSE",
]


bs_top   = Inches(2.30)
bs_h     = Inches(2.30)
bs_gap   = Inches(0.10)

# =========================================================================
# SLIDE 4 — Build spec for Segments 1 & 2
# =========================================================================
s4 = prs.slides.add_slide(BLANK)
add_header(s4, "Build spec  ·  Segments 1 & 2  ·  Families, Multi-Gen Reunion")
add_audience_base_strip(s4)

add_build_card(s4, Inches(0.45), bs_top, Inches(12.45), bs_h,
               "1", "Families (kids)", "858", SQL_S1, COL_S1)
add_build_card(s4, Inches(0.45), bs_top + bs_h + bs_gap, Inches(12.45), bs_h,
               "2", "Multi-Gen Reunion", "1,310", SQL_S2, COL_S2)
add_always_needed_strip(s4)
add_footer(s4, 4)


# =========================================================================
# SLIDE 5 — Build spec for Segments 3 & 4
# =========================================================================
s5 = prs.slides.add_slide(BLANK)
add_header(s5, "Build spec  ·  Segments 3 & 4  ·  Couples, Social / Crew")
add_audience_base_strip(s5)

add_build_card(s5, Inches(0.45), bs_top, Inches(12.45), bs_h,
               "3", "Couples", "2,016", SQL_S3, COL_S3)
add_build_card(s5, Inches(0.45), bs_top + bs_h + bs_gap, Inches(12.45), bs_h,
               "4", "Social / Crew  (catch-all)", "1,845", SQL_S4, COL_S4)
add_always_needed_strip(s5)
add_footer(s5, 5)


prs.save(OUT)
print(f"Wrote: {OUT}")
print(f"Slides: {len(prs.slides)}  ·  Size: 13.33 x 7.5 in (16:9)")
