"""
Homestand #1 Fan Experience Report Generator
Generates a PowerPoint deck comparing 2026 Homestand 1 (Apr 6-12) vs 2024 Season averages
with sentence-level qualitative feedback per department and NPS tier.
"""

import snowflake.connector
import pandas as pd
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import os

# ── Rays Dark Theme Colors ──────────────────────────────────────────────────
NAVY = RGBColor(0x09, 0x2C, 0x5C)
COLUMBIA_BLUE = RGBColor(0x8F, 0xBC, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xFF, 0x6B, 0x6B)
GOLD = RGBColor(0xF5, 0xC5, 0x18)
MEDIUM_BLUE = RGBColor(0x0D, 0x3B, 0x7A)
DARK_BLUE_BG = RGBColor(0x07, 0x20, 0x44)

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FILE = os.path.join(OUTPUT_DIR, "Homestand1_Report_2026_v7.pptx")

# ── Snowflake Connection ────────────────────────────────────────────────────
def get_connection():
    return snowflake.connector.connect(
        account='hta92307.east-us-2.azure',
        user='YTAKETANI@RAYSBASEBALL.COM',
        authenticator='externalbrowser',
        role='ACCOUNTADMIN',
        database='TBRDP_DW_DEV',
        schema='IM_RPT',
        warehouse='TBRDP_DW_CORTEX_XS_WH',
        client_store_temporary_credential=True,
    )

# ── Data Queries ────────────────────────────────────────────────────────────

SURVEY_COLS = """
    QUALTRICS_ID, GAME_DATE, SEASON, BUYER_TYPE, OVERALL_NUMRAT,
    CONCESS_NUMRAT, PARKING_NUMRAT, ENTERTAIN_NUMRAT, MERCH_NUMRAT,
    SEATVIEW_NUMRAT, STAFF_NUMRAT,
    -- Grid: Concessions
    CONCESS_GRID_VALUE_DESC, CONCESS_GRID_SELECTION_DESC,
    CONCESS_GRID_CLEAN_DESC,
    -- Grid: Merch
    MERCH_GRID_MERCHQUALITY_DESC, MERCH_GRID_SELECTION_DESC,
    MERCH_GRID_PRICE_DESC, MERCH_GRID_WAIT_DESC,
    -- Grid: Entertainment
    ENTERTAIN_GRID_MUSIC_DESC, ENTERTAIN_GRID_SCOREBOARD_DESC,
    ENTERTAIN_GRID_GAMES_DESC, ENTERTAIN_GRID_KIDS_ACTIVITIES_DESC,
    ENTERTAIN_GRID_PREGAME_CONTENT_DESC,
    -- Parking extras
    CLUB_PARKING_LOT_DESC,
    -- Security / Gate Entry
    GE_TIME_EXPECT_DESC,
    -- Concession extras
    CONCESS_SCREENER_DESC, CONCESS_WAIT_DESC, CONCESS_SPEND_DESC,
    CONCESS_WAIT_EXPECT_DESC, CONCESS_INTENT_NO_DESC, CONCESS_NO_DESC,
    -- Concession quality (17 food items)
    CONCESS_QUALITY_BURGERS_DESC, CONCESS_QUALITY_HOTDOG_DESC,
    CONCESS_QUALITY_PIZZA_DESC, CONCESS_QUALITY_CHICKEN_DESC,
    CONCESS_QUALITY_NACHOS_DESC, CONCESS_QUALITY_FRIES_DESC,
    CONCESS_QUALITY_POPCORN_DESC, CONCESS_QUALITY_PRETZELS_DESC,
    CONCESS_QUALITY_SANDWICH_DESC, CONCESS_QUALITY_SAUSAGE_DESC,
    CONCESS_QUALITY_ICECREAM_DESC, CONCESS_QUALITY_ALCOHOL_DESC,
    CONCESS_QUALITY_NONALCOHOL_DESC, CONCESS_QUALITY_NUTS_DESC,
    CONCESS_QUALITY_OTHER_DESSERT_DESC, CONCESS_QUALITY_OTHER_ENTREE_DESC,
    -- Merch extras
    MERCH_SCREENER_DESC, MERCH_NO_DESC,
    -- Seat view
    SEATVIEW_VALUE_SEAT,
    -- Ticketing
    PURCHASE_INTENT_DESC, PURCHASE_INTENT_NEXT_DESC, PREVIOUS_PURCHASE_DESC,
    INCENTIVES_RANK_1_DESC,
    -- Promotions
    GIVEAWAY_SAT_DESC, GIVEAWAY_RECEIVE_DESC,
    PA_PROMO_GRID_GIVEAWAY_DESC, PA_PROMO_GRID_THEME_DESC,
    PA_GIVEAWAY_INTENT_DESC,
    -- Brand health (12 items)
    BRANDHEALTH_GRID_EXCITING_DESC, BRANDHEALTH_GRID_FAMFRIENDLY_DESC,
    BRANDHEALTH_GRID_RIGHTDIRECTION_DESC, BRANDHEALTH_GRID_POSINFLUENCE_DESC,
    BRANDHEALTH_GRID_WELCOME_DESC, BRANDHEALTH_GRID_SAFE_DESC,
    BRANDHEALTH_GRID_ACCESSIBLE_DESC, BRANDHEALTH_GRID_SUSTAINABILITY_DESC,
    BRANDHEALTH_GRID_TRENDY_DESC, BRANDHEALTH_GRID_EMOTIONAL_DESC,
    BRANDHEALTH_GRID_CHAMPION_DESC, BRANDHEALTH_GRID_DIVERSITY_DESC,
    -- Travel
    TRAVELTO_METHOD_DESC, TRAVELTO_TIME_EXPECT_DESC,
    -- Staff type interacted with (text column, comma-separated)
    STAFF_TYPE,
    -- Staff type flags (may be NULL - kept for backwards compatibility)
    STAFF_TYPE_CONCESSIONS, STAFF_TYPE_SECURITY, STAFF_TYPE_USHER,
    STAFF_TYPE_PARKING, STAFF_TYPE_FAN_SERVICES, STAFF_TYPE_MERCH,
    STAFF_TYPE_TB, STAFF_TYPE_MIN_TB, STAFF_TYPE_NONE,
    -- Demographics
    AGE, GENDER_ID_DESC, HHI_ID_DESC, PARENT_ID_DESC,
    CIH_ID_DESC, TEAM_AVIDITY_DESC, COMMERCIAL_OPT_IN_DESC,
    -- Attendance
    ATTEND_NUM_PLAN_DESC,
    ATTEND_WITH_CATEGORY_FRIENDS, ATTEND_WITH_CATEGORY_SPOUSE,
    ATTEND_WITH_CATEGORY_ADULT_KIDS, ATTEND_WITH_CATEGORY_ALONE,
    ATTEND_WITH_CATEGORY_BUSINESS, ATTEND_WITH_CATEGORY_OTHERFAM,
    ATTEND_WITH_CATEGORY_OTHER
"""

def load_data():
    """Load all data from Snowflake into pandas DataFrames."""
    conn = get_connection()
    print("Connected to Snowflake.")

    # Homestand 1 data
    hs1 = pd.read_sql(f"""
        SELECT {SURVEY_COLS}
        FROM TBRDP_DW_DEV.IM_RPT.V_SBL_QUALTRICS_VOC_POST_ATTENDANCE_FULL_CORTEX_AI
        WHERE GAME_DATE::DATE BETWEEN '2026-04-06' AND '2026-04-12'
          AND OVERALL_NUMRAT IS NOT NULL AND OVERALL_NUMRAT < 80
    """, conn)
    print(f"Homestand 1: {len(hs1)} responses")

    # 2024 season data
    s24 = pd.read_sql(f"""
        SELECT {SURVEY_COLS}
        FROM TBRDP_DW_DEV.IM_RPT.V_SBL_QUALTRICS_VOC_POST_ATTENDANCE_FULL_CORTEX_AI
        WHERE SEASON = 2024 AND OVERALL_NUMRAT IS NOT NULL AND OVERALL_NUMRAT < 80
    """, conn)
    print(f"2024 Season: {len(s24)} responses")

    # Sentence-level feedback
    sentences = pd.read_sql("""
        SELECT QUALTRICS_ID, SATISFACTION_RATING, NPS_SEGMENT,
               SENTENCE_TEXT, SENTIMENT_CATEGORY, AI_CATEGORY,
               PARENT_CATEGORY, SENTENCE_LENGTH
        FROM TBRDP_DW_DEV.IM_RPT.T_OVERALL_FEEDBACK_SENTENCE_LEVEL
        WHERE GAME_DATE::DATE BETWEEN '2026-04-06' AND '2026-04-12'
    """, conn)
    print(f"Sentences: {len(sentences)} rows")

    # Generate AI theme summaries via Cortex
    print("Generating AI theme summaries via Cortex...")
    theme_summaries = load_theme_summaries(conn)
    print(f"Theme summaries: {len(theme_summaries)} generated")

    conn.close()
    return hs1, s24, sentences, theme_summaries


def load_theme_summaries(conn):
    """Batch-generate narrative theme summaries using Snowflake Cortex AI.

    For each top-3 AI_CATEGORY per (PARENT_CATEGORY, NPS_SEGMENT), aggregates
    the top sentences and calls CORTEX.COMPLETE to produce a 1-2 sentence
    narrative summary. Returns a dict keyed by (parent_cat, nps_segment, ai_category).
    """
    query = """
    WITH ranked_themes AS (
        SELECT
            PARENT_CATEGORY,
            NPS_SEGMENT,
            AI_CATEGORY,
            COUNT(*) AS mention_count,
            ROW_NUMBER() OVER (
                PARTITION BY PARENT_CATEGORY, NPS_SEGMENT
                ORDER BY COUNT(*) DESC
            ) AS rn
        FROM TBRDP_DW_DEV.IM_RPT.T_OVERALL_FEEDBACK_SENTENCE_LEVEL
        WHERE GAME_DATE::DATE BETWEEN '2026-04-06' AND '2026-04-12'
          AND LEN(AI_CATEGORY) <= 40
        GROUP BY PARENT_CATEGORY, NPS_SEGMENT, AI_CATEGORY
    ),
    top_themes AS (
        SELECT PARENT_CATEGORY, NPS_SEGMENT, AI_CATEGORY, mention_count
        FROM ranked_themes
        WHERE rn <= 3
    ),
    theme_comments AS (
        SELECT
            t.PARENT_CATEGORY,
            t.NPS_SEGMENT,
            t.AI_CATEGORY,
            t.mention_count,
            LISTAGG('- ' || LEFT(s.SENTENCE_TEXT, 180), '\\n')
                WITHIN GROUP (ORDER BY s.SENTENCE_LENGTH DESC) AS combined_comments
        FROM top_themes t
        JOIN (
            SELECT *, ROW_NUMBER() OVER (
                PARTITION BY PARENT_CATEGORY, NPS_SEGMENT, AI_CATEGORY
                ORDER BY SENTENCE_LENGTH DESC
            ) AS sent_rn
            FROM TBRDP_DW_DEV.IM_RPT.T_OVERALL_FEEDBACK_SENTENCE_LEVEL
            WHERE GAME_DATE::DATE BETWEEN '2026-04-06' AND '2026-04-12'
              AND LEN(AI_CATEGORY) <= 40
        ) s
          ON s.PARENT_CATEGORY = t.PARENT_CATEGORY
         AND s.NPS_SEGMENT = t.NPS_SEGMENT
         AND s.AI_CATEGORY = t.AI_CATEGORY
         AND s.sent_rn <= 15
        GROUP BY t.PARENT_CATEGORY, t.NPS_SEGMENT, t.AI_CATEGORY, t.mention_count
    )
    SELECT
        PARENT_CATEGORY,
        NPS_SEGMENT,
        AI_CATEGORY,
        mention_count,
        SNOWFLAKE.CORTEX.COMPLETE(
            'claude-sonnet-4-6',
            'You analyze Tampa Bay Rays fan survey feedback at Tropicana Field. '
            || 'Below are comments from ' || NPS_SEGMENT || 's (NPS '
            || CASE NPS_SEGMENT
                 WHEN 'Detractor' THEN '0-6'
                 WHEN 'Passive' THEN '7-8'
                 ELSE '9-10'
               END || ') about "' || AI_CATEGORY || '". '
            || 'Write exactly 1-2 sentences summarizing the most common specific '
            || CASE NPS_SEGMENT
                 WHEN 'Detractor' THEN 'complaints'
                 WHEN 'Passive' THEN 'feedback and concerns'
                 ELSE 'praises'
               END || '. '
            || 'Be concrete and specific — mention specific staff roles, locations, '
            || 'behaviors, or issues fans referenced. '
            || 'Do not use bullet points, quotation marks, or generic language.'
            || '\\n\\nComments:\\n' || LEFT(combined_comments, 3500)
            || '\\n\\nSummary:'
        ) AS ai_summary
    FROM theme_comments
    """
    df = pd.read_sql(query, conn)
    summaries = {}
    for _, row in df.iterrows():
        key = (row['PARENT_CATEGORY'], row['NPS_SEGMENT'], row['AI_CATEGORY'])
        summary_text = str(row['AI_SUMMARY']).strip().strip('"').strip()
        summaries[key] = summary_text
    return summaries


def clean_desc_columns(*dfs):
    """Strip numeric special codes (e.g. 85, 84, -10) from all _DESC columns.

    Snowflake sometimes stores numeric codes alongside text labels in _DESC
    columns, causing pandas mixed dtype issues. This converts any value that
    looks numeric to NaN so downstream string matching works correctly.
    Also cleans SEATVIEW_VALUE_SEAT which has the same issue.
    """
    target_cols = None  # compute once
    for df in dfs:
        if target_cols is None:
            target_cols = [c for c in df.columns
                          if c.endswith('_DESC') or c == 'SEATVIEW_VALUE_SEAT']
        for col in target_cols:
            if col not in df.columns:
                continue
            # Convert column to string first, then replace numeric-looking values
            s = df[col].copy()
            # Identify values that are numeric (int or float stored as object)
            numeric_mask = pd.to_numeric(s, errors='coerce').notna() & s.notna()
            df.loc[numeric_mask, col] = np.nan


# ── Helper Functions ────────────────────────────────────────────────────────

def safe_avg(series):
    """Average of non-null values below 80 (filter out special codes)."""
    valid = series.dropna()
    valid = valid[valid < 80]
    return round(valid.mean(), 2) if len(valid) > 0 else None

def safe_count(series):
    """Count of non-null values below 80."""
    valid = series.dropna()
    return int((valid < 80).sum())

def delta_pct(new_val, old_val):
    """Calculate percentage change."""
    if old_val is None or new_val is None or old_val == 0:
        return None
    return round(((new_val - old_val) / old_val) * 100, 1)

def nps_tier(rating):
    """Classify a 0-10 rating into NPS tier."""
    if rating is None or pd.isna(rating):
        return None
    if rating <= 6:
        return 'Detractor'
    elif rating <= 8:
        return 'Passive'
    else:
        return 'Promoter'

def grid_distribution(series, valid_responses=None):
    """Calculate % distribution for a satisfaction grid column.
    Returns dict with Highly satisfied, Somewhat satisfied, etc."""
    if valid_responses is None:
        valid_responses = ['Highly satisfied', 'Somewhat satisfied',
                          'Somewhat dissatisfied', 'Highly dissatisfied']
    s = series.dropna()
    s = s[s.isin(valid_responses)]
    total = len(s)
    if total == 0:
        return {r: 0.0 for r in valid_responses}, 0
    dist = {}
    for r in valid_responses:
        dist[r] = round(100.0 * (s == r).sum() / total, 1)
    return dist, total

def categorical_distribution(series, categories):
    """Calculate % distribution for categorical columns."""
    s = series.dropna()
    s = s[s != 'N/A']
    total = len(s)
    if total == 0:
        return {c: 0.0 for c in categories}, 0
    dist = {}
    for c in categories:
        dist[c] = round(100.0 * (s == c).sum() / total, 1)
    return dist, total


def multiselect_distribution(df, flag_cols, labels=None):
    """Calculate % distribution for multi-select flag columns.
    flag_cols: list of column names (binary 0/1 flags)
    labels: optional dict mapping col name -> display label
    Returns list of (label, pct, count) tuples sorted by count desc, and total respondents.
    """
    if labels is None:
        labels = {c: c.split('_')[-1].title() for c in flag_cols}
    # Total respondents = any row with at least one non-null flag
    mask = df[flag_cols].notna().any(axis=1)
    total = mask.sum()
    if total == 0:
        return [], 0
    results = []
    for col in flag_cols:
        count = int((df[col].fillna(0) > 0).sum())
        pct = round(100.0 * count / total, 1) if total > 0 else 0
        results.append((labels.get(col, col), pct, count))
    results.sort(key=lambda x: x[2], reverse=True)
    return results, total


def text_multiselect_distribution(series, categories, labels=None):
    """Calculate % distribution from a comma-separated text column.

    series: pandas Series containing comma-separated values (e.g. "parking, usher, concessions")
    categories: list of possible category keys (lowercase, matching the text values)
    labels: optional dict mapping category key -> display label
    Returns list of (label, pct, count) tuples sorted by count desc, and total respondents.
    """
    if labels is None:
        labels = {c: c.replace('_', ' ').title() for c in categories}
    s = series.dropna()
    s = s[s.str.strip() != '']
    total = len(s)
    if total == 0:
        return [], 0
    results = []
    for cat in categories:
        # Check if category appears in the comma-separated string
        count = int(s.str.contains(r'(?:^|,)\s*' + cat + r'\s*(?:,|$)', case=False, regex=True).sum())
        pct = round(100.0 * count / total, 1)
        results.append((labels.get(cat, cat), pct, count))
    results.sort(key=lambda x: x[2], reverse=True)
    return results, total


def age_bucket_distribution(series):
    """Bucket numeric AGE into ranges and return distribution.
    Returns dict of bucket: pct, and total count.
    """
    s = series.dropna()
    s = s[s > 0]  # filter invalid
    total = len(s)
    if total == 0:
        return {}, 0
    buckets = [
        ('Under 21', 0, 20),
        ('21-30', 21, 30),
        ('31-40', 31, 40),
        ('41-50', 41, 50),
        ('51-60', 51, 60),
        ('61+', 61, 200),
    ]
    dist = {}
    for label, low, high in buckets:
        count = ((s >= low) & (s <= high)).sum()
        dist[label] = round(100.0 * count / total, 1)
    return dist, total

def get_top_themes(sentences_df, parent_cats=None, ai_cats=None,
                   nps_segment=None, rating_col_tiers=None,
                   survey_df=None, rating_col=None, top_n=5):
    """Get top AI_CATEGORY themes from sentence-level data.

    Can filter by:
    - parent_cats: list of PARENT_CATEGORY values
    - ai_cats: list of specific AI_CATEGORY values
    - nps_segment: 'Detractor', 'Passive', 'Promoter'
    - rating_col_tiers: dict mapping rating column to tier bounds,
                        e.g. {'col': 'CONCESS_NUMRAT', 'min': 0, 'max': 6}
                        used with survey_df to get QUALTRICS_IDs in that tier
    """
    df = sentences_df.copy()

    # Filter by NPS segment
    if nps_segment:
        df = df[df['NPS_SEGMENT'] == nps_segment]

    # Filter by parent category
    if parent_cats:
        df = df[df['PARENT_CATEGORY'].isin(parent_cats)]

    # Filter by specific AI categories
    if ai_cats:
        df = df[df['AI_CATEGORY'].isin(ai_cats)]

    # Filter by department-specific rating tier
    if rating_col_tiers and survey_df is not None and rating_col:
        tier_min = rating_col_tiers['min']
        tier_max = rating_col_tiers['max']
        valid = survey_df[
            (survey_df[rating_col].notna()) &
            (survey_df[rating_col] < 80) &
            (survey_df[rating_col] >= tier_min) &
            (survey_df[rating_col] <= tier_max)
        ]['QUALTRICS_ID']
        df = df[df['QUALTRICS_ID'].isin(valid)]

    if len(df) == 0:
        return pd.DataFrame(columns=['AI_CATEGORY', 'count', 'pct']), 0

    total_in_tier = len(df)
    themes = (df.groupby('AI_CATEGORY')
              .size()
              .reset_index(name='count')
              .sort_values('count', ascending=False)
              .head(top_n))
    themes['pct'] = (themes['count'] / total_in_tier * 100).round(1)
    return themes, total_in_tier

def get_verbatims(sentences_df, parent_cats=None, ai_cats=None,
                  nps_segment=None, sentiment=None,
                  survey_df=None, rating_col=None, rating_col_tiers=None,
                  top_n=3):
    """Get representative verbatim quotes."""
    df = sentences_df.copy()

    if nps_segment:
        df = df[df['NPS_SEGMENT'] == nps_segment]
    if parent_cats:
        df = df[df['PARENT_CATEGORY'].isin(parent_cats)]
    if ai_cats:
        df = df[df['AI_CATEGORY'].isin(ai_cats)]
    if sentiment:
        df = df[df['SENTIMENT_CATEGORY'] == sentiment]
    if rating_col_tiers and survey_df is not None and rating_col:
        tier_min = rating_col_tiers['min']
        tier_max = rating_col_tiers['max']
        valid = survey_df[
            (survey_df[rating_col].notna()) &
            (survey_df[rating_col] < 80) &
            (survey_df[rating_col] >= tier_min) &
            (survey_df[rating_col] <= tier_max)
        ]['QUALTRICS_ID']
        df = df[df['QUALTRICS_ID'].isin(valid)]

    if len(df) == 0:
        return []

    # Pick longest, most detailed sentences
    df = df.sort_values('SENTENCE_LENGTH', ascending=False)
    # Deduplicate by taking one sentence per QUALTRICS_ID
    df = df.drop_duplicates(subset='QUALTRICS_ID')
    return df['SENTENCE_TEXT'].head(top_n).tolist()


def get_theme_summary(theme_summaries, parent_cats, nps_segment, ai_category,
                      rating_col=None, survey_df=None, rating_col_tiers=None):
    """Look up a pre-generated Cortex AI narrative summary for a theme.

    Searches the theme_summaries dict by (parent_category, nps_segment, ai_category).
    Since a department may map to multiple parent_categories, tries each one.
    """
    if not parent_cats:
        return "No summary available."

    for pc in parent_cats:
        key = (pc, nps_segment, ai_category)
        if key in theme_summaries:
            return theme_summaries[key]

    return "No summary available for this theme."


# ── Slide Building Helpers ──────────────────────────────────────────────────

def set_slide_bg(slide, color=NAVY):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=12,
                 bold=False, color=WHITE, alignment=PP_ALIGN.LEFT,
                 font_name='Calibri'):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_scorecard_table(slide, left, top, width, rows_data, col_widths=None):
    """Add a scorecard table to a slide.
    rows_data: list of dicts with keys: metric, hs1_val, s24_val, delta, n
    """
    num_rows = len(rows_data) + 1  # +1 for header
    num_cols = 5  # Metric, 2024 Avg, 2026 HS1, Delta %, n
    table_height = Inches(0.3 * num_rows)

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, table_height)
    table = table_shape.table

    # Set column widths
    if col_widths is None:
        col_widths = [Inches(2.8), Inches(1.1), Inches(1.1), Inches(1.1), Inches(0.7)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    headers = ['Metric', '2024 Avg', '2026 HS1', 'Delta %', 'n']
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = MEDIUM_BLUE

    # Data rows
    for r_idx, row in enumerate(rows_data):
        row_num = r_idx + 1
        values = [
            row.get('metric', ''),
            f"{row['s24_val']:.2f}" if row.get('s24_val') is not None else 'N/A',
            f"{row['hs1_val']:.2f}" if row.get('hs1_val') is not None else 'N/A',
            f"{row['delta']:+.1f}%" if row.get('delta') is not None else 'N/A',
            str(row.get('n', '')) if row.get('n') else '',
        ]
        for c_idx, val in enumerate(values):
            cell = table.cell(row_num, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.name = 'Calibri'
                p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
                # Color delta column
                if c_idx == 3 and row.get('delta') is not None:
                    if row['delta'] > 0:
                        p.font.color.rgb = GREEN
                    elif row['delta'] < 0:
                        p.font.color.rgb = RED
                    else:
                        p.font.color.rgb = GOLD
                else:
                    p.font.color.rgb = WHITE
            # Alternate row background
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE_BG if r_idx % 2 == 0 else NAVY

    return table_shape


def add_grid_table(slide, left, top, width, rows_data):
    """Add a satisfaction grid distribution table.
    rows_data: list of dicts with keys: metric, highly_sat, somewhat_sat,
               somewhat_dissat, highly_dissat, n
    """
    num_rows = len(rows_data) + 1
    num_cols = 6
    table_height = Inches(0.28 * num_rows)

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, table_height)
    table = table_shape.table

    col_widths = [Inches(2.2), Inches(1.1), Inches(1.1), Inches(1.2), Inches(1.2), Inches(0.5)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    headers = ['Metric', '% Highly\nSatisfied', '% Somewhat\nSatisfied',
               '% Somewhat\nDissatisfied', '% Highly\nDissatisfied', 'n']
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Calibri'
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = MEDIUM_BLUE

    for r_idx, row in enumerate(rows_data):
        row_num = r_idx + 1
        vals = [
            row.get('metric', ''),
            f"{row.get('highly_sat', 0):.1f}%",
            f"{row.get('somewhat_sat', 0):.1f}%",
            f"{row.get('somewhat_dissat', 0):.1f}%",
            f"{row.get('highly_dissat', 0):.1f}%",
            str(row.get('n', '')),
        ]
        for c_idx, val in enumerate(vals):
            cell = table.cell(row_num, c_idx)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.name = 'Calibri'
                p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
                p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE_BG if r_idx % 2 == 0 else NAVY

    return table_shape


def add_grid_bar_chart(slide, left, top, width, height, grid_rows, comp_rows=None):
    """Add a horizontal bar chart showing satisfaction distribution.

    grid_rows: list of dicts from build_grid_row() with keys:
        metric, highly_sat, somewhat_sat, somewhat_dissat, highly_dissat, n
    comp_rows: optional list of dicts from build_grid_comparison_row() with keys:
        metric, hs1_val, s24_val, delta, n
    """
    chart_data = CategoryChartData()
    chart_data.categories = [r['metric'] for r in grid_rows]

    # Series in order: Highly Satisfied → Highly Dissatisfied
    chart_data.add_series('Highly Satisfied',
                          [r.get('highly_sat', 0) for r in grid_rows])
    chart_data.add_series('Somewhat Satisfied',
                          [r.get('somewhat_sat', 0) for r in grid_rows])
    chart_data.add_series('Somewhat Dissatisfied',
                          [r.get('somewhat_dissat', 0) for r in grid_rows])
    chart_data.add_series('Highly Dissatisfied',
                          [r.get('highly_dissat', 0) for r in grid_rows])

    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True

    # Style the chart for dark background
    plot = chart.plots[0]
    plot.gap_width = 80

    # Series colors: green → yellow-green → orange → red
    series_colors = [
        RGBColor(0x2E, 0xCC, 0x71),  # Highly Sat - green
        RGBColor(0x8F, 0xBC, 0xE6),  # Somewhat Sat - columbia blue
        RGBColor(0xFF, 0xD7, 0x00),  # Somewhat Dissat - gold
        RGBColor(0xFF, 0x6B, 0x6B),  # Highly Dissat - red
    ]
    for i, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = series_colors[i]
        # Add data labels
        series.has_data_labels = True
        data_labels = series.data_labels
        data_labels.font.size = Pt(7)
        data_labels.font.color.rgb = WHITE
        data_labels.number_format = '0.0"%"'

    # Style axes for dark theme
    value_axis = chart.value_axis
    value_axis.has_title = False
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x33, 0x33, 0x55)
    value_axis.format.line.color.rgb = RGBColor(0x33, 0x33, 0x55)
    value_axis.tick_labels.font.size = Pt(8)
    value_axis.tick_labels.font.color.rgb = LIGHT_GRAY

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(9)
    category_axis.tick_labels.font.color.rgb = WHITE
    category_axis.format.line.color.rgb = RGBColor(0x33, 0x33, 0x55)

    # Legend styling
    legend = chart.legend
    legend.include_in_layout = False
    legend.font.size = Pt(8)
    legend.font.color.rgb = WHITE

    # If we have comparison data, add it as a text summary below the chart
    if comp_rows:
        comp_top = top + height + Inches(0.05)
        comp_text = "% Highly Satisfied — 2024 vs 2026:  "
        comp_parts = []
        for r in comp_rows:
            s24 = f"{r['s24_val']*100:.1f}%" if r.get('s24_val') else 'N/A'
            hs1 = f"{r['hs1_val']*100:.1f}%" if r.get('hs1_val') else 'N/A'
            delta = r.get('delta')
            delta_str = f"{delta:+.1f}pp" if delta is not None else ''
            comp_parts.append(f"{r['metric']}: {s24}→{hs1} ({delta_str})")
        comp_text += "  |  ".join(comp_parts)
        add_text_box(slide, left, comp_top, width, Inches(0.3),
                     comp_text, font_size=7, color=LIGHT_GRAY)

    return chart_shape


def add_qualitative_section(slide, left, top, width, sentences_df, survey_df,
                            parent_cats, rating_col=None, dept_label='',
                            theme_summaries=None):
    """Add a qualitative NPS tier breakdown section to a slide.
    Shows top 3 themes per tier with a 1-2 sentence summary of what fans said.
    """
    current_top = top
    tiers = [
        ('Detractor', '0-6', RED),
        ('Passive', '7-8', GOLD),
        ('Promoter', '9-10', GREEN),
    ]

    # Rating tier bounds for department-specific filtering
    tier_bounds = {
        'Detractor': {'min': 0, 'max': 6},
        'Passive': {'min': 7, 'max': 8},
        'Promoter': {'min': 9, 'max': 10},
    }

    for tier_name, tier_range, tier_color in tiers:
        # Get top 3 themes for this tier
        themes, total_in_tier = get_top_themes(
            sentences_df, parent_cats=parent_cats,
            nps_segment=tier_name, top_n=3
        )

        if rating_col:
            themes_by_rating, total_by_rating = get_top_themes(
                sentences_df, parent_cats=parent_cats,
                survey_df=survey_df, rating_col=rating_col,
                rating_col_tiers=tier_bounds[tier_name], top_n=3
            )
            themes = pd.concat([themes, themes_by_rating]).drop_duplicates(
                subset='AI_CATEGORY').sort_values('count', ascending=False).head(3)
            # Use the larger total for percentage denominator
            total_in_tier = max(total_in_tier, total_by_rating)

        # Tier header
        add_text_box(slide, left, current_top, width, Inches(0.18),
                     f"{tier_name}s ({tier_range})", font_size=10,
                     bold=True, color=tier_color)
        current_top += Inches(0.18)

        # Show top 3 themes with summaries
        if len(themes) > 0:
            for rank, (_, theme_row) in enumerate(themes.iterrows(), 1):
                cat_name = theme_row['AI_CATEGORY']
                cat_count = int(theme_row['count'])
                cat_pct = round(cat_count / total_in_tier * 100, 1) if total_in_tier > 0 else 0

                # Theme name + count + percentage as bold label
                add_text_box(slide, left + Inches(0.1), current_top, width - Inches(0.1),
                             Inches(0.16),
                             f"{rank}. {cat_name} ({cat_count} mentions, {cat_pct}%)",
                             font_size=8, bold=True, color=COLUMBIA_BLUE)
                current_top += Inches(0.16)

                # Get the theme summary
                summary = get_theme_summary(
                    theme_summaries, parent_cats=parent_cats,
                    nps_segment=tier_name, ai_category=cat_name,
                    rating_col=rating_col, survey_df=survey_df,
                    rating_col_tiers=tier_bounds.get(tier_name)
                )

                # Summary text (1-2 sentences) — compact height
                add_text_box(slide, left + Inches(0.2), current_top, width - Inches(0.25),
                             Inches(0.28), summary, font_size=7, color=LIGHT_GRAY)
                current_top += Inches(0.28)
        else:
            add_text_box(slide, left + Inches(0.1), current_top, width - Inches(0.1),
                         Inches(0.16), "No themes identified", font_size=8,
                         color=LIGHT_GRAY)
            current_top += Inches(0.16)

        current_top += Inches(0.03)

    return current_top


def build_scorecard_row(metric_name, hs1_series, s24_series):
    """Build a single scorecard row dict from two pandas series."""
    hs1_val = safe_avg(hs1_series)
    s24_val = safe_avg(s24_series)
    n = safe_count(hs1_series)
    return {
        'metric': metric_name,
        'hs1_val': hs1_val,
        's24_val': s24_val,
        'delta': delta_pct(hs1_val, s24_val),
        'n': n,
    }

def build_grid_row(metric_name, hs1_series, s24_series=None):
    """Build a grid distribution row dict. Shows HS1 distribution."""
    sat_levels = ['Highly satisfied', 'Somewhat satisfied',
                  'Somewhat dissatisfied', 'Highly dissatisfied']
    dist, n = grid_distribution(hs1_series, sat_levels)
    row = {
        'metric': metric_name,
        'highly_sat': dist.get('Highly satisfied', 0),
        'somewhat_sat': dist.get('Somewhat satisfied', 0),
        'somewhat_dissat': dist.get('Somewhat dissatisfied', 0),
        'highly_dissat': dist.get('Highly dissatisfied', 0),
        'n': n,
    }
    return row

def build_grid_comparison_row(metric_name, hs1_series, s24_series):
    """Build a scorecard row comparing % Highly Satisfied between HS1 and 2024."""
    sat_levels = ['Highly satisfied', 'Somewhat satisfied',
                  'Somewhat dissatisfied', 'Highly dissatisfied']
    hs1_dist, hs1_n = grid_distribution(hs1_series, sat_levels)
    s24_dist, s24_n = grid_distribution(s24_series, sat_levels)
    hs1_hs = hs1_dist.get('Highly satisfied', 0)
    s24_hs = s24_dist.get('Highly satisfied', 0)
    d = round(hs1_hs - s24_hs, 1) if hs1_n > 0 and s24_n > 0 else None
    return {
        'metric': f"{metric_name} (% Highly Sat)",
        'hs1_val': hs1_hs / 100 if hs1_hs else None,  # store as decimal for table
        's24_val': s24_hs / 100 if s24_hs else None,
        'delta': d,  # this is pp difference, not %
        'n': hs1_n,
        '_is_pct': True,  # flag for formatting
    }


# ── Slide Builders ──────────────────────────────────────────────────────────

def build_title_slide(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.0),
                 "2026 HOMESTAND #1", font_size=36, bold=True,
                 color=COLUMBIA_BLUE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(2.4), Inches(8.4), Inches(0.6),
                 "Fan Experience Report", font_size=28, bold=False,
                 color=WHITE, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(3.3), Inches(8.4), Inches(0.4),
                 "April 6 \u2013 April 12, 2026  |  6 Games", font_size=16,
                 color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(4.2), Inches(8.4), Inches(0.4),
                 "Tampa Bay Rays  |  Voice of the Customer Survey Analysis",
                 font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(8.4), Inches(0.3),
                 "2026 Homestand #1 vs 2024 Full Season Benchmark",
                 font_size=11, color=GOLD, alignment=PP_ALIGN.CENTER)


def build_exec_summary(prs, hs1, s24, sentences):
    """Slide 2: Executive Summary."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_text_box(slide, Inches(0.4), Inches(0.2), Inches(9.2), Inches(0.4),
                 "EXECUTIVE SUMMARY", font_size=22, bold=True, color=COLUMBIA_BLUE)

    # Response count and overall score
    hs1_overall = safe_avg(hs1['OVERALL_NUMRAT'])
    s24_overall = safe_avg(s24['OVERALL_NUMRAT'])
    overall_delta = delta_pct(hs1_overall, s24_overall)
    total_responses = len(hs1)

    # NPS distribution
    hs1_nps = hs1['OVERALL_NUMRAT'].apply(nps_tier)
    nps_counts = hs1_nps.value_counts()
    pct_promoter = round(100 * nps_counts.get('Promoter', 0) / total_responses, 1)
    pct_passive = round(100 * nps_counts.get('Passive', 0) / total_responses, 1)
    pct_detractor = round(100 * nps_counts.get('Detractor', 0) / total_responses, 1)
    nps_score = round(pct_promoter - pct_detractor, 1)

    summary_text = (
        f"Total Responses: {total_responses:,}   |   "
        f"Overall Avg: {hs1_overall:.2f} (2024: {s24_overall:.2f}, {overall_delta:+.1f}%)   |   "
        f"NPS: {nps_score:+.1f}"
    )
    add_text_box(slide, Inches(0.4), Inches(0.65), Inches(9.2), Inches(0.3),
                 summary_text, font_size=11, bold=True, color=WHITE)

    nps_text = (
        f"Promoters (9-10): {pct_promoter:.1f}%   |   "
        f"Passives (7-8): {pct_passive:.1f}%   |   "
        f"Detractors (0-6): {pct_detractor:.1f}%"
    )
    add_text_box(slide, Inches(0.4), Inches(0.95), Inches(9.2), Inches(0.25),
                 nps_text, font_size=10, color=LIGHT_GRAY)

    # Department scorecard
    dept_rows = []
    dept_metrics = [
        ('Overall Experience', 'OVERALL_NUMRAT'),
        ('Concessions', 'CONCESS_NUMRAT'),
        ('Parking', 'PARKING_NUMRAT'),
        ('Entertainment', 'ENTERTAIN_NUMRAT'),
        ('Merchandise', 'MERCH_NUMRAT'),
    ]
    for label, col in dept_metrics:
        dept_rows.append(build_scorecard_row(label, hs1[col], s24[col]))

    add_scorecard_table(slide, Inches(0.4), Inches(1.35), Inches(6.8), dept_rows)

    # Top positive & negative themes
    pos_themes, _ = get_top_themes(sentences, nps_segment='Promoter', top_n=5)
    neg_themes, _ = get_top_themes(sentences, nps_segment='Detractor', top_n=5)

    add_text_box(slide, Inches(0.4), Inches(3.6), Inches(4.5), Inches(0.25),
                 "TOP PROMOTER THEMES (9-10)", font_size=10, bold=True, color=GREEN)
    if len(pos_themes) > 0:
        pos_text = "\n".join(f"  {row['AI_CATEGORY']}  ({row['count']})" for _, row in pos_themes.iterrows())
        add_text_box(slide, Inches(0.4), Inches(3.85), Inches(4.5), Inches(1.5),
                     pos_text, font_size=8, color=LIGHT_GRAY)

    add_text_box(slide, Inches(5.2), Inches(3.6), Inches(4.5), Inches(0.25),
                 "TOP DETRACTOR THEMES (0-6)", font_size=10, bold=True, color=RED)
    if len(neg_themes) > 0:
        neg_text = "\n".join(f"  {row['AI_CATEGORY']}  ({row['count']})" for _, row in neg_themes.iterrows())
        add_text_box(slide, Inches(5.2), Inches(3.85), Inches(4.5), Inches(1.5),
                     neg_text, font_size=8, color=LIGHT_GRAY)


def build_department_slides(prs, hs1, s24, sentences, dept_name, dept_config, theme_summaries):
    """Build 1-2 slides for a department."""

    # ── Slide 1: Quantitative (skip if configured) ──
    if not dept_config.get('skip_quantitative'):
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide1)

        add_text_box(slide1, Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.4),
                     dept_name.upper(), font_size=22, bold=True, color=COLUMBIA_BLUE)

        current_top = Inches(0.55)

        # Scorecard rows
        scorecard_rows = dept_config.get('scorecard', [])
        if scorecard_rows:
            rows_data = []
            for label, col in scorecard_rows:
                rows_data.append(build_scorecard_row(label, hs1[col], s24[col]))
            add_scorecard_table(slide1, Inches(0.4), current_top, Inches(6.8), rows_data)
            current_top += Inches(0.3 * (len(rows_data) + 1)) + Inches(0.15)

        # Grid rows
        grid_metrics = dept_config.get('grid', [])
        use_bar_charts = dept_config.get('use_bar_charts', False)
        if grid_metrics:
            grid_rows = []
            for label, col in grid_metrics:
                grid_rows.append(build_grid_row(label, hs1[col]))

            # Build comparison rows (needed for both chart and table paths)
            comparison_rows = []
            for label, col in grid_metrics:
                comparison_rows.append(build_grid_comparison_row(label, hs1[col], s24[col]))

            if use_bar_charts:
                # Render as horizontal bar chart
                add_text_box(slide1, Inches(0.4), current_top, Inches(9.0), Inches(0.25),
                             "Satisfaction Distribution (2026 Homestand #1)",
                             font_size=10, bold=True, color=WHITE)
                current_top += Inches(0.25)
                chart_height = Inches(0.7 * len(grid_rows) + 0.5)
                add_grid_bar_chart(slide1, Inches(0.4), current_top, Inches(8.8),
                                   chart_height, grid_rows, comparison_rows)
                current_top += chart_height + Inches(0.4)  # extra space for comparison text
            else:
                # Render as table (original behavior)
                add_text_box(slide1, Inches(0.4), current_top, Inches(9.0), Inches(0.25),
                             "Satisfaction Distribution (2026 Homestand #1)",
                             font_size=10, bold=True, color=WHITE)
                current_top += Inches(0.25)
                add_grid_table(slide1, Inches(0.4), current_top, Inches(7.3), grid_rows)
                current_top += Inches(0.28 * (len(grid_rows) + 1)) + Inches(0.15)

                # Grid comparison table (% Highly Satisfied: 2024 vs 2026)
                comp_formatted = []
                for row in comparison_rows:
                    comp_formatted.append({
                        'metric': row['metric'],
                        'hs1_val': row['hs1_val'] * 100 if row['hs1_val'] else None,
                        's24_val': row['s24_val'] * 100 if row['s24_val'] else None,
                        'delta': row['delta'],
                        'n': row['n'],
                    })
                add_text_box(slide1, Inches(0.4), current_top, Inches(9.0), Inches(0.25),
                             "% Highly Satisfied: 2024 Season vs 2026 HS1",
                             font_size=10, bold=True, color=WHITE)
                current_top += Inches(0.25)
                add_scorecard_table(slide1, Inches(0.4), current_top, Inches(6.8), comp_formatted)
                current_top += Inches(0.3 * (len(comp_formatted) + 1)) + Inches(0.1)

        # Extra content (categorical distributions, multi-select, grids, etc.)
        extras = dept_config.get('extras', [])
        for extra in extras:
            if extra['type'] == 'categorical':
                dist, n = categorical_distribution(hs1[extra['col']], extra['categories'])
                if n == 0:
                    # Skip sections with no valid data (e.g. all NULL columns)
                    continue
                add_text_box(slide1, Inches(0.4), current_top, Inches(9.0), Inches(0.25),
                             extra['title'], font_size=10, bold=True, color=WHITE)
                current_top += Inches(0.25)
                dist_text = "  |  ".join(f"{k}: {v:.1f}%" for k, v in dist.items() if v > 0)
                add_text_box(slide1, Inches(0.5), current_top, Inches(9.0), Inches(0.2),
                             f"{dist_text}  (n={n})", font_size=8, color=LIGHT_GRAY)
                current_top += Inches(0.25)

            elif extra['type'] == 'multiselect':
                results, n = multiselect_distribution(hs1, extra['cols'], extra.get('labels'))
                if n == 0:
                    continue
                add_text_box(slide1, Inches(0.4), current_top, Inches(9.0), Inches(0.25),
                             extra['title'], font_size=10, bold=True, color=WHITE)
                current_top += Inches(0.25)
                dist_text = "  |  ".join(f"{lbl}: {pct:.1f}%" for lbl, pct, _ in results if pct > 0)
                add_text_box(slide1, Inches(0.5), current_top, Inches(9.0), Inches(0.2),
                             f"{dist_text}  (n={n})", font_size=8, color=LIGHT_GRAY)
                current_top += Inches(0.25)

            elif extra['type'] == 'text_multiselect':
                results, n = text_multiselect_distribution(
                    hs1[extra['col']], extra['categories'], extra.get('labels'))
                if n == 0:
                    continue
                add_text_box(slide1, Inches(0.4), current_top, Inches(9.0), Inches(0.25),
                             extra['title'], font_size=10, bold=True, color=WHITE)
                current_top += Inches(0.25)
                dist_text = "  |  ".join(f"{lbl}: {pct:.1f}%" for lbl, pct, _ in results if pct > 0)
                add_text_box(slide1, Inches(0.5), current_top, Inches(9.0), Inches(0.2),
                             f"{dist_text}  (n={n})", font_size=8, color=LIGHT_GRAY)
                current_top += Inches(0.25)

            elif extra['type'] == 'food_quality_grid':
                add_text_box(slide1, Inches(0.4), current_top, Inches(9.0), Inches(0.25),
                             extra['title'], font_size=10, bold=True, color=WHITE)
                current_top += Inches(0.25)
                food_grid_rows = []
                for label, col in extra['items']:
                    food_grid_rows.append(build_grid_row(label, hs1[col]))
                # Filter out items with 0 responses
                food_grid_rows = [r for r in food_grid_rows if r['n'] > 0]
                if food_grid_rows:
                    add_grid_table(slide1, Inches(0.4), current_top, Inches(7.3), food_grid_rows)
                    current_top += Inches(0.28 * (len(food_grid_rows) + 1)) + Inches(0.1)

            elif extra['type'] == 'age_buckets':
                add_text_box(slide1, Inches(0.4), current_top, Inches(9.0), Inches(0.25),
                             extra['title'], font_size=10, bold=True, color=WHITE)
                current_top += Inches(0.25)
                dist, n = age_bucket_distribution(hs1[extra['col']])
                dist_text = "  |  ".join(f"{k}: {v:.1f}%" for k, v in dist.items() if v > 0)
                add_text_box(slide1, Inches(0.5), current_top, Inches(9.0), Inches(0.2),
                             f"{dist_text}  (n={n})", font_size=8, color=LIGHT_GRAY)
                current_top += Inches(0.25)

    # ── Slide 2: Qualitative (skip for departments with no qualitative data) ──
    if dept_config.get('skip_qualitative'):
        return

    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide2)

    add_text_box(slide2, Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.4),
                 f"{dept_name.upper()} \u2014 FAN FEEDBACK BY NPS TIER",
                 font_size=18, bold=True, color=COLUMBIA_BLUE)

    parent_cats = dept_config.get('parent_categories', [])
    rating_col = dept_config.get('primary_rating_col')

    add_qualitative_section(
        slide2, Inches(0.4), Inches(0.6), Inches(9.2),
        sentences, hs1, parent_cats, rating_col, dept_name, theme_summaries
    )


# ── Department Configurations ───────────────────────────────────────────────

def get_department_configs():
    """Define what metrics each department section should show."""

    BRAND_CATEGORIES = ['Strongly agree', 'Somewhat agree',
                        'Somewhat disagree', 'Strongly disagree']

    return {
        'Fan Experience': {
            'scorecard': [
                ('Overall Experience Rating', 'OVERALL_NUMRAT'),
            ],
            'grid': [],
            'parent_categories': ['Fan Experience'],
            'primary_rating_col': 'OVERALL_NUMRAT',
            'extras': [
                {
                    'type': 'categorical',
                    'title': 'Seat View Value Perception',
                    'col': 'SEATVIEW_VALUE_SEAT',
                    'categories': ['Good value for the money I spent on the ticket',
                                   'Fair value for the money I spent on the ticket',
                                   'Poor value for the money I spent on the ticket'],
                },
                {
                    'type': 'categorical',
                    'title': 'Brand: Exciting',
                    'col': 'BRANDHEALTH_GRID_EXCITING_DESC',
                    'categories': BRAND_CATEGORIES,
                },
                {
                    'type': 'categorical',
                    'title': 'Brand: Family Friendly',
                    'col': 'BRANDHEALTH_GRID_FAMFRIENDLY_DESC',
                    'categories': BRAND_CATEGORIES,
                },
                {
                    'type': 'categorical',
                    'title': 'Brand: Right Direction',
                    'col': 'BRANDHEALTH_GRID_RIGHTDIRECTION_DESC',
                    'categories': BRAND_CATEGORIES,
                },
                {
                    'type': 'categorical',
                    'title': 'Brand: Welcoming',
                    'col': 'BRANDHEALTH_GRID_WELCOME_DESC',
                    'categories': BRAND_CATEGORIES,
                },
                {
                    'type': 'categorical',
                    'title': 'Brand: Safe',
                    'col': 'BRANDHEALTH_GRID_SAFE_DESC',
                    'categories': BRAND_CATEGORIES,
                },
                {
                    'type': 'categorical',
                    'title': 'Brand: Champion of Community',
                    'col': 'BRANDHEALTH_GRID_CHAMPION_DESC',
                    'categories': BRAND_CATEGORIES,
                },
                {
                    'type': 'categorical',
                    'title': 'Brand: Diversity & Inclusion',
                    'col': 'BRANDHEALTH_GRID_DIVERSITY_DESC',
                    'categories': BRAND_CATEGORIES,
                },
            ],
        },

        'Ticketing': {
            'scorecard': [],
            'grid': [],
            'parent_categories': ['Ticketing'],
            'primary_rating_col': 'OVERALL_NUMRAT',
            'extras': [
                {
                    'type': 'categorical',
                    'title': 'Purchase Intent (This Season)',
                    'col': 'PURCHASE_INTENT_DESC',
                    'categories': ['Yes, I do', 'No, I do not', 'I am not sure yet'],
                },
                {
                    'type': 'categorical',
                    'title': 'Purchase Intent (Next Season)',
                    'col': 'PURCHASE_INTENT_NEXT_DESC',
                    'categories': ['Highly likely', 'Somewhat likely',
                                   'Somewhat unlikely', 'Highly unlikely'],
                },
                {
                    'type': 'categorical',
                    'title': 'Already Purchased Future Tickets',
                    'col': 'PREVIOUS_PURCHASE_DESC',
                    'categories': ['Yes', 'No'],
                },
                {
                    'type': 'categorical',
                    'title': 'Top Purchase Incentive (#1 Rank)',
                    'col': 'INCENTIVES_RANK_1_DESC',
                    'categories': [
                        'Receiving a ticket discount / offer',
                        'Other',
                        'More affordable food / beverage options',
                        'Receiving a concessions credit',
                        'Receiving a parking credit',
                        'Giveaways (e.g., bobbleheads)',
                        'Shorter lines throughout the ballpark',
                        'Theme nights (e.g. 80s Night, fireworks)',
                    ],
                },
            ],
        },

        'Concessions': {
            'scorecard': [
                ('Concessions Rating', 'CONCESS_NUMRAT'),
            ],
            'grid': [
                ('Value for Money', 'CONCESS_GRID_VALUE_DESC'),
                ('Food/Bev Selection', 'CONCESS_GRID_SELECTION_DESC'),
                ('Cleanliness', 'CONCESS_GRID_CLEAN_DESC'),
            ],
            'parent_categories': ['Concessions'],
            'primary_rating_col': 'CONCESS_NUMRAT',
            'extras': [
                {
                    'type': 'categorical',
                    'title': 'Wait Time',
                    'col': 'CONCESS_WAIT_DESC',
                    'categories': ['Less than 5 minutes', '5-15 minutes',
                                   '16-30 minutes', 'More than 30 minutes'],
                },
                {
                    'type': 'categorical',
                    'title': 'Spend Per Person',
                    'col': 'CONCESS_SPEND_DESC',
                    'categories': ['Less than $20', '$20-$30', '$31-$40',
                                   '$41-$50', 'More than $50'],
                },
                {
                    'type': 'categorical',
                    'title': 'Wait Time vs Expectations',
                    'col': 'CONCESS_WAIT_EXPECT_DESC',
                    'categories': ['Much less than what I expected',
                                   'Slightly less what I expected',
                                   'About what I expected',
                                   'Slightly more than what I expected',
                                   'Much more than what I expected'],
                },
                {
                    'type': 'categorical',
                    'title': 'Why Not Returning for Concessions',
                    'col': 'CONCESS_INTENT_NO_DESC',
                    'categories': [
                        'Shorter wait times',
                        'Additional budget-friendly options',
                        'Better value of the food / beverages for the money',
                        'Better options near my seat',
                        'More options that meet my dietary needs (e.g., gluten-free, vegan)',
                        'More variety',
                        'Other (please specify below)',
                    ],
                },
                {
                    'type': 'food_quality_grid',
                    'title': 'Food & Beverage Quality Ratings',
                    'items': [
                        ('Hot Dogs', 'CONCESS_QUALITY_HOTDOG_DESC'),
                        ('Alcohol', 'CONCESS_QUALITY_ALCOHOL_DESC'),
                        ('Non-Alcohol Beverages', 'CONCESS_QUALITY_NONALCOHOL_DESC'),
                        ('Pretzels', 'CONCESS_QUALITY_PRETZELS_DESC'),
                        ('Popcorn', 'CONCESS_QUALITY_POPCORN_DESC'),
                        ('Nuts', 'CONCESS_QUALITY_NUTS_DESC'),
                        ('Fries', 'CONCESS_QUALITY_FRIES_DESC'),
                        ('Chicken', 'CONCESS_QUALITY_CHICKEN_DESC'),
                        ('Ice Cream', 'CONCESS_QUALITY_ICECREAM_DESC'),
                        ('Other Entrees', 'CONCESS_QUALITY_OTHER_ENTREE_DESC'),
                        ('Burgers', 'CONCESS_QUALITY_BURGERS_DESC'),
                        ('Pizza', 'CONCESS_QUALITY_PIZZA_DESC'),
                        ('Nachos', 'CONCESS_QUALITY_NACHOS_DESC'),
                        ('Sandwiches', 'CONCESS_QUALITY_SANDWICH_DESC'),
                        ('Other Desserts', 'CONCESS_QUALITY_OTHER_DESSERT_DESC'),
                        ('Sausage', 'CONCESS_QUALITY_SAUSAGE_DESC'),
                    ],
                },
            ],
        },

        'Retail': {
            'scorecard': [
                ('Merchandise Rating', 'MERCH_NUMRAT'),
            ],
            'grid': [
                ('Quality', 'MERCH_GRID_MERCHQUALITY_DESC'),
                ('Selection', 'MERCH_GRID_SELECTION_DESC'),
                ('Pricing', 'MERCH_GRID_PRICE_DESC'),
                ('Wait Time', 'MERCH_GRID_WAIT_DESC'),
            ],
            'use_bar_charts': True,
            'parent_categories': ['Retail'],
            'primary_rating_col': 'MERCH_NUMRAT',
            'extras': [
                {
                    'type': 'categorical',
                    'title': 'Purchased Merchandise?',
                    'col': 'MERCH_SCREENER_DESC',
                    'categories': ['Yes, I did', 'No, but I was interested in purchasing merchandise',
                                   'No, and I was not interested in purchasing merchandise',
                                   "I don't remember"],
                },
                {
                    'type': 'categorical',
                    'title': 'Reason for Not Purchasing Merchandise',
                    'col': 'MERCH_NO_DESC',
                    'categories': [
                        'The lines were too long',
                        'Merchandise was too expensive',
                        'The merchandise store / stands were too far away from my seat',
                        'Merchandise selection was limited / low inventory',
                        'There was no merchandise for my favorite team',
                        'Customer service was poor',
                        'Other (please specify below)',
                    ],
                },
            ],
        },

        'Security': {
            'scorecard': [],
            'grid': [],
            'parent_categories': ['Stadium Operations'],
            'primary_rating_col': None,
            'extras': [
                {
                    'type': 'categorical',
                    'title': 'Gate Entry Time vs Expectations',
                    'col': 'GE_TIME_EXPECT_DESC',
                    'categories': ['More than what I expected', 'About what I expected',
                                   'Less than what I expected'],
                },
            ],
        },

        'Parking': {
            'scorecard': [
                ('Parking Rating', 'PARKING_NUMRAT'),
            ],
            'grid': [],
            'parent_categories': ['Parking'],
            'primary_rating_col': 'PARKING_NUMRAT',
            'extras': [
                {
                    'type': 'categorical',
                    'title': 'Parked in Club Lot?',
                    'col': 'CLUB_PARKING_LOT_DESC',
                    'categories': ['Yes', 'No', "I'm not sure"],
                },
                {
                    'type': 'categorical',
                    'title': 'Travel Method to Stadium',
                    'col': 'TRAVELTO_METHOD_DESC',
                    'categories': ['Car/Personal Vehicle',
                                   'Rideshare service (e.g. Uber, Lyft)',
                                   'Walk/Bike/Scooter', 'Bus'],
                },
                {
                    'type': 'categorical',
                    'title': 'Travel Time vs Expectations',
                    'col': 'TRAVELTO_TIME_EXPECT_DESC',
                    'categories': ['Less than what I expected',
                                   'About what I expected',
                                   'More than what I expected'],
                },
            ],
        },

        'Entertainment': {
            'scorecard': [
                ('Entertainment Rating', 'ENTERTAIN_NUMRAT'),
            ],
            'grid': [
                ('In-Game Music', 'ENTERTAIN_GRID_MUSIC_DESC'),
                ('Scoreboard/Video', 'ENTERTAIN_GRID_SCOREBOARD_DESC'),
                ('In-Game Activities', 'ENTERTAIN_GRID_GAMES_DESC'),
                ('Kids Activities', 'ENTERTAIN_GRID_KIDS_ACTIVITIES_DESC'),
                ('Pregame Content', 'ENTERTAIN_GRID_PREGAME_CONTENT_DESC'),
            ],
            'use_bar_charts': True,
            'parent_categories': ['Game Entertainment'],
            'primary_rating_col': 'ENTERTAIN_NUMRAT',
            'extras': [],
        },

        'Staff / Operations': {
            'scorecard': [],
            'grid': [],
            'parent_categories': ['Fan Experience', 'Stadium Operations'],
            'primary_rating_col': None,
            'extras': [
                {
                    'type': 'text_multiselect',
                    'title': 'Staff Types Interacted With',
                    'col': 'STAFF_TYPE',
                    'categories': [
                        'concessions', 'security', 'usher', 'parking',
                        'fan_services', 'merch', 'ticket_scanner', 'none',
                    ],
                    'labels': {
                        'concessions': 'Concessions',
                        'security': 'Security',
                        'usher': 'Ushers',
                        'parking': 'Parking',
                        'fan_services': 'Fan Services',
                        'merch': 'Merchandise',
                        'ticket_scanner': 'Ticket Scanner',
                        'none': 'None / Did Not Interact',
                    },
                },
            ],
        },

        'Promotions': {
            'scorecard': [],
            'grid': [],
            'skip_quantitative': True,
            'parent_categories': ['Game Entertainment', 'Marketing'],
            'primary_rating_col': None,
            'extras': [
                {
                    'type': 'categorical',
                    'title': 'Giveaway Satisfaction',
                    'col': 'GIVEAWAY_SAT_DESC',
                    'categories': ['Highly satisfied', 'Somewhat satisfied',
                                   'Somewhat dissatisfied', 'Highly dissatisfied'],
                },
                {
                    'type': 'categorical',
                    'title': 'Received Giveaway?',
                    'col': 'GIVEAWAY_RECEIVE_DESC',
                    'categories': ['Yes', 'No'],
                },
                {
                    'type': 'categorical',
                    'title': 'Awareness: Giveaway (Before Arriving)',
                    'col': 'PA_PROMO_GRID_GIVEAWAY_DESC',
                    'categories': ['I cared about this',
                                   'I knew about this, but I did not care about it',
                                   'I did not know about this'],
                },
                {
                    'type': 'categorical',
                    'title': 'Awareness: Theme Night (Before Arriving)',
                    'col': 'PA_PROMO_GRID_THEME_DESC',
                    'categories': ['I cared about this',
                                   'I knew about this, but I did not care about it',
                                   'I did not know about this'],
                },
                {
                    'type': 'categorical',
                    'title': 'Giveaway Impact on Attendance Decision',
                    'col': 'PA_GIVEAWAY_INTENT_DESC',
                    'categories': ['I still would have gone to this game',
                                   'I would have chosen a game on another day',
                                   'I would not have purchased / attended at this time'],
                },
            ],
        },

        'Demographics & Attendance': {
            'scorecard': [],
            'grid': [],
            'parent_categories': [],
            'primary_rating_col': None,
            'skip_qualitative': True,
            'extras': [
                {
                    'type': 'categorical',
                    'title': 'Gender',
                    'col': 'GENDER_ID_DESC',
                    'categories': ['Man', 'Woman', 'Prefer to self-describe',
                                   'I prefer not to say'],
                },
                {
                    'type': 'age_buckets',
                    'title': 'Age Distribution',
                    'col': 'AGE',
                },
                {
                    'type': 'categorical',
                    'title': 'Parent Status',
                    'col': 'PARENT_ID_DESC',
                    'categories': ['Yes, I am', 'No, I am not', 'I prefer not to say'],
                },
                {
                    'type': 'categorical',
                    'title': 'Children in Household',
                    'col': 'CIH_ID_DESC',
                    'categories': ['None', '1', '2', '3', '4', 'I prefer not to say'],
                },
                {
                    'type': 'categorical',
                    'title': 'Team Avidity',
                    'col': 'TEAM_AVIDITY_DESC',
                    'categories': ['5 (passionate fan)', '4', '3', '2',
                                   '1 (occassional fan)'],
                },
                {
                    'type': 'categorical',
                    'title': 'Household Income',
                    'col': 'HHI_ID_DESC',
                    'categories': [
                        'Under $50,000', 'Between $50,000 and $75,000',
                        'Between $75,001 and $100,000', 'Between $100,001 and $125,000',
                        'Between $125,001 and $150,000', 'Between $150,001 and $200,000',
                        'Between $200,001 and $250,000', 'Between $250,001 and $300,000',
                        'Between $300,001 and $400,000', 'Between $400,001 and $500,000',
                        'Between $500,001 and $1,000,000', '$1,000,001 or more',
                        'I prefer not to say',
                    ],
                },
                {
                    'type': 'multiselect',
                    'title': 'Who Did You Attend With?',
                    'cols': [
                        'ATTEND_WITH_CATEGORY_FRIENDS', 'ATTEND_WITH_CATEGORY_SPOUSE',
                        'ATTEND_WITH_CATEGORY_ADULT_KIDS', 'ATTEND_WITH_CATEGORY_ALONE',
                        'ATTEND_WITH_CATEGORY_BUSINESS', 'ATTEND_WITH_CATEGORY_OTHERFAM',
                        'ATTEND_WITH_CATEGORY_OTHER',
                    ],
                    'labels': {
                        'ATTEND_WITH_CATEGORY_FRIENDS': 'Friends',
                        'ATTEND_WITH_CATEGORY_SPOUSE': 'Spouse/Partner',
                        'ATTEND_WITH_CATEGORY_ADULT_KIDS': 'Adult Children',
                        'ATTEND_WITH_CATEGORY_ALONE': 'Alone',
                        'ATTEND_WITH_CATEGORY_BUSINESS': 'Business/Colleagues',
                        'ATTEND_WITH_CATEGORY_OTHERFAM': 'Other Family',
                        'ATTEND_WITH_CATEGORY_OTHER': 'Other',
                    },
                },
                {
                    'type': 'categorical',
                    'title': 'Games Planned to Attend This Season',
                    'col': 'ATTEND_NUM_PLAN_DESC',
                    'categories': ['6', '7', 'Between 6 and 10 games',
                                   'Between 11 and 15 games',
                                   'Between 16 and 20 games', 'Over 20 games'],
                },
            ],
        },
    }


# ── Category Summary Slides ────────────────────────────────────────────────

def build_category_summary_slides(prs, sentences_df):
    """Build slides showing every AI_CATEGORY grouped by PARENT_CATEGORY
    with Positive / Negative / Neutral / Total sentence counts."""

    df = sentences_df.copy()

    # Filter out malformed / hallucinated AI_CATEGORY values
    df = df[df['AI_CATEGORY'].notna()]
    df = df[df['AI_CATEGORY'].str.len() <= 40]

    # Order parent categories by total sentence count (descending)
    parent_order = (df.groupby('PARENT_CATEGORY')
                    .size()
                    .sort_values(ascending=False)
                    .index.tolist())

    # Build summary data: for each parent, list AI categories with sentiment counts
    all_groups = []
    for parent in parent_order:
        pdf = df[df['PARENT_CATEGORY'] == parent]
        parent_total = len(pdf)

        agg = (pdf.groupby(['AI_CATEGORY', 'SENTIMENT_CATEGORY'])
               .size()
               .unstack(fill_value=0)
               .reset_index())

        # Ensure all sentiment columns exist
        for col in ['Positive', 'Negative', 'Neutral']:
            if col not in agg.columns:
                agg[col] = 0

        # Include Mixed in total but don't display it as a column
        mixed = agg['Mixed'] if 'Mixed' in agg.columns else 0
        agg['Total'] = agg['Positive'] + agg['Negative'] + agg['Neutral'] + mixed
        agg = agg.sort_values('Negative', ascending=False)

        all_groups.append({
            'parent': parent,
            'parent_total': parent_total,
            'categories': agg[['AI_CATEGORY', 'Positive', 'Negative',
                               'Neutral', 'Total']].values.tolist(),
        })

    # Render slides — fit multiple parent groups per slide, overflow to new slides
    MAX_ROWS_PER_SLIDE = 16  # data rows per slide (excluding headers)
    slide = None
    rows_on_slide = 0
    current_top = Inches(0)

    def new_slide():
        nonlocal slide, rows_on_slide, current_top
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide)
        add_text_box(slide, Inches(0.4), Inches(0.15), Inches(9.2), Inches(0.35),
                     "SENTENCE CATEGORY SUMMARY", font_size=18, bold=True,
                     color=COLUMBIA_BLUE)
        add_text_box(slide, Inches(0.4), Inches(0.45), Inches(9.2), Inches(0.22),
                     "AI-identified themes from fan comments  |  Homestand #1 2026",
                     font_size=9, color=LIGHT_GRAY)
        current_top = Inches(0.75)
        rows_on_slide = 0

    col_headers = ['Category', 'Positive', 'Negative', 'Neutral', 'Total']
    col_widths = [Inches(3.8), Inches(1.1), Inches(1.1), Inches(1.1), Inches(1.0)]
    num_cols = len(col_headers)
    table_left = Inches(0.4)
    table_width = Inches(8.4)

    for group in all_groups:
        cat_rows = group['categories']
        # Rows needed: 1 parent header + 1 table header + data rows + 0.1" gap
        rows_needed = len(cat_rows) + 1  # +1 for column header row

        # If this group won't fit, start a new slide
        if slide is None or rows_on_slide + rows_needed + 1 > MAX_ROWS_PER_SLIDE:
            new_slide()

        # Parent category label
        parent_label = f"{group['parent']}  ({group['parent_total']} sentences)"
        add_text_box(slide, Inches(0.4), current_top, Inches(9.0), Inches(0.25),
                     parent_label, font_size=11, bold=True, color=GOLD)
        current_top += Inches(0.27)

        # Build table for this parent group
        num_data_rows = len(cat_rows)
        num_table_rows = num_data_rows + 1  # +1 header
        row_height = 0.22
        table_height = Inches(row_height * num_table_rows)

        table_shape = slide.shapes.add_table(
            num_table_rows, num_cols, table_left, current_top,
            table_width, table_height)
        table = table_shape.table

        for i, w in enumerate(col_widths):
            table.columns[i].width = w

        # Header row
        for i, h in enumerate(col_headers):
            cell = table.cell(0, i)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = 'Calibri'
                p.alignment = PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT
            cell.fill.solid()
            cell.fill.fore_color.rgb = MEDIUM_BLUE

        # Data rows
        for r_idx, row_data in enumerate(cat_rows):
            row_num = r_idx + 1
            ai_cat = str(row_data[0])
            pos = int(row_data[1])
            neg = int(row_data[2])
            neu = int(row_data[3])
            tot = int(row_data[4])

            values = [ai_cat, str(pos), str(neg), str(neu), str(tot)]
            for c_idx, val in enumerate(values):
                cell = table.cell(row_num, c_idx)
                cell.text = val
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.font.name = 'Calibri'
                    p.alignment = PP_ALIGN.CENTER if c_idx > 0 else PP_ALIGN.LEFT
                    # Color-code sentiment count columns
                    if c_idx == 1 and pos > 0:
                        p.font.color.rgb = GREEN
                    elif c_idx == 2 and neg > 0:
                        p.font.color.rgb = RED
                    elif c_idx == 3:
                        p.font.color.rgb = LIGHT_GRAY
                    elif c_idx == 4:
                        p.font.color.rgb = WHITE
                    else:
                        p.font.color.rgb = WHITE
                # Alternating row bg
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE_BG if r_idx % 2 == 0 else NAVY

        current_top += table_height + Inches(0.12)
        rows_on_slide += rows_needed + 2  # account for parent label + gap


# ── Main ────────────────────────────────────────────────────────────────────

def main():
    print("Loading data from Snowflake...")
    hs1, s24, sentences, theme_summaries = load_data()

    # Normalize column names to uppercase
    hs1.columns = [c.upper() for c in hs1.columns]
    s24.columns = [c.upper() for c in s24.columns]
    sentences.columns = [c.upper() for c in sentences.columns]

    # Clean numeric codes from _DESC columns (fixes mixed dtype matching)
    clean_desc_columns(hs1, s24)

    print("Building PowerPoint...")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 widescreen

    # Title slide
    build_title_slide(prs)
    print("  Title slide done.")

    # Executive summary
    build_exec_summary(prs, hs1, s24, sentences)
    print("  Executive summary done.")

    # Sentence category summary (after exec summary, before departments)
    build_category_summary_slides(prs, sentences)
    print("  Category summary slides done.")

    # Department slides
    dept_configs = get_department_configs()
    for dept_name, config in dept_configs.items():
        print(f"  Building {dept_name}...")
        build_department_slides(prs, hs1, s24, sentences, dept_name, config, theme_summaries)

    # Save
    prs.save(OUTPUT_FILE)
    print(f"\nReport saved to: {OUTPUT_FILE}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
