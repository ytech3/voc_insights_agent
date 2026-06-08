"""
build_voc_hs2_vs_hs3.py
Tampa Bay Rays - VOC HS2 vs HS3 Homestand Comparison with MLB Benchmarks

Produces a 7-slide PPTX replicating the VOC_HS1_vs_HS2_v8.pptx layout,
adding MLB AVG and Rank columns from the March & April 2026 VOC Program report.

9-column layout:
  Metric | MLB AVG | Rank | '24 Score | '26 AVG | HS2 Score | HS3 Score | HS3 vs HS2 % Diff | '26 vs '24 % Diff

Usage:
    python build_voc_hs2_vs_hs3.py
"""

import json
import os
import sys
import openpyxl
import snowflake.connector
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
CONFIG_PATH = os.path.join(SCRIPT_DIR, "homestand_config.json")
MLB_DATA_PATH = os.path.join(SCRIPT_DIR, "Monthly", "Raw Data - March & April 2026.xlsx")
CONCESSIONS_QUALITY_PATH = os.path.join(SCRIPT_DIR, "Monthly", "Concessions Quality Results by Club - March & April 2026.xlsx")
OUTPUT_PATH = os.path.join(SCRIPT_DIR, "VOC_HS2_vs_HS3_v2.pptx")

# ── Snowflake connection ──
SF_ACCOUNT = "hta92307.east-us-2.azure"
SF_WAREHOUSE = "TBRDP_DW_CORTEX_XS_WH"
SF_DATABASE = "TBRDP_DW_DEV"
SF_SCHEMA = "IM_RPT"
SF_ROLE = "TBRDP_DW_PROD_CORTEX_USER"
SF_AUTHENTICATOR = "externalbrowser"

# ── Views ──
VOC_VIEW = "TBRDP_DW_DEV.IM_RPT.V_SBL_QUALTRICS_VOC_POST_ATTENDANCE_FULL_CORTEX_AI"
PARKING_VIEW = "TBRDP_DW_DEV.IM_RPT.V_SBL_QUALTRICS_VOC_POST_ATTENDANCE_FULL"

# ── Theme colors (Rays brand palette) ──
NAVY = RGBColor(0x09, 0x2C, 0x5C)
SKY = RGBColor(0x8F, 0xBC, 0xE6)
YELLOW = RGBColor(0xF5, 0xD1, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x4B, 0x4B, 0x4B)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xFF, 0x6B, 0x6B)
MEDIUM_BLUE = RGBColor(0x0D, 0x3B, 0x7A)
DARK_BLUE_BG = RGBColor(0x07, 0x20, 0x44)
MUTED_GRAY = RGBColor(0x99, 0x99, 0x99)
LIGHT_BG = RGBColor(0xF2, 0xF4, 0xF8)
GOLD = RGBColor(0xD4, 0xA5, 0x17)

# ── Logo ──
LOGO_PATH = os.path.join(SCRIPT_DIR, "MLB Logos", "TB_White.png")


# ═══════════════════════════════════════════════════════════════════════
# CONFIG & CONNECTION
# ═══════════════════════════════════════════════════════════════════════

def load_config():
    with open(CONFIG_PATH, "r") as f:
        return json.load(f)


def get_homestand(config, season, hs_num):
    for hs in config["homestands"]:
        if hs["season"] == season and hs["homestand_num"] == hs_num:
            return hs
    raise ValueError(f"Homestand {hs_num} for season {season} not found")


def get_connection():
    return snowflake.connector.connect(
        user="YTAKETANI@RAYSBASEBALL.COM",
        account=SF_ACCOUNT,
        authenticator=SF_AUTHENTICATOR,
        warehouse=SF_WAREHOUSE,
        database=SF_DATABASE,
        schema=SF_SCHEMA,
        role=SF_ROLE,
    )


def run_query(conn, sql):
    cur = conn.cursor()
    try:
        cur.execute(sql)
        return cur.fetchall()
    finally:
        cur.close()


# ═══════════════════════════════════════════════════════════════════════
# MLB BENCHMARK DATA LOADER
# ═══════════════════════════════════════════════════════════════════════

def load_mlb_benchmarks():
    """Load MLB benchmark data from the Raw Data Excel file.
    Returns a dict mapping metric names to (mlb_avg, tb_rank, total_clubs)."""
    import statistics

    wb = openpyxl.load_workbook(MLB_DATA_PATH, read_only=True, data_only=True)
    ws = wb[wb.sheetnames[0]]  # 'Mar & Apr - All'
    rows = list(ws.iter_rows(min_row=1, max_row=35, values_only=True))
    header = rows[0]
    data_rows = rows[1:]
    wb.close()

    benchmarks = {}

    # Numeric ratings (0-10 scale) — column index: rating, rank
    numeric_map = {
        "Overall Experience": (3, 4),
        "Parking": (5, 6),       # Parking Ingress
        "Gate Entry": (7, 8),
        "Seatview": (9, 10),
        "Concessions": (11, 12),
        "Merchandise": (13, 14),
        "Entertainment": (15, 16),  # Non-Game Entertainment
        "Staff": (17, 18),
        "Parking Egress": (19, 20),
    }

    for metric, (val_idx, rank_idx) in numeric_map.items():
        vals = [r[val_idx] for r in data_rows if r[val_idx] is not None]
        tb_row = next((r for r in data_rows if r[0] == 'TB'), None)
        if vals and tb_row:
            mlb_avg = round(statistics.mean(vals), 2)
            tb_rank = int(tb_row[rank_idx]) if tb_row[rank_idx] is not None else None
            total_clubs = len(vals)
            benchmarks[metric] = (mlb_avg, tb_rank, total_clubs)

    # Concessions grid — % Highly Satisfied
    concess_grid_map = {
        "Cleanliness": (27, 28),         # Highly Satisfied col, Rank col
        "Selection": (31, 32),
        "Value": (35, 36),
        "Speed of Service": (23, 24),    # maps to Customer Service in our data
    }

    for metric, (hs_idx, rank_idx) in concess_grid_map.items():
        vals = [r[hs_idx] for r in data_rows if r[hs_idx] is not None]
        tb_row = next((r for r in data_rows if r[0] == 'TB'), None)
        if vals and tb_row:
            mlb_avg = round(statistics.mean(vals) * 100, 1)  # convert to %
            tb_val = round(tb_row[hs_idx] * 100, 1) if tb_row[hs_idx] else None
            tb_rank = int(tb_row[rank_idx]) if tb_row[rank_idx] is not None else None
            total_clubs = len(vals)
            benchmarks[f"Concess_{metric}"] = (mlb_avg, tb_rank, total_clubs)

    # Concessions Wait — rank
    wait_vals = [r[40] for r in data_rows if r[40] is not None]
    tb_row = next((r for r in data_rows if r[0] == 'TB'), None)
    if wait_vals and tb_row:
        # "Better" % is col 39, rank is col 40
        better_vals = [r[39] for r in data_rows if r[39] is not None]
        mlb_avg = round(statistics.mean(better_vals) * 100, 1) if better_vals else None
        tb_rank = int(tb_row[40]) if tb_row[40] is not None else None
        total_clubs = len(wait_vals)
        benchmarks["Concess_Wait"] = (mlb_avg, tb_rank, total_clubs)

    # Entertainment grid — % Highly Satisfied
    entertain_map = {
        "Pregame Content": (43, 44),
        "Scoreboard": (47, 48),
        "Music": (51, 52),
        "Games/Contests": (55, 56),
        "Kids Activities": (59, 60),
    }

    for metric, (hs_idx, rank_idx) in entertain_map.items():
        vals = [r[hs_idx] for r in data_rows if r[hs_idx] is not None]
        tb_row = next((r for r in data_rows if r[0] == 'TB'), None)
        if vals and tb_row:
            mlb_avg = round(statistics.mean(vals) * 100, 1)
            tb_rank = int(tb_row[rank_idx]) if tb_row[rank_idx] is not None else None
            total_clubs = len(vals)
            benchmarks[f"Entertain_{metric}"] = (mlb_avg, tb_rank, total_clubs)

    return benchmarks


def load_concessions_quality():
    """Load food-item-level quality data from the Concessions Quality Excel.
    Returns a list of tuples: (item_name, tb_pct, mlb_avg_pct, rank, total_clubs)
    sorted by rank (best first)."""
    import statistics

    wb = openpyxl.load_workbook(CONCESSIONS_QUALITY_PATH, read_only=True, data_only=True)
    ws = wb[wb.sheetnames[0]]  # 'Mar & Apr'
    rows = list(ws.iter_rows(min_row=1, max_row=36, values_only=True))
    wb.close()

    header_row = rows[0]  # Category names in row 1
    data_rows = rows[3:]  # Data starts at row 4 (after header, blank, sub-header)

    # Parse category positions from header row
    categories = []
    for i, val in enumerate(header_row):
        if val and i > 0:
            categories.append((val, i))

    # Find TB row
    tb_row = next((r for r in data_rows if r[0] == 'TB'), None)
    if not tb_row:
        return []

    results = []
    for cat_name, start_idx in categories:
        # Each category has 3 columns: % Highly Satisfied, Count, Rank
        tb_pct = tb_row[start_idx]
        tb_count = tb_row[start_idx + 1]
        tb_rank = tb_row[start_idx + 2]

        # Calculate MLB average for this item
        all_pcts = [r[start_idx] for r in data_rows if r[start_idx] is not None and r[0] is not None]
        total_clubs = len(all_pcts)
        mlb_avg = statistics.mean(all_pcts) if all_pcts else None

        if tb_pct is not None and tb_rank is not None:
            results.append((
                cat_name,
                round(tb_pct * 100, 1),
                round(mlb_avg * 100, 1) if mlb_avg else None,
                int(tb_rank),
                total_clubs,
            ))

    # Sort by rank (best first)
    results.sort(key=lambda x: x[3])
    return results


# ── Monthly Rank Trends (hardcoded from MLB VOC PDF pages 4, 6, 8) ──
# Format: category -> (2024 Mar/Apr rank, 2025 Mar/Apr rank, 2026 Mar/Apr rank)
RANK_TRENDS = {
    # Numerical Report Card (PDF page 4)
    "Overall Experience":      (4,  6,   8),
    "Parking Ingress":         (1,  None, 4),
    "Gate Entry":              (None, 1,  2),
    "Seatview":                (4,  None, 9),
    "Concessions":             (2,  7,   8),
    "Merchandise":             (None, 1,  2),
    "Non-Game Entertainment":  (None, None, 17),
    "Staff":                   (None, 5,  7),
    "Parking Egress":          (2,  1,   2),
    # Concessions Report Card (PDF page 6)
    "Conc. Wait Times":        (3,  1,  17),
    "Conc. Cleanliness":       (3,  1,   6),
    "Conc. Speed of Service":  (3,  1,  13),
    "Conc. Selection":         (5, 18,  11),
    "Conc. Value":             (8,  1,  13),
    # Entertainment Report Card (PDF page 8)
    "Pregame Content":         (15, 24, 15),
    "Scoreboard":              (14, 29, 18),
    "Music":                   (14, 21, 23),
    "Games/Contests":          (16, 26, 18),
    "Kids Activities":         (10, 28, 18),
}


# ═══════════════════════════════════════════════════════════════════════
# SQL HELPERS
# ═══════════════════════════════════════════════════════════════════════

def _baseline_avg(col, baseline_season):
    return f"AVG(CASE WHEN SEASON = {baseline_season} AND {col} < 80 THEN {col} END)"


def _baseline_pct_highly(col, baseline_season, valid_labels):
    return (
        f"ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END) "
        f"/ NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)"
    )


def _current_pct_highly(col, current_season, valid_labels):
    return (
        f"ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END) "
        f"/ NULLIF(SUM(CASE WHEN SEASON = {current_season} AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)"
    )


def _hs_pct_highly(col, hs, valid_labels):
    return (
        f"ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs['start_date']}' AND '{hs['end_date']}' "
        f"AND {col} = 'Highly satisfied' THEN 1 ELSE 0 END) "
        f"/ NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs['start_date']}' AND '{hs['end_date']}' "
        f"AND {col} IN {valid_labels} THEN 1 ELSE 0 END), 0), 1)"
    )


def _round2(val):
    return round(float(val), 2) if val is not None else None


def _round1(val):
    return round(float(val), 1) if val is not None else None


def _float(val):
    return float(val) if val is not None else None


# ═══════════════════════════════════════════════════════════════════════
# QUERY BUILDERS
# ═══════════════════════════════════════════════════════════════════════

def query_core_ratings(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks):
    """Slide 1: Core Satisfaction Ratings (0-10 scale AVG)."""
    # Metrics from VOC_VIEW (CORTEX_AI view)
    voc_metrics = [
        ("Overall Experience", "OVERALL_NUMRAT", "Overall Experience"),
        ("Concessions", "CONCESS_NUMRAT", "Concessions"),
        ("Entertainment", "ENTERTAIN_NUMRAT", "Entertainment"),
        ("Parking Ingress", "PARKING_NUMRAT", "Parking"),
        ("Gate Entry", "GE_NUMRAT", "Gate Entry"),
    ]
    rows = []
    for label, col, bench_key in voc_metrics:
        sql = f"""
        SELECT
            {_baseline_avg(col, baseline_season)},
            AVG(CASE WHEN SEASON = {current_season} AND {col} < 80 THEN {col} END) AS s_current,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND {col} < 80 THEN {col} END) AS hs_prev,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND {col} < 80 THEN {col} END) AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round2(r[0])
        s_cur = _round2(r[1])
        hp = _round2(r[2])
        hc = _round2(r[3])
        bench = benchmarks.get(bench_key, (None, None, None))
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, bench[0], bench[1], bench[2], s_base, s_cur, hp, hc, False, False))

    # Metrics from PARKING_VIEW (SBL view) — Seatview, Merchandise, Parking Egress
    parking_metrics = [
        ("Seatview", "REAL_SEATVIEW_NUMRAT", "Seatview"),
        ("Merchandise", "MERCH_NUMRAT", "Merchandise"),
        ("Parking Egress", "PARKING_EXIT_NUMRAT", "Parking Egress"),
    ]
    for label, col, bench_key in parking_metrics:
        sql = f"""
        SELECT
            {_baseline_avg(col, baseline_season)},
            AVG(CASE WHEN SEASON = {current_season} AND {col} < 80 THEN {col} END) AS s_current,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND {col} < 80 THEN {col} END) AS hs_prev,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND {col} < 80 THEN {col} END) AS hs_curr
        FROM {PARKING_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round2(r[0])
        s_cur = _round2(r[1])
        hp = _round2(r[2])
        hc = _round2(r[3])
        bench = benchmarks.get(bench_key, (None, None, None))
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, bench[0], bench[1], bench[2], s_base, s_cur, hp, hc, False, False))

    # Staff: 2026 uses REAL_STAFF_NUMRAT (0-10); 2024 uses TB_ADDON_4_* (1-5 normalized to 0-10)
    staff_addon_cols = "TB_ADDON_4_1, TB_ADDON_4_3, TB_ADDON_4_5, TB_ADDON_4_6, TB_ADDON_4_11, TB_ADDON_4_38, TB_ADDON_4_10"
    sql = f"""
    SELECT
        (SELECT AVG((col_val - 1) * 2.5)
         FROM {VOC_VIEW}
         UNPIVOT (col_val FOR col_name IN ({staff_addon_cols}))
         WHERE SEASON = {baseline_season} AND STADIUM = 'Tropicana Field'
           AND col_val BETWEEN 1 AND 5) AS s_base,
        AVG(CASE WHEN SEASON = {current_season} AND REAL_STAFF_NUMRAT < 80 THEN REAL_STAFF_NUMRAT END) AS s_current,
        AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND REAL_STAFF_NUMRAT < 80 THEN REAL_STAFF_NUMRAT END) AS hs_prev,
        AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND REAL_STAFF_NUMRAT < 80 THEN REAL_STAFF_NUMRAT END) AS hs_curr
    FROM {PARKING_VIEW}
    WHERE STADIUM = 'Tropicana Field'
      AND SEASON IN ({baseline_season}, {current_season})
    """
    result = run_query(conn, sql)
    r = result[0]
    s_base = _round2(r[0])
    s_cur = _round2(r[1])
    hp = _round2(r[2])
    hc = _round2(r[3])
    bench = benchmarks.get("Staff", (None, None, None))
    if s_cur is not None and (hp is not None or hc is not None):
        rows.append(("Staff", bench[0], bench[1], bench[2], s_base, s_cur, hp, hc, False, False))

    return ("Core Satisfaction Ratings",
            "On a scale from 0 to 10, how would you rate your overall experience at Tropicana Field?",
            rows)


def query_staff_ratings(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks):
    """Slide 2: Staff Satisfaction (normalized 0-10)."""
    staff_pairs = [
        ("Parking Staff", "_1", "_1"),
        ("Fan Host Usher", "_3", "_3"),
        ("Concessions Staff", "_5", "_5"),
        ("Retail Staff", "_6", "_6"),
        ("Security", "_11", "_60"),
        ("Tech Team", "_38", "_38"),
        ("Fan Host Ticket Taker/Scanner", "_10", "_61"),
    ]
    rows = []
    # Individual staff benchmarks are NOT available for TB in the MLB data
    # (TB's individual staff columns are all NULL in the Raw Data Excel).
    # Only the aggregate Staff rating (rank 7/29) exists.
    for label, s4, s23 in staff_pairs:
        col_old = f"TB_ADDON_4{s4}"
        col_new = f"TB_ADDON_23{s23}"
        sql = f"""
        SELECT
            AVG(CASE WHEN SEASON = {baseline_season} AND {col_old} BETWEEN 1 AND 5
                      THEN ({col_old} - 1) * 2.5 END) AS s_base,
            AVG(CASE WHEN SEASON = {current_season} AND {col_new} < 80
                      THEN {col_new} END) AS s_current,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}'
                          AND {col_new} < 80
                      THEN {col_new} END) AS hs_prev,
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}'
                          AND {col_new} < 80
                      THEN {col_new} END) AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round2(r[0])
        s_cur = _round2(r[1])
        hp = _round2(r[2])
        hc = _round2(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            # No individual staff MLB benchmark available for TB — show N/A
            rows.append((label, None, None, None,
                        s_base, s_cur, hp, hc, False, False))
    return ("Staff Satisfaction Ratings",
            "On a scale from 0 to 10, please rate your overall satisfaction with each staff member:",
            rows)


def query_concessions_grid(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks):
    """Concessions Grid (% Highly Satisfied)."""
    metrics = [
        ("Customer Service", "CONCESS_GRID_CUSTSERV_DESC", "Concess_Speed of Service"),
        ("Value", "CONCESS_GRID_VALUE_DESC", "Concess_Value"),
        ("Selection", "CONCESS_GRID_SELECTION_DESC", "Concess_Selection"),
        ("Cleanliness", "CONCESS_GRID_CLEAN_DESC", "Concess_Cleanliness"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    rows = []
    for label, col, bench_key in metrics:
        sql = f"""
        SELECT
            {_baseline_pct_highly(col, baseline_season, valid_labels)},
            {_current_pct_highly(col, current_season, valid_labels)} AS s_current,
            {_hs_pct_highly(col, hs_prev, valid_labels)} AS hs_prev,
            {_hs_pct_highly(col, hs_curr, valid_labels)} AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round1(r[0])
        s_cur = _round1(r[1])
        hp = _round1(r[2])
        hc = _round1(r[3])
        bench = benchmarks.get(bench_key, (None, None, None))
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, bench[0], bench[1], bench[2], s_base, s_cur, hp, hc, True, False))
    return rows


def query_concessions_wait(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks):
    """Concessions Wait-Time Expectation Extremes."""
    rows = []
    valid_wait = "('Much more than what I expected','Slightly more than what I expected','About what I expected','Slightly less what I expected','Much less than what I expected')"
    # MLB Wait benchmark measures "% Better than expected" (Better/Equal/Worse)
    # which is not directly comparable to our "Much Less" / "Much More" breakout.
    # Use the aggregate wait rank for the "% Much Less" row (positive metric).
    wait_bench = benchmarks.get("Concess_Wait", (None, None, None))
    for label_name, text_val, invert, use_bench in [
        ("% Much Less Wait Than Expected", "Much less than what I expected", False, True),
        ("% Much More Wait Than Expected", "Much more than what I expected", True, False),
    ]:
        sql = f"""
        SELECT
            ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND CONCESS_WAIT_EXPECT_DESC = '{text_val}' THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND CONCESS_WAIT_EXPECT_DESC IN {valid_wait} THEN 1 ELSE 0 END), 0), 1)
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        r = run_query(conn, sql)[0]
        s_base = _float(r[0])
        s_cur = _float(r[1])
        hp = _float(r[2])
        hc = _float(r[3])
        if s_cur is not None or hp is not None or hc is not None:
            if use_bench:
                rows.append((label_name, wait_bench[0], wait_bench[1], wait_bench[2],
                            s_base, s_cur, hp, hc, True, invert))
            else:
                rows.append((label_name, None, None, None,
                            s_base, s_cur, hp, hc, True, invert))
    return rows


def query_entertainment_grid(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks):
    """Entertainment Grid (% Highly Satisfied)."""
    metrics = [
        ("Pregame Content", "ENTERTAIN_GRID_PREGAME_CONTENT_DESC", "Entertain_Pregame Content"),
        ("Scoreboard", "ENTERTAIN_GRID_SCOREBOARD_DESC", "Entertain_Scoreboard"),
        ("Music", "ENTERTAIN_GRID_MUSIC_DESC", "Entertain_Music"),
        ("Games/Contests", "ENTERTAIN_GRID_GAMES_DESC", "Entertain_Games/Contests"),
        ("Kids Activities", "ENTERTAIN_GRID_KIDS_ACTIVITIES_DESC", "Entertain_Kids Activities"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    rows = []
    for label, col, bench_key in metrics:
        sql = f"""
        SELECT
            {_baseline_pct_highly(col, baseline_season, valid_labels)},
            {_current_pct_highly(col, current_season, valid_labels)} AS s_current,
            {_hs_pct_highly(col, hs_prev, valid_labels)} AS hs_prev,
            {_hs_pct_highly(col, hs_curr, valid_labels)} AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round1(r[0])
        s_cur = _round1(r[1])
        hp = _round1(r[2])
        hc = _round1(r[3])
        bench = benchmarks.get(bench_key, (None, None, None))
        if s_cur is not None and (hp is not None or hc is not None):
            rows.append((label, bench[0], bench[1], bench[2], s_base, s_cur, hp, hc, True, False))
    return rows


def query_food_quality(conn, baseline_season, current_season, hs_prev, hs_curr, concess_quality):
    """Food Quality (% Highly Satisfied). MLB benchmark from Concessions Quality Excel."""
    # Map our internal food item names to the Concessions Quality Excel item names
    food_to_mlb_map = {
        "Hotdog": "Hot Dogs",
        "Burgers": "Burgers",
        "Chicken": "Chicken Tenders",
        "Pizza": "Pizza",
        "Nachos": "Nachos",
        "Fries": "Fries",
        "Alcohol": "Alcoholic Beverages",
        "Non-Alcohol Beverages": "Non-alcoholic Beverages",
        "Ice Cream": "Ice Cream",
        "Nuts": "Nuts",
        "Popcorn": "Popcorn",
        "Pretzels": "Pretzels",
        "Other Dessert": "Other Desserts",
        "Other Entree": "Other Entrees",
        "Sandwich": "Sandwiches",
        "Sausage": "Sausages",
    }
    # Build lookup from concessions quality data: item_name -> (tb_pct, mlb_avg, rank, total)
    cq_lookup = {item: (tb_pct, mlb_avg, rank, total)
                 for item, tb_pct, mlb_avg, rank, total in concess_quality}

    food_items = [
        ("Hotdog", "CONCESS_QUALITY_HOTDOG_DESC"),
        ("Burgers", "CONCESS_QUALITY_BURGERS_DESC"),
        ("Chicken", "CONCESS_QUALITY_CHICKEN_DESC"),
        ("Pizza", "CONCESS_QUALITY_PIZZA_DESC"),
        ("Nachos", "CONCESS_QUALITY_NACHOS_DESC"),
        ("Fries", "CONCESS_QUALITY_FRIES_DESC"),
        ("Alcohol", "CONCESS_QUALITY_ALCOHOL_DESC"),
        ("Non-Alcohol Beverages", "CONCESS_QUALITY_NONALCOHOL_DESC"),
        ("Ice Cream", "CONCESS_QUALITY_ICECREAM_DESC"),
        ("Nuts", "CONCESS_QUALITY_NUTS_DESC"),
        ("Popcorn", "CONCESS_QUALITY_POPCORN_DESC"),
        ("Pretzels", "CONCESS_QUALITY_PRETZELS_DESC"),
        ("Other Dessert", "CONCESS_QUALITY_OTHER_DESSERT_DESC"),
        ("Other Entree", "CONCESS_QUALITY_OTHER_ENTREE_DESC"),
        ("Sandwich", "CONCESS_QUALITY_SANDWICH_DESC"),
        ("Sausage", "CONCESS_QUALITY_SAUSAGE_DESC"),
    ]
    valid_labels = "('Highly satisfied','Somewhat satisfied','Somewhat dissatisfied','Highly dissatisfied')"
    rows = []
    for label, col in food_items:
        sql = f"""
        SELECT
            {_baseline_pct_highly(col, baseline_season, valid_labels)},
            {_current_pct_highly(col, current_season, valid_labels)} AS s_current,
            {_hs_pct_highly(col, hs_prev, valid_labels)} AS hs_prev,
            {_hs_pct_highly(col, hs_curr, valid_labels)} AS hs_curr
        FROM {VOC_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        result = run_query(conn, sql)
        r = result[0]
        s_base = _round1(r[0])
        s_cur = _round1(r[1])
        hp = _round1(r[2])
        hc = _round1(r[3])
        if s_cur is not None and (hp is not None or hc is not None):
            # Look up MLB benchmark for this food item
            mlb_key = food_to_mlb_map.get(label)
            cq_data = cq_lookup.get(mlb_key) if mlb_key else None
            if cq_data:
                mlb_avg, rank, total = cq_data[1], cq_data[2], cq_data[3]
                rows.append((label, mlb_avg, rank, total, s_base, s_cur, hp, hc, True, False))
            else:
                rows.append((label, None, None, None, s_base, s_cur, hp, hc, True, False))
    return rows


def query_gate_entry(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks):
    """Gate Entry Time Expectation."""
    rows = []
    # The MLB benchmark for Gate Entry is a 0-10 numerical rating,
    # not comparable to our % expectation breakdown. Show N/A.
    for label_name, code_val, invert in [
        ("% Less Than Expected", 3, False),
        ("% More Than Expected", 1, True),
    ]:
        sql = f"""
        SELECT
            ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND GE_TIME_EXPECT < 80 AND GE_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND GE_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND GE_TIME_EXPECT < 80 AND GE_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND GE_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND GE_TIME_EXPECT < 80 AND GE_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND GE_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1),
            ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND GE_TIME_EXPECT < 80 AND GE_TIME_EXPECT = {code_val} THEN 1 ELSE 0 END)
                / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND GE_TIME_EXPECT < 80 THEN 1 ELSE 0 END), 0), 1)
        FROM {PARKING_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        r = run_query(conn, sql)[0]
        s_base = _float(r[0])
        s_cur = _float(r[1])
        hp = _float(r[2])
        hc = _float(r[3])
        if s_cur is not None or hp is not None or hc is not None:
            rows.append((label_name, None, None, None,
                        s_base, s_cur, hp, hc, True, invert))
    return rows


def query_parking(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks):
    """Parking Experience — numeric ratings."""
    rows = []
    park_bench = benchmarks.get("Parking", (None, None, None))
    egress_bench = benchmarks.get("Parking Egress", (None, None, None))

    for label, col, bench in [
        ("Finding & Parking (AVG 0-10)", "PARKING_NUMRAT", park_bench),
        ("Exiting Parking Lot (AVG 0-10)", "PARKING_EXIT_NUMRAT", egress_bench),
    ]:
        sql = f"""
        SELECT
            {_baseline_avg(col, baseline_season)},
            AVG(CASE WHEN SEASON = {current_season} AND {col} < 80 THEN {col} END),
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND {col} < 80 THEN {col} END),
            AVG(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND {col} < 80 THEN {col} END)
        FROM {PARKING_VIEW}
        WHERE STADIUM = 'Tropicana Field'
          AND SEASON IN ({baseline_season}, {current_season})
        """
        r = run_query(conn, sql)[0]
        s_base = _round2(r[0])
        s_cur = _round2(r[1])
        hp = _round2(r[2])
        hc = _round2(r[3])
        if s_cur is not None or hp is not None or hc is not None:
            rows.append((label, bench[0], bench[1], bench[2], s_base, s_cur, hp, hc, False, False))
    return rows


def query_seat_value(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks):
    """Seat View Value: % Good Value."""
    rows = []
    # Seatview benchmark is a 0-10 numerical rating, not comparable to % Good Value
    sql = f"""
    SELECT
        ROUND(100.0 * SUM(CASE WHEN SEASON = {baseline_season} AND SEATVIEW_VALUE < 80 AND SEATVIEW_VALUE = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN SEASON = {baseline_season} AND SEATVIEW_VALUE < 80 THEN 1 ELSE 0 END), 0), 1),
        ROUND(100.0 * SUM(CASE WHEN SEASON = {current_season} AND SEATVIEW_VALUE < 80 AND SEATVIEW_VALUE = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN SEASON = {current_season} AND SEATVIEW_VALUE < 80 THEN 1 ELSE 0 END), 0), 1),
        ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND SEATVIEW_VALUE < 80 AND SEATVIEW_VALUE = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_prev["start_date"]}' AND '{hs_prev["end_date"]}' AND SEATVIEW_VALUE < 80 THEN 1 ELSE 0 END), 0), 1),
        ROUND(100.0 * SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND SEATVIEW_VALUE < 80 AND SEATVIEW_VALUE = 1 THEN 1 ELSE 0 END)
            / NULLIF(SUM(CASE WHEN GAME_DATE BETWEEN '{hs_curr["start_date"]}' AND '{hs_curr["end_date"]}' AND SEATVIEW_VALUE < 80 THEN 1 ELSE 0 END), 0), 1)
    FROM {PARKING_VIEW}
    WHERE STADIUM = 'Tropicana Field'
      AND SEASON IN ({baseline_season}, {current_season})
    """
    r = run_query(conn, sql)[0]
    s_base = _float(r[0])
    s_cur = _float(r[1])
    hp = _float(r[2])
    hc = _float(r[3])
    if s_cur is not None or hp is not None or hc is not None:
        rows.append(("% Good Seat Value", None, None, None,
                    s_base, s_cur, hp, hc, True, False))
    return rows


# ═══════════════════════════════════════════════════════════════════════
# FORMATTING HELPERS
# ═══════════════════════════════════════════════════════════════════════

def pct_diff(a, b):
    if a is None or b is None or a == 0:
        return None
    return round((b - a) / a * 100, 1)


def fmt_val(val, is_pct):
    if val is None:
        return "N/A"
    return f"{val:.1f}%" if is_pct else f"{val:.2f}"


def fmt_diff(diff):
    if diff is None:
        return "N/A"
    sign = "+" if diff > 0 else ""
    return f"{sign}{diff:.1f}%"


def fmt_rank(rank, total):
    if rank is None:
        return "—"
    if total:
        return f"{rank}/{total}"
    return str(rank)


def get_rank_trend(r24, r26):
    """Return (trend_text, trend_color) comparing '24 rank to '26 rank."""
    if r24 is None:
        return ("N/A", MUTED_GRAY)
    if r26 < r24:
        return (f"\u2191 {r24}\u2192{r26}", GREEN)
    elif r26 > r24:
        return (f"\u2193 {r24}\u2192{r26}", RED)
    else:
        return (f"\u2192 {r26}", DARK_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# PPTX BUILDER
# ═══════════════════════════════════════════════════════════════════════

def set_cell_text(cell, text, font_size=10, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def set_cell_fill(cell, color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = color


def add_rect(slide, prs, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=10, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    run.font.italic = italic
    return tb


def add_branded_header(slide, prs, title_text):
    add_rect(slide, prs, 0, 0, prs.slide_width, Inches(0.95), NAVY)
    add_rect(slide, prs, 0, Inches(0.95), prs.slide_width, Inches(0.06), YELLOW)
    add_text(slide, Inches(0.45), Inches(0.10), Inches(11), Inches(0.30),
             "VOC POST-ATTENDANCE SURVEY  \u00b7  HOMESTAND COMPARISON",
             size=10, bold=True, color=SKY)
    add_text(slide, Inches(0.45), Inches(0.36), Inches(11), Inches(0.55),
             title_text, size=22, bold=True, color=WHITE)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(12.35), Inches(0.18), height=Inches(0.65))


def add_branded_footer(slide, prs, page_num, total_pages):
    add_rect(slide, prs, 0, Inches(7.20), prs.slide_width, Inches(0.30), NAVY)
    add_text(slide, Inches(0.45), Inches(7.23), Inches(8), Inches(0.25),
             "Tampa Bay Rays  \u00b7  Strategy & Analytics",
             size=9, color=SKY)
    add_text(slide, Inches(10.0), Inches(7.23), Inches(3.0), Inches(0.25),
             f"Source: Qualtrics VOC + MLB VOC Program  |  {page_num} / {total_pages}",
             size=9, color=SKY, align=PP_ALIGN.RIGHT)


def build_table(slide, data_rows, top, left, table_width, col_headers, col_widths,
                include_mlb=True, row_height_override=None, rank_trend_map=None):
    """Build a branded data table on the slide.

    data_rows: list of tuples:
      (metric, mlb_avg, rank, total_clubs, s_base, s_cur, hp, hc, is_pct, invert_diff)
    rank_trend_map: optional dict mapping metric name -> (trend_text, trend_color)
      If provided, adds a "Trend" column at the end.
    """
    # If rank_trend_map is provided, add Trend column
    if rank_trend_map:
        col_headers = list(col_headers) + ["Trend"]
        col_widths = list(col_widths)
        # Shrink existing columns slightly to make room
        col_widths = [Inches(w.inches * 0.88) for w in col_widths] + [Inches(1.5)]

    n_cols = len(col_headers)
    n_rows = len(data_rows) + 1  # +1 for header

    if row_height_override:
        row_height = row_height_override
    else:
        row_height = Inches(0.38) if len(data_rows) <= 8 else Inches(0.32)
    header_height = Inches(0.45)
    table_height = header_height + row_height * len(data_rows)

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, int(table_height))
    table = table_shape.table

    for i, w in enumerate(col_widths):
        table.columns[i].width = int(w)

    # Header row
    table.rows[0].height = int(header_height)
    for ci, hdr in enumerate(col_headers):
        cell = table.cell(0, ci)
        set_cell_fill(cell, MEDIUM_BLUE)
        set_cell_text(cell, hdr, font_size=9, bold=True, color=WHITE)

    # Data rows
    for ri, row_data in enumerate(data_rows):
        metric, mlb_avg, rank, total_clubs, s_base, s_cur, hp, hc, is_pct, invert_diff = row_data
        row_idx = ri + 1
        table.rows[row_idx].height = int(row_height)

        hs_diff = pct_diff(hp, hc)
        baseline_diff = pct_diff(s_base, s_cur)

        if invert_diff:
            if hs_diff is not None:
                hs_diff = -hs_diff
            if baseline_diff is not None:
                baseline_diff = -baseline_diff

        bg_color = LIGHT_BG if ri % 2 == 0 else WHITE

        if include_mlb:
            # 9-column: Metric | MLB AVG | Rank | '24 Score | '26 AVG | HS2 | HS3 | HS3vsHS2 | '26vs'24
            cell_values = [
                (metric, 9, False, DARK_GRAY, PP_ALIGN.LEFT, None),
                (fmt_val(mlb_avg, is_pct) if mlb_avg is not None else "—", 9, False, MUTED_GRAY, PP_ALIGN.CENTER, None),
                (fmt_rank(rank, total_clubs), 9, True if rank and rank <= 5 else False, GOLD if rank and rank <= 5 else DARK_GRAY, PP_ALIGN.CENTER, None),
                (fmt_val(s_base, is_pct), 9, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                (fmt_val(s_cur, is_pct), 9, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                (fmt_val(hp, is_pct), 9, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                (fmt_val(hc, is_pct), 9, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                (None, 9, True, None, PP_ALIGN.CENTER, hs_diff),
                (None, 9, True, None, PP_ALIGN.CENTER, baseline_diff),
            ]
        else:
            # 7-column without MLB: Metric | '24 Score | '26 AVG | HS2 | HS3 | HS3vsHS2 | '26vs'24
            cell_values = [
                (metric, 9, False, DARK_GRAY, PP_ALIGN.LEFT, None),
                (fmt_val(s_base, is_pct), 9, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                (fmt_val(s_cur, is_pct), 9, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                (fmt_val(hp, is_pct), 9, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                (fmt_val(hc, is_pct), 9, False, DARK_GRAY, PP_ALIGN.CENTER, None),
                (None, 9, True, None, PP_ALIGN.CENTER, hs_diff),
                (None, 9, True, None, PP_ALIGN.CENTER, baseline_diff),
            ]

        for ci, (text, fsize, is_bold_or_diff, fcolor, align, diff_val) in enumerate(cell_values):
            cell = table.cell(row_idx, ci)
            set_cell_fill(cell, bg_color)
            if diff_val is not None:
                # This is a diff column
                dc = GREEN if diff_val > 0 else (RED if diff_val < 0 else DARK_GRAY)
                set_cell_text(cell, fmt_diff(diff_val), font_size=fsize, bold=True, color=dc)
            elif text is not None:
                set_cell_text(cell, text, font_size=fsize, bold=is_bold_or_diff, color=fcolor, alignment=align)
            else:
                set_cell_text(cell, "N/A", font_size=fsize, color=MUTED_GRAY)

        # Render Trend column if rank_trend_map is provided
        if rank_trend_map:
            trend_ci = len(cell_values)
            cell = table.cell(row_idx, trend_ci)
            set_cell_fill(cell, bg_color)
            trend_info = rank_trend_map.get(metric)
            if trend_info:
                t_text, t_color = trend_info
                set_cell_text(cell, t_text, font_size=8, bold=True, color=t_color)
            else:
                set_cell_text(cell, "N/A", font_size=8, color=MUTED_GRAY)

    return table_height


# ═══════════════════════════════════════════════════════════════════════
# MAIN BUILD
# ═══════════════════════════════════════════════════════════════════════

def build_report():
    config = load_config()
    hs_prev = get_homestand(config, 2026, 2)
    hs_curr = get_homestand(config, 2026, 3)
    baseline_season = 2024
    current_season = 2026

    print("Loading MLB benchmarks...")
    benchmarks = load_mlb_benchmarks()
    print(f"  Loaded {len(benchmarks)} benchmark metrics")

    print("Loading concessions quality data...")
    concess_quality = load_concessions_quality()
    print(f"  Loaded {len(concess_quality)} food item benchmarks")

    print("Connecting to Snowflake (SSO)...")
    conn = get_connection()
    print("Connected. Querying VOC data...")

    # Get response counts
    sql = f"""
    SELECT
        COUNT(CASE WHEN SEASON = {baseline_season} THEN 1 END),
        COUNT(DISTINCT CASE WHEN SEASON = {baseline_season} THEN GAME_DATE END),
        COUNT(CASE WHEN SEASON = {current_season} THEN 1 END),
        COUNT(DISTINCT CASE WHEN SEASON = {current_season} THEN GAME_DATE END)
    FROM {VOC_VIEW}
    WHERE STADIUM = 'Tropicana Field'
      AND SEASON IN ({baseline_season}, {current_season})
    """
    r = run_query(conn, sql)[0]
    base_resp, base_games, cur_resp, cur_games = int(r[0]), int(r[1]), int(r[2]), int(r[3])

    print("  [1/7] Core Satisfaction Ratings...")
    core_rows = query_core_ratings(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks)

    print("  [2/7] Staff Satisfaction Ratings...")
    staff_rows = query_staff_ratings(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks)

    print("  [3/7] Concessions Grid + Wait...")
    concess_grid_rows = query_concessions_grid(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks)
    concess_wait_rows = query_concessions_wait(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks)

    print("  [4/7] Food Quality...")
    food_rows = query_food_quality(conn, baseline_season, current_season, hs_prev, hs_curr, concess_quality)

    print("  [5/7] Entertainment Grid...")
    entertain_rows = query_entertainment_grid(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks)

    print("  [6/7] Gate Entry + Parking...")
    gate_rows = query_gate_entry(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks)
    parking_rows = query_parking(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks)

    print("  [7/7] Seat Value...")
    seat_rows = query_seat_value(conn, baseline_season, current_season, hs_prev, hs_curr, benchmarks)

    conn.close()
    print("Queries complete. Building PPTX...")

    # Build PPTX
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total_pages = 8

    # Context note
    context_note = (
        f"'24 = Full Season ({base_games} games, {base_resp:,} responses)  |  "
        f"'26 = {cur_games} games ({cur_resp:,} responses)  |  "
        f"HS2 = 04/20-04/26 ({hs_prev['games']}g)  |  "
        f"HS3 = 05/01-05/06 ({hs_curr['games']}g)  |  "
        f"MLB Benchmark = Mar & Apr 2026 (30 clubs)"
    )

    # 9-column headers
    col_headers_9 = [
        "Metric", "MLB\nAVG", "Rank\n/30", "'24\nScore", "'26\nAVG",
        "HS2\nScore", "HS3\nScore", "HS3 vs HS2\n% Diff", "'26 vs '24\n% Diff"
    ]
    col_widths_9 = [
        Inches(2.4), Inches(0.9), Inches(0.8), Inches(1.1), Inches(1.1),
        Inches(1.1), Inches(1.1), Inches(1.8), Inches(1.8)
    ]

    # 7-column headers (for food quality which has no MLB benchmarks)
    col_headers_7 = [
        "Metric", "'24\nScore", "'26\nAVG", "HS2\nScore", "HS3\nScore",
        "HS3 vs HS2\n% Diff", "'26 vs '24\n% Diff"
    ]
    col_widths_7 = [
        Inches(2.8), Inches(1.3), Inches(1.3), Inches(1.3), Inches(1.3),
        Inches(2.0), Inches(2.0)
    ]

    # ── Build rank trend maps for each slide (maps metric label -> (text, color)) ──
    # Slide 1: Core Ratings
    core_trend_map = {}
    for label, trend_key in [("Overall Experience", "Overall Experience"),
                              ("Concessions", "Concessions"),
                              ("Entertainment", "Non-Game Entertainment"),
                              ("Parking Ingress", "Parking Ingress"),
                              ("Gate Entry", "Gate Entry"),
                              ("Seatview", "Seatview"),
                              ("Merchandise", "Merchandise"),
                              ("Staff", "Staff"),
                              ("Parking Egress", "Parking Egress")]:
        td = RANK_TRENDS.get(trend_key)
        if td:
            core_trend_map[label] = get_rank_trend(td[0], td[2])

    # Slide 3: Concessions Grid
    concess_trend_map = {}
    for label, trend_key in [("Customer Service", "Conc. Speed of Service"),
                              ("Value", "Conc. Value"),
                              ("Selection", "Conc. Selection"),
                              ("Cleanliness", "Conc. Cleanliness")]:
        td = RANK_TRENDS.get(trend_key)
        if td:
            concess_trend_map[label] = get_rank_trend(td[0], td[2])
    # Wait time trends
    td = RANK_TRENDS.get("Conc. Wait Times")
    if td:
        concess_trend_map["% Much Less Wait Than Expected"] = get_rank_trend(td[0], td[2])

    # Slide 4: Food Quality — no '24 rank data available from MLB at item level

    # Slide 5: Entertainment Grid
    entertain_trend_map = {}
    for label, trend_key in [("Pregame Content", "Pregame Content"),
                              ("Scoreboard", "Scoreboard"),
                              ("Music", "Music"),
                              ("Games/Contests", "Games/Contests"),
                              ("Kids Activities", "Kids Activities")]:
        td = RANK_TRENDS.get(trend_key)
        if td:
            entertain_trend_map[label] = get_rank_trend(td[0], td[2])

    # ─── SLIDE 1: Core Satisfaction Ratings ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Core Satisfaction Ratings")
    _, subtitle, rows = core_rows
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             subtitle, size=11, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             context_note, size=8, color=MUTED_GRAY)

    # Sort by '26 vs '24 % diff
    def sort_by_yoy(row):
        _, _, _, _, s_base, s_cur, _, _, is_pct, invert = row
        d = pct_diff(s_base, s_cur)
        if d is None: return float('-inf')
        if invert: d = -d
        return -d
    rows = sorted(rows, key=sort_by_yoy)

    build_table(slide, rows, Inches(1.65), Inches(0.45), Inches(12.4),
                col_headers_9, col_widths_9, include_mlb=True, rank_trend_map=core_trend_map)
    add_branded_footer(slide, prs, 1, total_pages)

    # ─── SLIDE 2: Staff Satisfaction Ratings ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Staff Satisfaction Ratings")
    _, subtitle, rows = staff_rows
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             subtitle, size=11, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             context_note, size=8, color=MUTED_GRAY)
    rows = sorted(rows, key=sort_by_yoy)
    build_table(slide, rows, Inches(1.65), Inches(0.45), Inches(12.4),
                col_headers_9, col_widths_9, include_mlb=True)
    add_branded_footer(slide, prs, 2, total_pages)

    # ─── SLIDE 3: Concessions Grid + Wait (combined) ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Concessions Grid (% Highly Satisfied) + Wait Expectations")
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             "Please rate your overall satisfaction with food and/or beverages.",
             size=11, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             context_note, size=8, color=MUTED_GRAY)

    concess_grid_rows_sorted = sorted(concess_grid_rows, key=sort_by_yoy)
    tbl_h = build_table(slide, concess_grid_rows_sorted, Inches(1.65), Inches(0.45), Inches(12.4),
                col_headers_9, col_widths_9, include_mlb=True, rank_trend_map=concess_trend_map)

    # Second table: Wait expectations
    wait_top = Inches(1.65) + tbl_h + Inches(0.35)
    add_text(slide, Inches(0.45), wait_top - Inches(0.25), Inches(12), Inches(0.25),
             "How did waiting for concessions compare to your expectations?",
             size=10, italic=True, color=GRAY)
    build_table(slide, concess_wait_rows, wait_top + Inches(0.05), Inches(0.45), Inches(12.4),
                col_headers_9, col_widths_9, include_mlb=True, rank_trend_map=concess_trend_map)
    add_branded_footer(slide, prs, 3, total_pages)

    # ─── SLIDE 4: Food Quality ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Food Quality (% Highly Satisfied)")
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             "How would you rate the quality (taste, temperature, etc.) of each food / beverage?",
             size=11, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             context_note, size=8, color=MUTED_GRAY)

    food_rows_sorted = sorted(food_rows, key=sort_by_yoy)
    build_table(slide, food_rows_sorted, Inches(1.60), Inches(0.45), Inches(12.4),
                col_headers_9, col_widths_9, include_mlb=True,
                row_height_override=Inches(0.28))
    add_branded_footer(slide, prs, 4, total_pages)

    # ─── SLIDE 5: Entertainment Grid ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Entertainment Grid (% Highly Satisfied)")
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             "Please rate your satisfaction with the following:",
             size=11, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.35), Inches(12.5), Inches(0.25),
             context_note, size=8, color=MUTED_GRAY)

    entertain_rows_sorted = sorted(entertain_rows, key=sort_by_yoy)
    build_table(slide, entertain_rows_sorted, Inches(1.65), Inches(0.45), Inches(12.4),
                col_headers_9, col_widths_9, include_mlb=True, rank_trend_map=entertain_trend_map)
    add_branded_footer(slide, prs, 5, total_pages)

    # ─── SLIDE 6: Gate Entry + Parking + Seat Value (combined) ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Gate Entry + Parking + Seat Value")

    # Gate Entry table
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.25),
             "How did the time it took to pass through security and gate entry compare to expectations?",
             size=10, italic=True, color=GRAY)
    add_text(slide, Inches(0.45), Inches(1.30), Inches(12.5), Inches(0.22),
             context_note, size=7, color=MUTED_GRAY)

    gate_h = build_table(slide, gate_rows, Inches(1.55), Inches(0.45), Inches(12.4),
                col_headers_9, col_widths_9, include_mlb=True,
                row_height_override=Inches(0.32))

    # Parking table
    park_top = Inches(1.55) + gate_h + Inches(0.30)
    add_text(slide, Inches(0.45), park_top - Inches(0.22), Inches(12), Inches(0.22),
             "Parking Experience (Ingress & Egress Ratings)",
             size=10, italic=True, color=GRAY)
    park_h = build_table(slide, parking_rows, park_top + Inches(0.05), Inches(0.45), Inches(12.4),
                col_headers_9, col_widths_9, include_mlb=True,
                row_height_override=Inches(0.32))

    # Seat Value table
    seat_top = park_top + Inches(0.05) + park_h + Inches(0.30)
    add_text(slide, Inches(0.45), seat_top - Inches(0.22), Inches(12), Inches(0.22),
             "Thinking of the money you spent on your ticket, your view was a... (% Good Value)",
             size=10, italic=True, color=GRAY)
    build_table(slide, seat_rows, seat_top + Inches(0.05), Inches(0.45), Inches(12.4),
                col_headers_9, col_widths_9, include_mlb=True,
                row_height_override=Inches(0.32))
    add_branded_footer(slide, prs, 6, total_pages)

    # ─── SLIDE 7: Concessions Food Quality — MLB Rank by Item ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "Concessions Food Quality — MLB Rank by Item")
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             "% Highly Satisfied with food/beverage quality — Tampa Bay vs. 30 MLB clubs (March & April 2026)",
             size=11, italic=True, color=GRAY)

    # Build food quality by item table
    fq_headers = ["Food Item", "TB %\nHighly Sat.", "MLB AVG\n%", "Rank\n/30", "vs. MLB"]
    fq_widths = [Inches(2.8), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.6)]
    n_fq_rows = len(concess_quality) + 1
    fq_row_height = Inches(0.28)
    fq_header_height = Inches(0.45)
    fq_table_height = fq_header_height + fq_row_height * len(concess_quality)

    table_shape = slide.shapes.add_table(n_fq_rows, 5, Inches(0.45), Inches(1.45),
                                          Inches(12.0), int(fq_table_height))
    table = table_shape.table
    for i, w in enumerate(fq_widths):
        table.columns[i].width = int(w)

    # Header
    table.rows[0].height = int(fq_header_height)
    for ci, hdr in enumerate(fq_headers):
        cell = table.cell(0, ci)
        set_cell_fill(cell, MEDIUM_BLUE)
        set_cell_text(cell, hdr, font_size=9, bold=True, color=WHITE)

    # Data rows
    for ri, (item_name, tb_pct, mlb_avg_pct, rank, total_clubs) in enumerate(concess_quality):
        row_idx = ri + 1
        table.rows[row_idx].height = int(fq_row_height)
        bg_color = LIGHT_BG if ri % 2 == 0 else WHITE

        # Food Item name
        cell = table.cell(row_idx, 0)
        set_cell_fill(cell, bg_color)
        set_cell_text(cell, item_name, font_size=9, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

        # TB % Highly Satisfied
        cell = table.cell(row_idx, 1)
        set_cell_fill(cell, bg_color)
        set_cell_text(cell, f"{tb_pct:.1f}%", font_size=9, color=DARK_GRAY)

        # MLB AVG %
        cell = table.cell(row_idx, 2)
        set_cell_fill(cell, bg_color)
        set_cell_text(cell, f"{mlb_avg_pct:.1f}%" if mlb_avg_pct else "—", font_size=9, color=MUTED_GRAY)

        # Rank
        cell = table.cell(row_idx, 3)
        set_cell_fill(cell, bg_color)
        rank_color = GOLD if rank <= 5 else (GREEN if rank <= 10 else (RED if rank >= 25 else DARK_GRAY))
        set_cell_text(cell, f"{rank}/{total_clubs}", font_size=9, bold=True, color=rank_color)

        # vs. MLB
        cell = table.cell(row_idx, 4)
        set_cell_fill(cell, bg_color)
        if mlb_avg_pct:
            vs_mlb = round(tb_pct - mlb_avg_pct, 1)
            diff_color = GREEN if vs_mlb > 0 else (RED if vs_mlb < 0 else DARK_GRAY)
            sign = "+" if vs_mlb > 0 else ""
            set_cell_text(cell, f"{sign}{vs_mlb:.1f}%", font_size=9, bold=True, color=diff_color)
        else:
            set_cell_text(cell, "—", font_size=9, color=MUTED_GRAY)

    add_branded_footer(slide, prs, 7, total_pages)

    # ─── SLIDE 8: Monthly Rank Trends ───
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    add_branded_header(slide, prs, "MLB Rank Trends — Report Card Categories")
    add_text(slide, Inches(0.45), Inches(1.08), Inches(12), Inches(0.30),
             "Tampa Bay Rays VOC Program rank among MLB clubs — April 2024 vs. April 2026",
             size=11, italic=True, color=GRAY)

    # Build rank trends table (4 columns: Category | '24 Rank | '26 Rank | Trend)
    trend_headers = ["Category", "'24\nMar/Apr Rank", "'26\nMar/Apr Rank", "Trend"]
    trend_widths = [Inches(3.8), Inches(2.5), Inches(2.5), Inches(3.2)]
    trend_categories = list(RANK_TRENDS.items())
    n_trend_rows = len(trend_categories) + 1
    trend_row_height = Inches(0.28)
    trend_header_height = Inches(0.45)
    trend_table_height = trend_header_height + trend_row_height * len(trend_categories)

    table_shape = slide.shapes.add_table(n_trend_rows, 4, Inches(0.45), Inches(1.45),
                                          Inches(12.0), int(trend_table_height))
    table = table_shape.table
    for i, w in enumerate(trend_widths):
        table.columns[i].width = int(w)

    # Header
    table.rows[0].height = int(trend_header_height)
    for ci, hdr in enumerate(trend_headers):
        cell = table.cell(0, ci)
        set_cell_fill(cell, MEDIUM_BLUE)
        set_cell_text(cell, hdr, font_size=9, bold=True, color=WHITE)

    # Data rows
    for ri, (category, (r24, _r25, r26)) in enumerate(trend_categories):
        row_idx = ri + 1
        table.rows[row_idx].height = int(trend_row_height)
        bg_color = LIGHT_BG if ri % 2 == 0 else WHITE

        # Category name
        cell = table.cell(row_idx, 0)
        set_cell_fill(cell, bg_color)
        set_cell_text(cell, category, font_size=9, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)

        # '24 rank
        cell = table.cell(row_idx, 1)
        set_cell_fill(cell, bg_color)
        if r24 is not None:
            r24_color = GOLD if r24 <= 5 else (GREEN if r24 <= 10 else DARK_GRAY)
            set_cell_text(cell, str(r24), font_size=9, bold=True, color=r24_color)
        else:
            set_cell_text(cell, "N/A", font_size=9, color=MUTED_GRAY)

        # '26 rank
        cell = table.cell(row_idx, 2)
        set_cell_fill(cell, bg_color)
        r26_color = GOLD if r26 <= 5 else (GREEN if r26 <= 10 else (RED if r26 >= 20 else DARK_GRAY))
        set_cell_text(cell, str(r26), font_size=9, bold=True, color=r26_color)

        # Trend arrow (compare '24 to '26 only)
        cell = table.cell(row_idx, 3)
        set_cell_fill(cell, bg_color)
        if r24 is not None:
            if r26 < r24:
                # Lower rank number = better
                trend_text = f"\u2191 Improved ({r24}\u2192{r26})"
                trend_color = GREEN
            elif r26 > r24:
                trend_text = f"\u2193 Declined ({r24}\u2192{r26})"
                trend_color = RED
            else:
                trend_text = f"\u2192 Same ({r26})"
                trend_color = DARK_GRAY
        else:
            trend_text = "New in '26"
            trend_color = MUTED_GRAY
        set_cell_text(cell, trend_text, font_size=9, bold=True, color=trend_color)

    add_branded_footer(slide, prs, 8, total_pages)

    # Save
    prs.save(OUTPUT_PATH)
    print(f"\nDone! Output: {OUTPUT_PATH}")
    print(f"  Slides: {total_pages}")
    print(f"  Layout: 9 columns (with MLB benchmarks)")


if __name__ == "__main__":
    build_report()
